#!/usr/bin/env python3
# predict_developability.py
# MLAbDev — IPI Antibody Developability Prediction Platform
# Author: Hoan Nguyen | IPI Biologics
#
# ══════════════════════════════════════════════════════════════════════════════
# MODELS & PLMs SUPPORTED
# ══════════════════════════════════════════════════════════════════════════════
#
#   --model  transformer_lm      Transformer + PLM embeddings  (best AUC)
#            transformer_onehot  Transformer + one-hot sequences (no PLM)
#            rf                  Random Forest (fast, SHAP, mutagenesis)
#            xgboost             XGBoost
#            cnn                 CNN
#
#   --lm     biophysical         charge, pI, hydrophobicity, R/K/W counts (26d)
#            kmer                1-mer + 2-mer AA frequencies (~440d)
#            onehot              VH+VL one-hot position encoding
#            onehot_vh           VH one-hot only
#            onehot_cdr3         HCDR3 one-hot only
#            ablang              ABlang2 480-dim   pip install ablang2
#            antiberty           AntiBERTy 512-dim  pip install antiberty
#            antiberta2          AntiBERTa2 1024-dim pip install transformers
#            antiberta2-cssp     AntiBERTa2-CSSP 1024-dim
#            igbert              IgBERT 1024-dim
#
# ══════════════════════════════════════════════════════════════════════════════
# TRAINING MODES (transformer_lm)
# ══════════════════════════════════════════════════════════════════════════════
#
#   MODE 1 — Frozen embeddings  (DEFAULT — recommended)
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (frozen) → pre-computed .emb.csv → classifier trains
#   PLM weights: never updated  |  Trained: ~200k classifier params only
#   Best for: n < 10,000  |  CPU  |  same domain as pretraining
#   Your results: ablang ρ_OVA = −0.66 on GDPa3 (n=80)
#
#   MODE 2 — PLM layer unfreezing  (--train --finetune_plm)
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (top layers update via backprop) → classifier trains
#   PLM weights: top N layers updated  |  Trained: ~20M params
#   Best for: n > 20,000  |  GPU  |  domain-shifted sequences
#   No .emb.csv needed — sequences processed in batches during training
#
#   MODE 3 — LoRA  (--train --finetune_plm --peft lora)  [RECOMMENDED]
#   ─────────────────────────────────────────────────────────────────────────
#   sequences → PLM (W frozen, only LoRA A×B trained) → classifier trains
#   PLM weights: W NEVER changes, only low-rank A×B matrices (~400k params)
#   Best for: n > 1,000  |  CPU feasible  |  low forgetting risk
#   No .emb.csv needed — sequences processed in batches during training
#
#   LEVEL 2 — Collaborator fine-tune  (--finetune --pretrained path.pt)
#   ─────────────────────────────────────────────────────────────────────────
#   Load YOUR pretrained .pt → fine-tune classifier on THEIR small dataset
#   PLM: embedded in their local MLAbDev install (ablang2/transformers)
#   Best for: 50–2,000 new antibodies from a collaborator lab
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 1 — RF biophysical (fastest, interpretable, no PLM needed)
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Step 1: K-fold validation (find optimal threshold)
#   python predict_developability.py --kfold 10 \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx \
#       --cost_fn 3.0                # catch FAILs: missing 1 bad Ab costs 3× FP
#   → kfold log + CV_ROC + recommended threshold embedded in BEST_*.pkl
#
#   # Step 2: Train on full dataset
#   python predict_developability.py --train \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx
#   → FINAL_psr_filter_biophysical_rf_ipi_psr_trainset.pkl
#   → SHAP bar + beeswarm + heatmap on train set
#
#   # Step 3: Predict + SHAP waterfall + CDR3 mutagenesis
#   python predict_developability.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm biophysical --model rf \
#       --db data/ipi_psr_trainset.xlsx \
#       --test_target psr_smp_filter \  # ground truth col if different from --target
#       --mutagenesis 50                 # CDR3 heatmaps for first 50 antibodies
#   → predictions .xlsx + ROC + KDE + histogram + SHAP waterfalls PPT + mutagenesis PPT
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 2 — transformer_lm Mode 1 (frozen ABlang, best overall result)
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Step 1: K-fold (CDR3-cluster stratified — no sequence leakage)
#   python predict_developability.py --kfold 10 \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --db data/ipi_psr_trainset.xlsx \
#       --cluster 0.8                # Levenshtein threshold for CDR3 grouping
#   → mean AUC ± std across 10 folds
#   → recommended threshold embedded in BEST_*.pt  (used automatically at predict)
#
#   # Step 2: Train on full dataset (Mode 1 — default, frozen embeddings)
#   python predict_developability.py --train \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --db data/ipi_psr_trainset.xlsx
#   → ipi_psr_trainset.xlsx.ablang.emb.csv  (auto-generated if missing)
#   → FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt
#
#   # Step 3: Predict (uses threshold embedded from kfold automatically)
#   python predict_developability.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --db data/ipi_psr_trainset.xlsx
#   → new_cohort_pred_psr_filter_ablang_transformer_lm_ipi_psr_trainset.xlsx
#   → ROC, KDE, histogram, SHAP beeswarm, waterfall PPT
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 3 — transformer_lm Mode 3 LoRA (PLM adaptation from sequences)
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Sequences → IgBERT → LoRA A×B → classifier
#   # No .emb.csv needed. W_original frozen. Only A,B matrices trained.
#
#   python predict_developability.py --train --finetune_plm \
#       --target psr_filter --lm igbert --model transformer_lm \
#       --db data/ipi_psr_trainset.xlsx \
#       --peft lora        \    # Mode 3 LoRA
#       --lora_r 8         \    # rank: 4=small n, 8=medium(default), 16=large n
#       --lora_alpha 16    \    # scaling = alpha/r = 2.0 (keep 2×rank)
#       --lora_layers 10 11 \   # which IgBERT layers get LoRA (default: last 2)
#       --lr_plm 1e-5      \    # low lr for LoRA matrices (careful update)
#       --lr_classifier 1e-4 \  # normal lr for classifier head
#       --finetune_epochs 20
#   → FINAL_psr_filter_igbert_transformer_lm_ipi_psr_trainset_lora8.pt
#   → lora_weights_psr_filter_igbert_transformer_lm_ipi_psr_trainset.pt (~2 MB)
#
#   # Predict with LoRA model (explicit path needed — multiple checkpoints may exist)
#   python predict_developability.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm igbert --model transformer_lm \
#       --db data/ipi_psr_trainset.xlsx \
#       --model_path models/FINAL_psr_filter_igbert_transformer_lm_ipi_lora8.pt
#
#   # LoRA trainable param count (IgBERT, r=8, layers 10-11):
#   #   2 layers × 2 matrices (Q,V) × 2 × (1024×8) = 65,536 PLM params updated
#   #   vs 110,000,000 total IgBERT params → 0.06% of PLM updated
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 4 — transformer_lm Mode 2 (full layer unfreezing, large dataset)
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Best for n > 20,000, GPU available, domain-shifted sequences
#   python predict_developability.py --train --finetune_plm \
#       --target psr_filter --lm igbert --model transformer_lm \
#       --db data/large_dataset.xlsx \
#       --freeze_plm_layers 10  \   # freeze first 10 of 12 IgBERT layers
#       --lr_plm 1e-6           \   # very slow PLM update (avoids forgetting)
#       --lr_classifier 1e-4    \
#       --finetune_epochs 20
#   → FINAL_psr_filter_igbert_transformer_lm_large_dataset_plmft.pt
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 5 — Level 2: collaborator fine-tunes your pretrained model
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Collaborator has: MLAbDev installed + ablang2 + their 300 antibodies
#   # You share: FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt
#
#   python predict_developability.py --finetune \
#       --pretrained FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --finetune_db their_cohort_300.xlsx \
#       --freeze_layers 1   \   # freeze first classifier layer (safe for n=300)
#       --finetune_lr 1e-6  \   # 100× lower than original lr (avoids forgetting)
#       --finetune_epochs 10
#   → FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_ft_their_cohort_300.pt
#
#   # Predict with their fine-tuned model
#   python predict_developability.py --predict new_targets.xlsx \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --db ipi_psr_trainset.xlsx \
#       --model_path FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_ft_their_cohort_300.pt
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 6 — Selecting which checkpoint to use at predict time
# ══════════════════════════════════════════════════════════════════════════════
#
#   Multiple checkpoints may exist for the same target + lm + model:
#     FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt          (Mode 1)
#     FINAL_psr_filter_igbert_transformer_lm_ipi_psr_trainset_lora8.pt    (Mode 3)
#     FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_ft_lab.pt   (fine-tuned)
#
#   Option A — Auto-discover (default)
#   Uses FINAL_{target}_{lm}_{model}_{db}.pt naming convention.
#   python predict_developability.py --predict cohort.xlsx \
#       --target psr_filter --lm ablang --model transformer_lm \
#       --db ipi_psr_trainset.xlsx
#
#   Option B — Pass checkpoint directly via --db  ← SIMPLEST
#   If --db ends with .pt or .pkl, it is treated as the model checkpoint.
#   No --model_path flag needed. db_stem is inferred from the checkpoint name.
#   python predict_developability.py --predict cohort.xlsx \
#       --target psr_filter --lm igbert --model transformer_lm \
#       --db models/FINAL_psr_filter_igbert_transformer_lm_ipi_lora8.pt
#   → [db] --db is a checkpoint file → using as --model_path
#   → [db] Inferred db_stem from checkpoint: 'ipi_psr_trainset'
#
#   Option C — Explicit --model_path (when --db is still a dataset)
#   python predict_developability.py --predict cohort.xlsx \
#       --target psr_filter --lm igbert --model transformer_lm \
#       --db ipi_psr_trainset.xlsx \
#       --model_path models/FINAL_psr_filter_igbert_transformer_lm_ipi_lora8.pt
#
#   Option D — Auto-fallback (if exact checkpoint not found, lists candidates)
#   If exact FINAL_*.pt is missing, the platform searches and prompts:
#     [load] Found 3 matching checkpoint(s):
#       [0] FINAL_..._lora8.pt      ← most recent, used automatically
#       [1] FINAL_..._plmft.pt
#       [2] FINAL_..._ft_lab300.pt
#     [load] To use a different one: --model_path <path>
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE 7 — Correlation analysis across all models
# ══════════════════════════════════════════════════════════════════════════════
#
#   # Predict with multiple models, then compare all against assay data
#   python utils/developability_correlation.py \
#       --files pred_ablang.xlsx pred_igbert_lora8.xlsx pred_biophysical.xlsx \
#               pred_kmer.xlsx pred_onehot_hcdr3.xlsx \
#       --assay polyreactivity_prscore_ova_avg polyreactivity_prscore_cho_avg \
#       --target psr_filter_ova
#   → Spearman/Pearson heatmap (all models vs all assays)
#   → Per-model scatter plots with ρ annotation
#   → Boxplots PASS vs FAIL per assay
#   → t-SNE coloured by model prediction vs true label
#
# ══════════════════════════════════════════════════════════════════════════════
# KEY FLAGS — QUICK REFERENCE
# ══════════════════════════════════════════════════════════════════════════════
#
#   TRAINING & KFOLD
#     --cluster 0.8            CDR3 Levenshtein clustering threshold
#     --cluster_col CDR3       Sequence for clustering (CDR3/HSEQ/VHVL)
#     --cost_fn 3.0            FN penalty (↑ = lower threshold = catch more FAILs)
#     --cost_fp 1.0            FP penalty (↑ = higher threshold = fewer false alarms)
#     --split 0.8              Train/val fraction for --train (0 = no split)
#
#   PREDICTION
#     --model_path path.pt     Explicit checkpoint path (overrides auto-discovery)
#     --threshold 0.42         Override prediction threshold (else from kfold)
#     --test_target col        Ground-truth column in predict file
#     --mutagenesis N          CDR3 mutagenesis heatmaps (N=50, 0=all antibodies)
#
#   PLM FINE-TUNING (Mode 2 / Mode 3)
#     --finetune_plm           Enable PLM weight updates during --train
#     --peft lora              Mode 3 LoRA (recommended over Mode 2)
#     --lora_r 8               LoRA rank: 4=small n, 8=medium, 16=large n
#     --lora_alpha 16          LoRA scaling = alpha/r (keep at 2×rank)
#     --lora_layers 10 11      Which PLM layers get LoRA (default: last 2)
#     --freeze_plm_layers 10   Mode 2: freeze first N of 12 PLM layers
#     --lr_plm 1e-5            Learning rate for PLM LoRA/unfrozen layers
#     --lr_classifier 1e-4     Learning rate for classifier head
#     --finetune_epochs 20     Epochs for PLM fine-tuning
#
#   LEVEL 2 FINE-TUNING (collaborator)
#     --finetune               Fine-tune from a pretrained checkpoint
#     --pretrained path.pt     Path to pretrained .pt to start from
#     --finetune_db file.xlsx  New dataset for fine-tuning
#     --freeze_layers 1        Classifier layer freezing (all/0/N)
#     --finetune_lr 1e-6       Fine-tuning lr (100× lower than original)
#     --finetune_epochs 10     Fine-tuning epochs
#
#   DECISION GUIDE
#     n < 1,000   any GPU?    → Mode 1 frozen + rf biophysical for mutagenesis
#     n 1–10k     no GPU      → Mode 1 frozen (current IPI datasets)
#     n 1–10k     GPU         → Mode 3 LoRA r=4 if domain-shifted
#     n 10–50k    GPU         → Mode 3 LoRA r=8  ← recommended
#     n > 50k     GPU         → Mode 3 LoRA r=16 or Mode 2
#     collaborator any        → Level 2 --finetune from your pretrained
#
# ══════════════════════════════════════════════════════════════════════════════
# ── Platform fixes — must be FIRST, before any imports ───────────────────────
import os
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("OBJC_DISABLE_INITIALIZE_FORK_SAFETY", "YES")
# ─────────────────────────────────────────────────────────────────────────────
"""
IPI Antibody Developability Prediction Platform
Final Production Version — DEC-2025
Supports: SEC & PSR | XGBoost & RF & CNN & Transformer (One-Hot) & Transformer (LM)
          | ablang, antiberty, antiberta2, antiberta2-cssp, onehot, onehot_vh

Changes (2026):
  * Auto-extract CDR3 from HSEQ when CDR3 column is missing (ANARCI -> regex fallback)
  * Auto-prepend 'C' to any CDR3 missing the conserved N-terminal cysteine
  * RF model: fixed kfold/predict/train calls to pass X_df + embeddings separately
"""

from config import MODEL_DIR, PREDICTION_DIR
import argparse
import os
import re
import warnings
import pandas as pd
import numpy as np
import torch
from pathlib import Path

from embedding_generator import generate_embedding

import sys
import datetime


class _Tee:
    def __init__(self, log_path: str, mode: str = 'w'):
        self._terminal = sys.__stdout__
        self._log = open(log_path, mode, buffering=1, encoding='utf-8')
        print(f"[log] Writing to: {log_path}" + (" (append)" if mode == 'a' else ""), flush=True)

    def write(self, msg):
        self._terminal.write(msg)
        self._log.write(msg)

    def flush(self):
        self._terminal.flush()
        self._log.flush()

    def close(self):
        if self._log and not self._log.closed:
            self._log.close()

    @property
    def encoding(self):
        return self._terminal.encoding

    def isatty(self):
        return False


def _setup_logging(args, db_path: str) -> str:
    ts      = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    lm_tag  = args.lm.replace(",", "_")
    # Normalise onehot_cdr3 → onehot_hcdr3 everywhere
    _LM_NORM = {"onehot_cdr3": "onehot_hcdr3"}
    args.lm  = _LM_NORM.get(args.lm, args.lm)
    lm_tag   = args.lm.replace(",", "_")
    db_stem = Path(db_path).stem if db_path else "default"

    if args.train:
        name     = f"train_{args.target}_{lm_tag}_{args.model}_{db_stem}_{ts}.log"
        log_path = os.path.join(MODEL_DIR, name)
        mode_w   = 'w'
    elif args.kfold:
        name     = f"kfold_{args.target}_{lm_tag}_{args.model}_{db_stem}_k{args.kfold}_{ts}.log"
        log_path = os.path.join(MODEL_DIR, name)
        mode_w   = 'w'
    elif args.predict:
        p        = Path(args.predict)
        name     = f"predict_{p.stem}_{args.target}_{lm_tag}_{args.model}_{db_stem}_{ts}.log"
        log_path = str(p.parent / name)
        mode_w   = 'w'
    elif args.build_embedding:
        p        = Path(args.build_embedding)
        name     = f"embedding_{p.stem}_{lm_tag}_{ts}.log"
        log_path = str(p.parent / name)
        mode_w   = 'w'
    elif getattr(args, 'split_dataset', False):
        name     = f"split_{Path(db_path).stem}_{ts}.log"
        log_path = str(Path(db_path).parent / name)
        mode_w   = 'w'
    else:
        log_path = f"ipi_{ts}.log"
        mode_w   = 'w'

    os.makedirs(os.path.dirname(log_path) or '.', exist_ok=True)
    sys.stdout = _Tee(log_path, mode=mode_w)
    ts_str = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    print(f"[log] {'='*58}")
    print(f"[log] Started  : {ts_str}")
    print(f"[log] Command  : {' '.join(sys.argv)}")
    print(f"[log] Platform : {sys.platform}  Python {sys.version.split()[0]}")
    print(f"[log] {'='*58}")
    print()
    return log_path


from models.xgboost import XGBoostModel
from models.randomforest import RandomForestModel
from models.cnn import CNNModel
from models.transformer_onehot import TransformerOneHotModel
from models.transformer_lm import TransformerLMModel

os.makedirs(MODEL_DIR, exist_ok=True)
os.makedirs(PREDICTION_DIR, exist_ok=True)


def get_default_db_path():
    data_dir = "data"
    if not os.path.exists(data_dir):
        return None
    files = [f for f in os.listdir(data_dir)
             if f.startswith("ipi_antibodydb") and f.endswith(".xlsx")]
    if not files:
        return None
    files.sort(key=lambda x: os.path.getmtime(os.path.join(data_dir, x)), reverse=True)
    return os.path.join(data_dir, files[0])


# ===========================================================================
# CDR3 EXTRACTION + C-PREFIX FIX
# ===========================================================================

def _extract_cdr3_single_anarci(hseq: str) -> str:
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            from anarci import anarci as _anarci
            res = _anarci([("s", hseq)], scheme="imgt", output=False)
        numbered_seqs = res[0]
        if not numbered_seqs or numbered_seqs[0] is None:
            return ""
        domains = numbered_seqs[0]
        if not domains:
            return ""
        numbering, chain_type = domains[0]
        return "".join(
            aa for (pos, ins), aa in numbering
            if 105 <= pos <= 117 and aa != "-"
        )
    except Exception:
        return ""


def _extract_cdr3_regex(hseq: str) -> str:
    hseq = hseq.upper().replace("-", "").replace(" ", "")
    for pat in [r"C([A-Z]{3,30}?)WG[A-Z]G", r"C([A-Z]{3,35}?)WG"]:
        m = list(re.finditer(pat, hseq))
        if m:
            return m[-1].group(1)
    return ""


def extract_cdr3_from_hseq(hseq_series: pd.Series, verbose: bool = True) -> pd.Series:
    n = len(hseq_series)
    cdr3s = []
    n_anarci = n_regex = n_fail = 0

    if verbose:
        print(f"\n  [CDR3] CDR3 column not found — extracting from {n:,} HSEQ sequences ...")
        print(f"  [CDR3] Method: ANARCI (IMGT 105-117) with regex fallback")

    for hseq in hseq_series:
        if not isinstance(hseq, str) or len(hseq.strip()) < 20:
            cdr3s.append("")
            n_fail += 1
            continue
        cdr3 = _extract_cdr3_single_anarci(hseq)
        if cdr3:
            n_anarci += 1
        else:
            cdr3 = _extract_cdr3_regex(hseq)
            if cdr3:
                n_regex += 1
            else:
                n_fail += 1
        cdr3s.append(cdr3)

    if verbose:
        print(f"  [CDR3] Done: ANARCI={n_anarci:,}  regex={n_regex:,}  failed={n_fail:,}")
        if n_fail > 0:
            pct = 100 * n_fail / max(n, 1)
            print(f"  [CDR3] WARNING: {n_fail:,} sequences ({pct:.1f}%) could not be extracted — CDR3 set to ''.")
        if n_regex > 0:
            print(f"  [CDR3] NOTE: {n_regex:,} sequences used regex fallback.")

    return pd.Series(cdr3s, index=hseq_series.index, name="CDR3")


def _ensure_cdr3(df: pd.DataFrame, verbose: bool = True) -> pd.DataFrame:
    if "HSEQ" not in df.columns:
        if verbose:
            print("  [CDR3] WARNING: No HSEQ column — cannot extract CDR3.")
        return df

    df = df.copy()
    _nan_like = {"nan", "none", "null", "n/a", "na", ""}

    if "CDR3" not in df.columns:
        if verbose:
            print("  [CDR3] 'CDR3' column not found.")
        df["CDR3"] = ""
    else:
        df["CDR3"] = df["CDR3"].fillna("").astype(str)
        df.loc[df["CDR3"].str.strip().str.lower().isin(_nan_like), "CDR3"] = ""

    empty_mask = df["CDR3"].str.len() == 0
    n_missing  = empty_mask.sum()

    if n_missing > 0:
        if verbose and n_missing < len(df):
            print(f"  [CDR3] {n_missing:,} rows have empty CDR3 — extracting from HSEQ.")
        extracted = extract_cdr3_from_hseq(df.loc[empty_mask, "HSEQ"], verbose=verbose)
        df.loc[empty_mask, "CDR3"] = extracted.values

    return df


def fix_cdr3_c_prefix(df: pd.DataFrame, cdr3_col: str = "CDR3",
                      verbose: bool = True) -> pd.DataFrame:
    if cdr3_col not in df.columns:
        return df

    df   = df.copy()
    mask = (df[cdr3_col].str.len() > 0) & (~df[cdr3_col].str.startswith("C"))
    n    = mask.sum()

    if n > 0 and verbose:
        print(f"\n  [CDR3-C] NOTICE: {n:,} CDR3 sequence(s) do not start with 'C'.")
        print(f"  [CDR3-C] Prepending 'C' to maintain consistency with training data.")
        print(f"  [CDR3-C] (Normal when CDR3 was extracted without the conserved N-terminal cysteine.)")
        shown = 0
        for bc, row in df[mask].iterrows():
            old = row[cdr3_col]
            print(f"    BARCODE={bc}  '{old}'  ->  'C{old}'")
            shown += 1
            if shown >= 5:
                break
        if n > 5:
            print(f"    ... and {n - 5:,} more")

    df.loc[mask, cdr3_col] = "C" + df.loc[mask, cdr3_col]
    return df


# ===========================================================================
# FILE I/O HELPERS
# ===========================================================================

def save_dataframe(df: pd.DataFrame, path: str) -> None:
    path = str(path)
    os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
    if path.lower().endswith('.csv'):
        df.to_csv(path, index=False)
    else:
        df.to_excel(path, index=False)
    print(f"  Saved: {path}  ({len(df):,} rows)")


def read_dataframe(path: str) -> pd.DataFrame:
    if path.lower().endswith('.csv'):
        return pd.read_csv(path)
    return pd.read_excel(path)


def _split_embedding(emb_path, train_barcodes, val_barcodes, train_out, val_out):
    emb = pd.read_csv(emb_path, index_col=0)
    emb.index = emb.index.astype(str).str.strip()
    db_base = str(Path(emb_path).name)
    _parts  = db_base.split(".")
    emb_idx = next((i for i, p in enumerate(_parts) if p == "emb"), None)
    if emb_idx is None:
        print(f"  [emb-split] Cannot parse LM tag from {db_base} — skipping")
        return
    lm_tag = ".".join(_parts[emb_idx-1:])
    tr_out = f"{train_out}.{lm_tag}"
    va_out = f"{val_out}.{lm_tag}"
    tr_bcs = [b for b in train_barcodes.astype(str) if b in emb.index]
    va_bcs = [b for b in val_barcodes.astype(str)   if b in emb.index]
    n_missing_tr = len(train_barcodes) - len(tr_bcs)
    n_missing_va = len(val_barcodes)   - len(va_bcs)
    if n_missing_tr or n_missing_va:
        print(f"  [emb-split] WARNING: {n_missing_tr} train + {n_missing_va} val BARCODEs not found — skipped")
    emb.loc[tr_bcs].to_csv(tr_out)
    emb.loc[va_bcs].to_csv(va_out)
    print(f"  [emb-split] {os.path.basename(tr_out)}  ({len(tr_bcs):,} rows)")
    print(f"  [emb-split] {os.path.basename(va_out)}  ({len(va_bcs):,} rows)")


def _find_embedding_files(db_path: str) -> list:
    import glob
    return sorted(glob.glob(f"{db_path}.*.emb.csv"))


def _align_embedding(df: pd.DataFrame, embedding: pd.DataFrame, context: str = "") -> tuple:
    tag = f"[{context}] " if context else ""
    df = df.copy()
    df.index = df.index.astype(str).str.strip()
    embedding = embedding.copy()
    embedding.index      = embedding.index.astype(str).str.strip()
    embedding.index.name = "BARCODE"
    _sample = embedding.index[:5].tolist()
    if all(str(v).isdigit() for v in _sample):
        print(f"  {tag}WARNING: embedding index looks numeric ({_sample}) — expected BARCODE strings.")
    merged = df.join(embedding, how="inner")
    n_missing = len(df) - len(merged)
    if len(merged) == 0:
        raise ValueError(
            f"{tag}No overlapping BARCODEs between data ({len(df):,} rows) "
            f"and embedding ({len(embedding):,} rows)."
        )
    if n_missing > 0:
        print(f"  {tag}WARNING: {n_missing:,} / {len(df):,} rows ({n_missing/len(df):.1%}) have no embedding — excluded.")
    emb_cols = embedding.columns.tolist()
    return merged[emb_cols], merged.drop(columns=emb_cols)


def split_and_save(db_path: str, split: float = 0.8,
                   cluster_thresh: float = 0.8,
                   cluster_col: str = "CDR3",
                   label_col: str = "psr_filter") -> tuple:
    from sklearn.model_selection import StratifiedGroupKFold
    from utils.clustering import greedy_clustering_by_levenshtein

    df = read_dataframe(db_path)
    p  = Path(db_path)
    print(f"\n[split] {p.name}  ({len(df):,} rows)  split={split:.0%}/{1-split:.0%}  cluster_col={cluster_col}")

    _cc = cluster_col.upper()
    if _cc == 'HSEQ':  _cc = 'VH'
    if _cc == 'HVHVL': _cc = 'VHVL'

    col_map = {'CDR3': f'HCDR3_CLUSTER_{cluster_thresh}',
               'VH':   f'VH_CLUSTER_{cluster_thresh}',
               'VHVL': f'VHVL_CLUSTER_{cluster_thresh}'}
    seq_map = {'CDR3': 'CDR3', 'VH': 'HSEQ', 'VHVL': 'HSEQ'}
    grp_col = col_map.get(_cc, f'HCDR3_CLUSTER_{cluster_thresh}')
    seq_col = seq_map.get(_cc, 'CDR3')

    if grp_col not in df.columns:
        if seq_col not in df.columns:
            raise ValueError(f"Sequence column '{seq_col}' not found.")
        print(f"[split] Computing {grp_col} (threshold={cluster_thresh}) ...")
        seqs = (df['HSEQ'].fillna('').astype(str) + '_' +
                df['LSEQ'].fillna('').astype(str)).tolist() if _cc == 'VHVL' \
               else df[seq_col].fillna('').astype(str).tolist()
        df[grp_col] = greedy_clustering_by_levenshtein(seqs, cluster_thresh)
        n_clust = df[grp_col].nunique()
        print(f"[split] {n_clust:,} clusters  (mean {len(df)/n_clust:.1f} sequences/cluster)")
        try:
            _df_save = read_dataframe(db_path)
            _df_save[grp_col] = df[grp_col].values
            save_dataframe(_df_save, db_path)
            print(f"[split] Saved {grp_col} → {os.path.basename(db_path)}")
        except Exception as _e:
            print(f"[split] WARNING: could not save {grp_col} back to {os.path.basename(db_path)}: {_e}")
    else:
        print(f"[split] Using existing {grp_col}  ({df[grp_col].nunique():,} clusters)")

    _lc     = label_col if label_col in df.columns else None
    has_lbl = _lc is not None

    if has_lbl:
        _lbl_mask = df[_lc].notna()
        n_unl     = (~_lbl_mask).sum()
        if n_unl > 0:
            print(f"[split] {_lbl_mask.sum():,} labelled + {n_unl:,} unlabelled rows — ALL preserved")
    else:
        print(f"[split] WARNING: '{label_col}' not found — splitting by cluster only")
        _lbl_mask = pd.Series([True] * len(df), index=df.index)

    df_lbl  = df[_lbl_mask].copy()
    y_arr   = df_lbl[_lc].values.astype(int) if has_lbl else np.zeros(len(df_lbl), int)
    groups  = df_lbl[grp_col].values

    n_splits = max(2, round(1.0 / (1.0 - split)))
    sgkf = StratifiedGroupKFold(n_splits=n_splits, shuffle=True, random_state=42)
    best_split, best_diff = None, float('inf')
    for tr, va in sgkf.split(np.arange(len(y_arr)), y_arr, groups):
        diff = abs(y_arr[va].mean() - y_arr.mean()) if has_lbl else 0.0
        if diff < best_diff: best_diff, best_split = diff, (tr, va)

    tr_idx, va_idx = best_split
    tr_clusters = set(groups[tr_idx])
    va_clusters = set(groups[va_idx])
    leaked = tr_clusters & va_clusters
    print(f"[split] {'[WARN] ' + str(len(leaked)) + ' cluster(s) leaked' if leaked else '✓ No cluster leakage'}")

    if not _lbl_mask.all():
        df_unl = df[~_lbl_mask].copy()
        unl_in_train = df_unl[grp_col].apply(lambda g: g in tr_clusters or g not in va_clusters)
        df_unl_train = df_unl[unl_in_train]
        df_unl_val   = df_unl[~unl_in_train]
    else:
        df_unl_train = df_unl_val = pd.DataFrame(columns=df.columns)

    train_df = pd.concat([df_lbl.iloc[tr_idx], df_unl_train], ignore_index=True)
    val_df   = pd.concat([df_lbl.iloc[va_idx],   df_unl_val], ignore_index=True)

    _pos_tr = f"  pos={y_arr[tr_idx].mean():.1%}" if has_lbl else ""
    _pos_va = f"  pos={y_arr[va_idx].mean():.1%}" if has_lbl else ""
    print(f"[split] Train={len(train_df):,}{_pos_tr}  Val={len(val_df):,}{_pos_va}")

    ext       = p.suffix
    train_out = str(p.parent / f"{p.stem}_train{ext}")
    val_out   = str(p.parent / f"{p.stem}_val{ext}")
    save_dataframe(train_df.reset_index(drop=True), train_out)
    save_dataframe(val_df.reset_index(drop=True),   val_out)

    bc_col = 'BARCODE'
    if bc_col in df.columns:
        train_bcs = train_df[bc_col].astype(str)
        val_bcs   = val_df[bc_col].astype(str)
    else:
        train_bcs = pd.Index(train_df.index.astype(str))
        val_bcs   = pd.Index(val_df.index.astype(str))

    emb_files = _find_embedding_files(db_path)
    if emb_files:
        print(f"\n[split] Found {len(emb_files)} embedding file(s) — splitting by BARCODE ...")
        for emb_path in emb_files:
            lm_name = os.path.basename(emb_path).split('.emb.csv')[0].split('.')[-1]
            print(f"  [emb-split] {lm_name}: {os.path.basename(emb_path)}")
            try:
                _split_embedding(emb_path, train_bcs, val_bcs, train_out, val_out)
            except Exception as e:
                print(f"  [emb-split] WARNING: failed for {lm_name} — {e}")
    else:
        print(f"\n[split] No embedding files found for {p.name}")

    return train_out, val_out, train_df, val_df


# ===========================================================================
# LOAD DATA
# ===========================================================================

def load_data(db_path, lm="antiberta2", label_col="psr_filter"):
    print(f"\nLoading database : {os.path.basename(db_path)}")
    print(f"Target           : {label_col}")
    print(f"Embedding        : {lm}")

    df = pd.read_excel(db_path)
    print(f"Total rows       : {len(df):,}  |  columns: {len(df.columns)}")

    required = ["BARCODE", "HSEQ", "LSEQ", label_col]
    if lm in ("onehot", "onehot_vh"):
        required += ["CDR3"]

    if "HSEQ" in df.columns:
        df = _ensure_cdr3(df, verbose=True)
        df = fix_cdr3_c_prefix(df, verbose=True)

    missing = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns: {missing}")

    df = df.dropna(subset=required).set_index("BARCODE")
    print(f"After dropna     : {len(df):,} rows")

    # Modes that need no PLM embedding — return sequence DataFrame directly.
    # RF FeatureBuilder computes kmer/biophysical features from HSEQ/LSEQ/CDR3.
    _SEQ_ONLY_MODES = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}

    if lm in _SEQ_ONLY_MODES:
        X    = df[["HSEQ", "LSEQ", "CDR3"]]
        y    = df[label_col].values
        data = df.copy()
        if lm not in ("onehot", "onehot_vh"):
            print(f"[load_data] lm='{lm}' — sequence-only mode, no PLM embedding loaded")
    else:
        # PLM embedding (ablang, antiberta2, antiberty, antiberta2-cssp, igbert ...)
        possible = [f"{db_path}.{lm}.emb.csv"]
        emb_file = next((f for f in possible if os.path.exists(f)), None)
        if not emb_file:
            print(f"Embedding not found -> generating {lm}...")
            emb_file = generate_embedding(db_path, lm=lm)
        print(f"Embedding file   : {emb_file}")
        print(f"Embedding size   : {os.path.getsize(emb_file) / 1024 / 1024:.1f} MB")
        embedding = pd.read_csv(emb_file, index_col=0)
        print(f"Embedding shape  : {embedding.shape[0]:,} samples × {embedding.shape[1]:,} dims")
        X, data = _align_embedding(df, embedding, context="load_data")
        y = data[label_col].values

    print(f"Samples loaded   : {len(y):,}")
    print(f"Target stats     : mean={y.mean():.4f}  std={y.std():.4f}  "
          f"min={y.min():.4f}  max={y.max():.4f}")
    return X, data, y





# ── Universal CDR3 mutagenesis ────────────────────────────────────────────────
# Works for ANY model with predict_proba().
# Called from auto_predict() when --mutagenesis is set.

def _run_cdr3_mutagenesis(
        data, model, model_type, lm, db_stem,
        input_file, target, test_target,
        embeddings_fn=None, n_override=None):
    """
    Run in-silico CDR3 mutagenesis for any model type.

    Parameters
    ----------
    data         : full prediction DataFrame (BARCODE index, HSEQ/CDR3 cols)
    model        : loaded model object with predict_proba()
    model_type   : 'rf' | 'xgboost' | 'transformer_lm' | 'transformer_onehot' | 'cnn'
    lm           : lm name used for this prediction
    db_stem      : training DB stem
    input_file   : path to input file (used to name output folder)
    target       : --target column name
    test_target  : --test_target column name (or None)
    embeddings_fn: callable(mutant_df) → embeddings array, or None for seq-only
    """
    import traceback as _tb
    import matplotlib.pyplot as plt
    import matplotlib.colors as _mc
    import numpy as np

    AMINO_ACIDS = 'ACDEFGHIKLMNPQRSTVWY'

    # ── Config ────────────────────────────────────────────────────────────────
    # Read from model config if available, else defaults
    _cfg = getattr(model, 'config', {}) if hasattr(model, 'config') else {}
    _mut_cfg  = _cfg.get('mutagenesis', {})
    _shap_cfg = _cfg.get('shap', {})
    _fmt      = _mut_cfg.get('format',   'tiff')
    _dpi      = _mut_cfg.get('pub_dpi',   300)
    _make_ppt = _mut_cfg.get('make_ppt',  True)
    _max_n    = _mut_cfg.get('max_samples', 50)
    # --mutagenesis N overrides YAML max_samples
    # n_override=0 means scan ALL antibodies
    if n_override is not None:
        _max_n = len(data) if n_override == 0 else int(n_override)

    # Waterfall folder stem = same as SHAP waterfall prefix
    _path     = Path(input_file)
    _stem     = f"{_path.stem}_{target}_{lm}_{model_type}_{db_stem}"
    _mut_dir  = str(_path.parent / f"{_stem}_mutagenesis")
    import os; os.makedirs(_mut_dir, exist_ok=True)

    # Actual label column
    _label_col = test_target if (test_target and test_target in data.columns) \
                 else (target if target in data.columns else None)

    # Limit samples
    _n = min(len(data), _max_n)
    if _n < len(data):
        print(f"[Mutagenesis] Limited to first {_n}/{len(data)} antibodies "
              f"(mutagenesis.max_samples={_max_n})")

    # ── Build PLM embedding function if model needs embeddings ──────────────
    # For PLM-based RF: each mutant needs a new embedding generated on the fly.
    # This is slow (~0.1-0.5s per mutant) but correct.
    _needs_emb = (hasattr(model, 'fb_') and model.fb_ is not None and
                  model.fb_.feat_cfg.get('embedding', False))

    if _needs_emb:
        print(f"\n[Mutagenesis] WARNING: PLM mode '{lm}' is NOT recommended for mutagenesis.")
        print(f"  Reason: PLM embeddings encode the whole sequence holistically.")
        print(f"  A single CDR3 point mutation shifts the {lm} vector by ~0.1%")
        print(f"  → RF score barely changes → all cells show the same score (flat heatmap).")
        print(f"  This is not a bug — it reflects the PLM model's insensitivity to")
        print(f"  single-residue CDR3 mutations at the sequence representation level.")
        print(f"")
        print(f"  Recommended: re-run with a sequence-based --lm for meaningful mutagenesis:")
        print(f"    --lm biophysical   (fastest, charge/pI/R-count driven)")
        print(f"    --lm kmer          (k-mer motif driven)")
        print(f"    --lm onehot_cdr3   (position-specific)")
        print(f"")
        print(f"  Proceeding anyway — heatmap will likely be flat (informational only).")
        try:
            from embedding_generator import EmbeddingGenerator as _EmbGen
            _emb_gen = _EmbGen(lm=lm)

            def _embeddings_fn(df_row):
                """Generate PLM embedding for a single mutant row in-memory."""
                # df_row has mutant HSEQ (with CDR3 spliced in)
                _hseq = str(df_row.iloc[0].get('HSEQ', '') or '')
                _lseq = str(df_row.iloc[0].get('LSEQ', '') or '')
                _bc   = str(df_row.index[0])
                _emb  = _emb_gen.embed_single(_hseq, _lseq, barcode=_bc)
                return _emb.reshape(1, -1).astype('float32')

        except Exception as _ee:
            print(f"[Mutagenesis] WARNING: Could not set up PLM embedder ({_ee})")
            print(f"[Mutagenesis] Trying temp-file approach ...")
            try:
                import tempfile, os as _os, pandas as _pd
                from embedding_generator import generate_embedding as _gen_emb

                def _embeddings_fn(df_row):
                    with tempfile.NamedTemporaryFile(
                            suffix='.xlsx', delete=False) as _tf:
                        _tmp = _tf.name
                    try:
                        df_row.reset_index().to_excel(_tmp, index=False)
                        _emb_csv = _gen_emb(_tmp, lm=lm)
                        _emb_df  = _pd.read_csv(_emb_csv, index_col=0)
                        return _emb_df.values.astype('float32')
                    finally:
                        _os.unlink(_tmp)
                        for _ext in ['.emb.csv', f'.{lm}.emb.csv']:
                            _c = _tmp + _ext
                            if _os.path.exists(_c): _os.unlink(_c)

            except Exception as _ee2:
                print(f"[Mutagenesis] WARNING: Both embedding approaches failed — {_ee2}")
                print(f"[Mutagenesis] All mutant scores will equal WT score.")
                print(f"[Mutagenesis] → Use --lm biophysical or --lm kmer for fast reliable mutagenesis.")
                _embeddings_fn = None
    else:
        _embeddings_fn = None

    print(f"[Mutagenesis] {model_type.upper()} | {lm} | {_n} antibodies | "
          f"output → {_mut_dir}/")

    task = getattr(model, 'task', 'classification')
    _saved = []

    for _s in range(_n):
        _row  = data.iloc[_s]
        _bc   = str(data.index[_s])
        _cdr3 = str(_row.get('CDR3', '') or '').upper().replace('-', '')
        if not _cdr3:
            print(f"  [{_s+1}/{_n}] {_bc}: CDR3 missing — skipped")
            continue

        _n_pos = len(_cdr3)
        _n_aa  = len(AMINO_ACIDS)

        # Actual label for title
        _actual = None
        if _label_col:
            try:
                _v = _row.get(_label_col)
                if _v is not None and str(_v) not in ('', 'nan'):
                    _actual = int(float(_v))
            except Exception:
                pass

        print(f"  [{_s+1}/{_n}] {_bc}  CDR3={_cdr3}  "
              f"({_n_pos}×{_n_aa}={_n_pos*_n_aa} mutants)  "
              + (f"Actual({_label_col})={'PASS' if _actual==1 else 'FAIL'}"
                 if _actual is not None else ""))

        # WT score
        _wt_df = data.iloc[[_s]].copy()
        try:
            _wt_score = float(_predict_single(model, model_type, _wt_df,
                                              _embeddings_fn, task))
        except Exception as _e:
            print(f"    WT score failed: {_e}"); _wt_score = float('nan')

        # Score matrix (n_aa × n_pos)
        _mat = np.full((_n_aa, _n_pos), np.nan, dtype=np.float32)
        _vh  = str(_row.get('HSEQ', '') or '')
        _cdr3_start = _vh.find(_cdr3) if _vh else -1

        _first_err  = True   # print first exception only — avoid log flood
        for _pi in range(_n_pos):
            for _ai, _mut_aa in enumerate(AMINO_ACIDS):
                _mcdr3 = _cdr3[:_pi] + _mut_aa + _cdr3[_pi+1:]
                _mrow  = _row.to_dict()
                _mrow['CDR3'] = _mcdr3
                if _cdr3_start >= 0:
                    _mrow['HSEQ'] = (_vh[:_cdr3_start] + _mcdr3 +
                                     _vh[_cdr3_start + _n_pos:])
                _mut_df = pd.DataFrame([_mrow], index=[_bc])
                try:
                    _mat[_ai, _pi] = float(_predict_single(
                        model, model_type, _mut_df, _embeddings_fn, task))
                except Exception as _me:
                    if _first_err:
                        import traceback as _tb
                        print(f"    [Mutagenesis] mutant scoring failed: {_me}")
                        print(_tb.format_exc())
                        _first_err = False
                    _mat[_ai, _pi] = _wt_score

        # ── Warn if all scores are identical (embedding fallback triggered) ────
        _unique_scores = np.unique(_mat[~np.isnan(_mat)])
        if len(_unique_scores) <= 1:
            print(f"    [Mutagenesis] WARNING: all mutant scores = WT score ({_wt_score:.4f})")
            print(f"    This means PLM re-embedding failed for all mutants.")
            print(f"    → Retrain or use --lm biophysical/kmer for reliable mutagenesis.")

        # ── Plot ─────────────────────────────────────────────────────────────
        _fw = max(9, _n_pos * 0.55 + 3)
        fig, ax = plt.subplots(figsize=(_fw, 7))

        if task == 'classification':
            _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0)
            _cmap = 'RdBu'
        else:
            _vmin = float(np.nanmin(_mat))
            _vmax = float(np.nanmax(_mat))
            _norm = _mc.TwoSlopeNorm(vmin=_vmin,
                                      vcenter=(_vmin+_vmax)/2,
                                      vmax=_vmax)
            _cmap = 'coolwarm'

        im = ax.imshow(_mat, cmap=_cmap, norm=_norm, aspect='auto')

        # Cell annotations
        _fsz_cell = max(4.0, min(7.0, 120.0 / max(_n_pos, 1)))
        for _ai in range(_n_aa):
            for _pi in range(_n_pos):
                _v = _mat[_ai, _pi]
                if np.isnan(_v): continue
                _tc = 'white' if (_v < 0.35 or _v > 0.65) else '#333'
                ax.text(_pi, _ai, f"{_v:.2f}", ha='center', va='center',
                        fontsize=_fsz_cell, color=_tc,
                        fontweight='bold' if abs(_v - 0.5) > 0.30 else 'normal')

        # Box WT residues
        for _pi, _wt in enumerate(_cdr3):
            if _wt in AMINO_ACIDS:
                ax.add_patch(plt.Rectangle(
                    (_pi - 0.5, AMINO_ACIDS.index(_wt) - 0.5), 1, 1,
                    fill=False, edgecolor='black', lw=2.0, zorder=5))

        # Axes
        ax.set_xticks(range(_n_pos))
        ax.set_xticklabels([f"{_cdr3[i]}\n{i+1}" for i in range(_n_pos)],
                           fontsize=8.5)
        ax.set_yticks(range(_n_aa))
        ax.set_yticklabels(list(AMINO_ACIDS), fontsize=8)
        ax.set_xlabel('CDR3 position  (WT residue shown above position number)',
                      fontsize=9)
        ax.set_ylabel('Substituted AA', fontsize=9)

        # Colorbar
        cbar = plt.colorbar(im, ax=ax, shrink=0.85, pad=0.02)
        cbar.set_label('P(PASS)' if task=='classification' else 'Score',
                       fontsize=9, labelpad=6)
        if task == 'classification':
            cbar.set_ticks([0.0, 0.25, 0.5, 0.75, 1.0])
            cbar.set_ticklabels(['0.0\nFAIL', '0.25', '0.50\nborder',
                                 '0.75', '1.0\nPASS'])
        cbar.ax.tick_params(labelsize=7.5)

        # Title
        _score_str = (f"WT P(PASS)={_wt_score:.4f}" if not np.isnan(_wt_score)
                      else "WT score=N/A")
        _act_str   = (f"  |  Actual ({_label_col}) = "
                      f"{'PASS' if _actual==1 else 'FAIL'}"
                      if _actual is not None else "")
        _task_str  = 'CL' if task == 'classification' else 'REG'
        # Sensitivity check — warn if heatmap is flat
        _mat_range = float(np.nanmax(_mat) - np.nanmin(_mat))
        _flat_warn  = ""
        if _mat_range < 0.01 and _needs_emb:
            _flat_warn = "\n⚠ Flat heatmap — PLM embeddings insensitive to single mutations"
            _flat_warn += "\n  Re-run with --lm biophysical for position-specific insights"

        ax.set_title(
            f"IPI MLAbDev · CDR3 Mutagenesis Heatmap\n"
            f"ID: {_bc}   {_score_str}{_act_str}\n"
            f"{model_type.upper()} | {lm} | {_task_str} | {db_stem}"
            f"{_flat_warn}",
            fontsize=9, loc='center', pad=8,
            color='#CC3300' if _flat_warn else 'black'
        )
        plt.tight_layout()

        _bc_safe = _bc.replace('/', '_').replace(' ', '_')
        _img_path = os.path.join(_mut_dir,
                                 f"{_s+1:04d}_{_bc_safe}_cdr3_mutagenesis.{_fmt}")
        _save_kw = dict(dpi=_dpi, bbox_inches='tight')
        if _fmt == 'tiff':    _save_kw['format'] = 'tiff'
        elif _fmt in ('jpeg','jpg'):
            _save_kw['format'] = 'jpeg'
            _save_kw['pil_kwargs'] = {'quality': 95}
        plt.savefig(_img_path, **_save_kw)
        plt.close()
        _saved.append((_bc, _img_path, _wt_score))
        print(f"    → {os.path.basename(_img_path)}")

    print(f"[Mutagenesis] {len(_saved)} heatmaps → {_mut_dir}/")
    if _saved and _needs_emb:
        print(f"\n[Mutagenesis] NOTE: PLM mutagenesis complete.")
        print(f"  If all scores look identical, this is expected for PLM-based RF models.")
        print(f"  For position-specific insights, re-run with --lm biophysical or --lm kmer.")

    # ── PPT ───────────────────────────────────────────────────────────────────
    if _make_ppt and _saved:
        try:
            from pptx import Presentation as _Prs
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = _Prs()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]

            for _bc, _img, _wt in _saved:
                slide = prs.slides.add_slide(blank)
                _iw = Inches(11); _ih = Inches(6.5)
                slide.shapes.add_picture(
                    _img,
                    (prs.slide_width - _iw) / 2, Inches(0.4),
                    width=_iw, height=_ih)
                txb = slide.shapes.add_textbox(
                    Inches(0.15), Inches(7.1),
                    Inches(13.0), Inches(0.35))
                tf  = txb.text_frame
                tf.text = (f"{_bc}  |  {model_type.upper()} | {lm} | "
                           f"{db_stem}  |  "
                           f"WT={'%.3f'%_wt if not np.isnan(_wt) else 'N/A'}")
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.runs[0].font.size = Pt(7)
                p.runs[0].font.color.rgb = RGBColor(0x88, 0x87, 0x80)

            _ppt = os.path.join(_mut_dir, "cdr3_mutagenesis_all.pptx")
            prs.save(_ppt)
            print(f"[Mutagenesis] PPT ({len(_saved)} slides) → {_ppt}")
        except ImportError:
            print("[Mutagenesis] pip install python-pptx for PPT")
        except Exception as _pe:
            print(f"[Mutagenesis] PPT failed — {_pe}")


def _predict_single(model, model_type, df_row, embeddings_fn, task):
    """Route predict_proba to the correct model interface."""
    import numpy as np
    if model_type == 'rf':
        # If model needs PLM embeddings, generate them via embeddings_fn
        _needs_emb = (hasattr(model, 'fb_') and model.fb_ is not None and
                      model.fb_.feat_cfg.get('embedding', False))
        if _needs_emb and embeddings_fn is not None:
            _emb = embeddings_fn(df_row)
        else:
            _emb = None
        return model.predict_proba(df_row, embeddings=_emb)[0]
    elif model_type == 'xgboost':
        _SEQ_ONLY_R = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3",
                       "biophysical","kmer","seq","none"}
        _xgb_row_df = df_row[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in df_row.columns for c in ['HSEQ','LSEQ','CDR3']) else df_row.copy()
        _emb_r = embeddings_fn(df_row) if embeddings_fn else None
        if _emb_r is not None and hasattr(_emb_r, 'reshape'):
            _emb_r = _emb_r.reshape(1, -1)
        if getattr(model, 'task', 'classification') == 'regression':
            return float(model.predict(_xgb_row_df, embeddings=_emb_r)[0])
        return model.predict_proba(_xgb_row_df, embeddings=_emb_r)[0]
    elif model_type in ('transformer_onehot',):
        return model.predict_proba(df_row)[0]
    elif model_type in ('transformer_lm', 'cnn'):
        _emb = embeddings_fn(df_row) if embeddings_fn else None
        if _emb is not None:
            return model.predict_proba(_emb)[0]
        return model.predict_proba(df_row)[0]
    else:
        return model.predict_proba(df_row)[0]



def _find_matching_checkpoints(model_dir: str, target: str, lm: str,
                                model_type: str, db_stem: str,
                                ext: str) -> list:
    """
    Search MODEL_DIR for checkpoints that match target + lm + model_type + db_stem.
    Returns a sorted list of matching paths (most recently modified first).

    Matches both exact and suffix-extended filenames:
      FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset.pt        (Mode 1)
      FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_plmft.pt  (Mode 2)
      FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_lora8.pt  (Mode 3)
      FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_ft_*.pt   (fine-tuned)
    """
    import glob, os
    from pathlib import Path

    # Base prefix that all variants share
    _prefix = f"FINAL_{target}_{lm}_{model_type}_{db_stem}"

    # Search for all files starting with prefix
    _pattern = os.path.join(model_dir, f"{_prefix}*{ext}")
    _found   = glob.glob(_pattern)

    # Sort by modification time — most recent first
    _found.sort(key=lambda p: os.path.getmtime(p), reverse=True)
    return _found

def auto_predict(input_file, target="sec_filter", lm="antiberta2",
                 model_type="xgboost", db_path=None, test_target=None,
                 run_mutagenesis=False, mutagenesis_n=None, threshold=None,
                 model_path=None, **kwargs):
    kwargs['model_path'] = model_path   # pass to inner checkpoint lookup
    print(f"\nPREDICTING: {os.path.basename(input_file)}")
    print(f"Target: {target.upper()} | Model: {model_type.upper()} | LM: {lm}")
    if db_path:
        print(f"Using model trained on: {os.path.basename(db_path)}")

    if input_file.lower().endswith((".xlsx", ".xls")):
        data = pd.read_excel(input_file)
    elif input_file.lower().endswith(".csv"):
        data = pd.read_csv(input_file)
    else:
        raise ValueError(f"Unsupported file format: {input_file}")

    if "BARCODE" not in data.columns:
        data["BARCODE"] = range(len(data))
    data["BARCODE"] = data["BARCODE"].astype(str).str.strip()
    data = data.set_index("BARCODE")

    if "HSEQ" in data.columns:
        data = data.reset_index()
        data = _ensure_cdr3(data, verbose=True)
        data = fix_cdr3_c_prefix(data, verbose=True)
        data = data.set_index("BARCODE")

    _SEQ_ONLY = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}

    if lm in _SEQ_ONLY and lm not in ("onehot", "onehot_vh"):
        # seq-only RF mode — no embedding file needed
        print(f"[predict] lm='{lm}' — sequence-only mode, no PLM embedding loaded")
        X = data[["HSEQ", "LSEQ", "CDR3"]] if all(
            c in data.columns for c in ["HSEQ", "LSEQ", "CDR3"]) else data
    elif lm not in ["onehot", "onehot_vh"]:
        emb_file = f"{input_file}.{lm}.emb.csv"
        if not os.path.exists(emb_file):
            print("Generating embedding from input file...")
            generate_embedding(input_file, lm=lm)
        print(f"Using embedding: {emb_file}")
        embedding = pd.read_csv(emb_file, index_col=0)
        X, data = _align_embedding(data, embedding, context="predict")
    else:
        required = ["HSEQ", "LSEQ", "CDR3"]
        missing  = [c for c in required if c not in data.columns]
        if missing:
            raise ValueError(f"Missing columns for one-hot: {missing}")

        _nan_like = {"nan", "none", "null", "n/a", "na", "#n/a", "#value!", "#ref!"}
        for _col in required:
            data[_col] = data[_col].fillna("").astype(str)
            _bad = data[_col].str.strip().str.lower().isin(_nan_like)
            if _bad.any():
                print(f"  [fix-nan] {_col}: {_bad.sum():,} nan-like -> ''")
                data.loc[_bad, _col] = ""

        _empty_lseq = data["LSEQ"].str.len() == 0
        if _empty_lseq.any():
            print(f"  [warn] {_empty_lseq.sum():,} antibodies have empty LSEQ ({_empty_lseq.mean():.1%}) — encoded with VL=zeros.")

        _no_hseq = data["HSEQ"].str.strip().str.len() < 5
        if _no_hseq.any():
            print(f"  [fix-nan] Dropping {_no_hseq.sum():,} rows with missing HSEQ")
            data = data[~_no_hseq]
        if len(data) == 0:
            raise ValueError("No valid sequences after cleaning")

        X = data[required]

    # ── Normalise LM name (consistent filename convention) ─────────────────
    # onehot_cdr3 is accepted as CLI alias for onehot_hcdr3
    _LM_ALIASES = {"onehot_cdr3": "onehot_hcdr3"}
    lm = _LM_ALIASES.get(lm, lm)

    # db_stem: derive from db_path OR from model_path checkpoint name
    if db_path:
        db_stem = Path(db_path).stem
    elif kwargs.get('model_path'):
        # Extract db portion from checkpoint name
        import re as _re2
        _ckpt_stem = Path(kwargs['model_path']).stem
        _m2 = _re2.search(
            r'(?:transformer_lm|transformer_onehot|rf|xgboost|cnn)_(.+?)(?:_ft_|_lora|_plmft|$)',
            _ckpt_stem)
        db_stem = _m2.group(1) if _m2 else _ckpt_stem
    else:
        db_stem = "default"
    ext        = ".pt" if model_type in ["cnn", "transformer_onehot", "transformer_lm"] else ".pkl"
    # For kmer/onehot RF: add chain tag to model path
    _chain_tag = ""
    if model_type == "rf":
        _SEQ_MODES = {"onehot", "kmer"}
        if lm in _SEQ_MODES or lm in ("biophysical", "none", "seq"):
            try:
                from models.randomforest import RandomForestModel as _RFM_tmp
                _tmp_cfg   = _RFM_tmp("config/random_forest.yaml").config
                _tmp_fb    = __import__('models.randomforest', fromlist=['FeatureBuilder']).FeatureBuilder(_tmp_cfg)
                _chain_tag = _tmp_fb.chain_tag
                if _chain_tag:
                    _chain_tag = f"_{_chain_tag}"
            except Exception:
                _chain_tag = ""
    # --model_path explicitly overrides auto-discovery
    _explicit_path = kwargs.get('model_path', None)
    if _explicit_path:
        model_path = _explicit_path
        print(f"[load] Using explicit model path: {model_path}")
    else:
        # Try regression variant first, then classification (no suffix)
        _base_path = f"{MODEL_DIR}/FINAL_{target}_{lm}{_chain_tag}_{model_type}_{db_stem}"
        if os.path.exists(f"{_base_path}_regression{ext}"):
            model_path = f"{_base_path}_regression{ext}"
            print(f"[load] Regression checkpoint: {Path(model_path).name}")
        else:
            model_path = f"{_base_path}{ext}"
        # Auto-fallback: if exact path not found, search for any matching checkpoint
        if not os.path.exists(model_path):
            _candidates = _find_matching_checkpoints(
                MODEL_DIR, target, lm, model_type, db_stem, ext)
            if _candidates:
                print(f"\n[load] Exact checkpoint not found: {Path(model_path).name}")
                print(f"[load] Found {len(_candidates)} matching checkpoint(s):")
                for i, c in enumerate(_candidates):
                    print(f"  [{i}] {Path(c).name}")
                print(f"[load] Using: [{0}] {Path(_candidates[0]).name}")
                print(f"[load] To use a different one: --model_path <path>")
                model_path = _candidates[0]
            else:
                raise FileNotFoundError(
                    f"Model not found: {model_path}\n"
                    f"Searched in: {MODEL_DIR}\n"
                    f"Run --train first, or use --model_path to specify explicitly."
                )

    if not os.path.exists(model_path):
        raise FileNotFoundError(f"Model not found: {model_path}")

    if model_type == "xgboost":
        model = XGBoostModel.load(model_path)
    elif model_type == "rf":
        model = RandomForestModel.load(model_path)
    elif model_type == "cnn":
        model = CNNModel.load(model_path,
                              embedding_dim=X.shape[1] if lm not in ["onehot", "onehot_vh", "onehot_hcdr3", "onehot_cdr3"] else None)
    elif model_type == "transformer_onehot":
        # transformer_onehot architecture is fixed:
        #   branch-1 = VH+VL (lm_mode='onehot') or VH (lm_mode='onehot_vh')
        #   branch-2 = CDR3  always
        # --lm onehot_cdr3/hcdr3 not applicable — no change to lm_mode
        model = TransformerOneHotModel.load(model_path)
    elif model_type == "transformer_lm":
        model = TransformerLMModel.load(model_path, embedding_dim=X.shape[1])

    # ── [FIX] RF uses (X_df, embeddings=array); others use raw array or DataFrame ──
    if model_type == "rf":
        _SEQ_ONLY_P = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}
        _rf_X_df_p  = data[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
        _rf_emb_p   = None if lm in _SEQ_ONLY_P else (
            X.values if hasattr(X, 'values') else X)
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(_rf_X_df_p, embeddings=_rf_emb_p)
        else:
            scores = model.predict_proba(_rf_X_df_p, embeddings=_rf_emb_p)
    elif model_type == "xgboost":
        _SEQ_ONLY_P = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3",
                       "biophysical","kmer","seq","none"}
        _xgb_X_df_p = data[['HSEQ','LSEQ','CDR3']].copy() if all(
            c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
        _xgb_emb_p  = None if lm in _SEQ_ONLY_P else (
            X.values if hasattr(X, 'values') else X)
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(_xgb_X_df_p, embeddings=_xgb_emb_p)
        else:
            scores = model.predict_proba(_xgb_X_df_p, embeddings=_xgb_emb_p)
    else:
        if getattr(model, 'task', 'classification') == 'regression':
            scores = model.predict(X)
        else:
            scores = model.predict_proba(X)

    # ── Threshold + labels ────────────────────────────────────────────────────
    _is_regression = getattr(model, 'task', 'classification') == 'regression'
    _reg_tag       = "_regression" if _is_regression else ""
    if _is_regression:
        # Regression: scores are continuous — no binary thresholding
        labels = scores
        print(f"[predict] regression mode — continuous scores  "
              f"(min={scores.min():.4f}  max={scores.max():.4f}  "
              f"mean={scores.mean():.4f})")
    elif threshold is not None:
        _thresh = float(threshold)
        labels  = (scores >= _thresh).astype(int)
        print(f"[predict] threshold={_thresh:.4f}  (--threshold override)")
    else:
        _thresh = getattr(model, "recommended_threshold", 0.5)
        labels  = (scores >= _thresh).astype(int)
        print(f"[predict] threshold={_thresh:.4f}  "
              f"({'embedded from kfold' if _thresh != 0.5 else 'default 0.5'})")

    data[f"{model_type}_{lm}_{db_stem}_score"] = scores
    data[f"{model_type}_{lm}_{db_stem}_label"] = labels

    if model_type in ("transformer_lm", "cnn") and hasattr(model, 'save_hidden_states'):
        try:
            _hidden_stem = str(Path(input_file).with_name(
                f"{Path(input_file).stem}_pred_{target}_{lm}_{model_type}_{db_stem}_hidden.emb"))
            _hidden_csv = f"{_hidden_stem}.csv"
            _X_for_hidden = X if hasattr(X, 'index') else pd.DataFrame(X, index=data.index)
            model.save_hidden_states(_X_for_hidden, _hidden_csv)
            print(f"[predict] Hidden states → {Path(_hidden_csv).name}")
        except Exception as _he:
            print(f"[predict] WARNING: hidden state extraction failed — {_he}")

    # ── RF + XGBoost SHAP on predict set ────────────────────────────────────
    if model_type in ("rf", "xgboost") and hasattr(model, 'shap_analysis'):
        # Always reload SHAP config from current YAML
        try:
            import yaml as _yaml
            _yaml_path = ("config/xgboost.yaml" if model_type == "xgboost"
                          else "config/random_forest.yaml")
            if os.path.exists(_yaml_path):
                with open(_yaml_path) as _yf:
                    _yaml_shap = (_yaml.safe_load(_yf) or {}).get('shap', {})
                if _yaml_shap:
                    model.config['shap'] = _yaml_shap
        except Exception:
            pass
        _shap_cfg = model.config.get('shap', {})
        # Default enabled=True so predict always runs SHAP unless explicitly disabled
        _shap_ok  = _shap_cfg.get('enabled', True)
        if not _shap_ok:
            print("[SHAP] Skipped — shap.enabled=false in config/random_forest.yaml")
        else:
            try:
                _SEQ_ONLY_SH = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}
                _shap_emb    = None if lm in _SEQ_ONLY_SH else (
                    X.values if hasattr(X, 'values') else X)
                _shap_X_df   = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
                _shap_top    = _shap_cfg.get('top_features', 30)
                _shap_prefix = f"{Path(input_file).stem}_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"

                print(f"\n[SHAP] Running on PREDICT set  n={len(_shap_X_df):,}  top={_shap_top}")
                print(f"[SHAP] fb_={'ready' if model.fb_ is not None else 'MISSING — model may need retraining'}")
                print(f"[SHAP] Output dir → {Path(input_file).parent}")

                if model.fb_ is None:
                    print("[SHAP] ERROR: FeatureBuilder not found in loaded model.")
                    print("       Retrain the model with the updated randomforest.py to embed fb_.")
                else:
                    # Override MODEL_DIR in the model's module so outputs
                    # go to the predict file's directory (not MODEL_DIR)
                    if model_type == "xgboost":
                        import models.xgboost as _shap_mod
                    else:
                        import models.randomforest as _shap_mod
                    _rfmod    = _shap_mod
                    _orig_mdir = _shap_mod.MODEL_DIR
                    _shap_mod.MODEL_DIR = str(Path(input_file).parent)
                    # Pass actual labels if target column exists in input file
                    # Also check test_target (--test_target flag)
                    _shap_actual = None
                    _label_col   = test_target if test_target and test_target in data.columns                                    else (target if target in data.columns else None)
                    if _label_col:
                        try:
                            _shap_actual = []
                            for v in data[_label_col].values:
                                if isinstance(v, float) and v != v:   # NaN
                                    _shap_actual.append(None)
                                else:
                                    try:    _shap_actual.append(int(v))
                                    except: _shap_actual.append(None)
                            print(f"[SHAP] Actual labels from '{_label_col}' "
                                  f"(shown in waterfall title)")
                        except Exception:
                            _shap_actual = None

                    # actual_col_name: --test_target takes priority, else --target
                    _shap_col = test_target if test_target else target
                    try:
                        model.shap_analysis(
                            _shap_X_df, _shap_emb,
                            output_prefix   = _shap_prefix,
                            split_tag       = "predict",
                            top_n           = _shap_top,
                            barcodes        = list(_shap_X_df.index.astype(str)),
                            actual_labels   = _shap_actual,
                            actual_col_name = _shap_col,
                            lm_name         = lm,
                            db_name         = db_stem,
                        )
                    finally:
                        _shap_mod.MODEL_DIR = _orig_mdir
            except Exception as _se:
                import traceback
                print(f"[SHAP] predict SHAP failed: {_se}")
                print(traceback.format_exc())

    # ── CDR3 Mutagenesis — universal (all model types) ───────────────────────
    if run_mutagenesis:
        try:
            # PLM modes: warn about speed but proceed
            _PLM_MODES = {"ablang","antiberta2","antiberta2-cssp","antiberty","igbert"}
            if lm in _PLM_MODES:
                print(f"\n[Mutagenesis] WARNING: PLM mode '{lm}' requires re-embedding "
                      f"each mutant — this may be slow for large CDR3s.")
                print(f"  Each antibody: {max(1,len(data.iloc[0].get('CDR3','') or ''))} positions "
                      f"× 20 AA = embeddings per antibody.")
                print(f"  Consider --lm biophysical or --lm kmer for faster mutagenesis.\n")
            _run_cdr3_mutagenesis(
                data         = data,
                model        = model,
                model_type   = model_type,
                lm           = lm,
                db_stem      = db_stem,
                input_file   = input_file,
                target       = target,
                test_target  = test_target,
                embeddings_fn= None,
                n_override   = mutagenesis_n,
            )
        except Exception as _me:
            import traceback
            print(f"[Mutagenesis] failed — {_me}")
            print(traceback.format_exc())

    path        = Path(input_file)
    output_file = path.with_name(
        f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}{path.suffix}")
    if path.suffix.lower() in [".xlsx", ".xls"]:
        data.reset_index().to_excel(output_file, index=False)
    else:
        data.reset_index().to_csv(output_file, index=False)

    print(f"Saved predictions to: {output_file}")
    if _is_regression:
        import numpy as _np_pr
        print(f"Score stats: mean={_np_pr.mean(scores):.4f}  "
              f"std={_np_pr.std(scores):.4f}  "
              f"min={_np_pr.min(scores):.4f}  max={_np_pr.max(scores):.4f}")
    else:
        print(f"Positive rate: {labels.mean():.1%}")

    # Evaluation — use test_target if given, else fall back to target
    _eval_label_col  = test_target if (test_target and test_target != target) else target
    _data_with_index = data.reset_index()
    _eval_col_found  = _eval_label_col in _data_with_index.columns
    if not _eval_col_found and target in _data_with_index.columns:
        # Fall back to --target if --test_target column not found
        _eval_label_col = target
        _eval_col_found = True
    if _eval_col_found:
        print(f"\n[eval] Ground-truth column '{_eval_label_col}' found — running evaluation...")
        if _is_regression:
            # Regression evaluation: Spearman / Pearson / R² / MAE
            try:
                import numpy as _enp
                from scipy.stats import spearmanr as _sp, pearsonr as _pe
                from sklearn.metrics import r2_score as _r2, mean_absolute_error as _mae
                _data_eval  = data.reset_index()
                _true_vals  = _data_eval[_eval_label_col].dropna().values.astype(float)
                _score_col_r = f"{model_type}_{lm}_{db_stem}_score"
                _pred_vals  = _data_eval.loc[_data_eval[_eval_label_col].notna(),
                                              _score_col_r].values.astype(float)
                if len(_true_vals) > 2:
                    _rho   = _sp(_true_vals, _pred_vals)[0]
                    _pear  = _pe(_true_vals, _pred_vals)[0]
                    _r2v   = _r2(_true_vals, _pred_vals)
                    _maev  = _mae(_true_vals, _pred_vals)
                    print(f"[eval] Regression metrics  n={len(_true_vals)}")
                    print(f"  Spearman ρ = {_rho:.4f}")
                    print(f"  Pearson  r = {_pear:.4f}")
                    print(f"  R²         = {_r2v:.4f}")
                    print(f"  MAE        = {_maev:.4f}")
                else:
                    print("[eval] Too few samples for regression metrics.")
            except Exception as _e:
                print(f"[eval] Regression eval failed — {_e}")
        else:
            try:
                from utils.evaluate_model import evaluate
                _score_col = f"{model_type}_{lm}_{db_stem}_score"
                _eval_stem = str(path.with_name(f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"))
                evaluate(file=str(output_file), target=_eval_label_col, score_col=_score_col,
                         cost_fp=1.0, cost_fn=3.0, out=_eval_stem, test_target=test_target,
                         model_type=model_type, lm=lm, db_stem=db_stem,
                         dataset_name=Path(input_file).stem)
            except ImportError:
                print("[eval] utils/evaluate_model.py not found — skipping.")
            except Exception as _e:
                import traceback as _etb
                print(f"[eval] WARNING: evaluation failed — {_e}")
                print(_etb.format_exc())
    else:
        print(f"\n[eval] Skipped — neither '{_eval_label_col}' nor '{target}' column found in predict file.")
        print(f"  Add --test_target with a valid column name for evaluation metrics.")

    # ── Biophysical property panels — classification only ────────────────────
    if _is_regression:
        print("\n[biophys] Skipped — regression mode (no PASS/FAIL labels)")
    else:
        try:
            from utils.plot_biophysical import plot_biophysical_report
            _bio_stem = str(path.with_name(
                f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"))
            print(f"\n[biophys] Generating biophysical property panels...")
            plot_biophysical_report(file=str(output_file), target=target,
                                    test_target=_eval_label_col if _eval_col_found else None,
                                    out=_bio_stem,
                                    dataset_name=Path(input_file).stem)
        except ImportError:
            print("[biophys] utils/plot_biophysical.py not found — skipping.")
        except Exception as _e:
            import traceback as _btb
            print(f"[biophys] WARNING: biophysical plot failed — {_e}")
            print(_btb.format_exc())

    # ── Summary of output files ───────────────────────────────────────────────
    _out_stem = str(path.with_name(f"{path.stem}_pred_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}"))
    print(f"\n{'─'*60}")
    print(f"  OUTPUT FILES  ({Path(input_file).name})")
    print(f"{'─'*60}")
    print(f"  Predictions     : {Path(output_file).name}")
    if _eval_col_found:
        if _is_regression:
            print(f"  Regression eval : Spearman ρ / Pearson r / R² / MAE (in log)")
        else:
            print(f"  ROC curve       : {Path(_out_stem).name}_roc_{target}.tiff")
            print(f"  KDE plot        : {Path(_out_stem).name}_kde_{target}.tiff")
            print(f"  Histogram       : {Path(_out_stem).name}_histogram_{target}.tiff")
            print(f"  Eval summary    : {Path(_out_stem).name}_eval_{target}.xlsx")
            print(f"  Biophys report  : {Path(_out_stem).name}_biophysical_*.tiff")
    _wf_dir = Path(output_file).parent / f"{path.stem}_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}_predict_shap_waterfalls"
    _mut_dir_out = Path(output_file).parent / f"{path.stem}_{target}_{lm}_{model_type}_{db_stem}{_reg_tag}_mutagenesis"
    print(f"  SHAP beeswarm   : {Path(_out_stem).name}_predict_shap_beeswarm.png")
    print(f"  SHAP waterfall  : {_wf_dir.name}/ + .pptx")
    if run_mutagenesis:
        print(f"  Mutagenesis     : {_mut_dir_out.name}/ + .pptx")
    print(f"{'─'*60}")


_ALL_PLM_LMS = ["ablang", "antiberty", "antiberta2", "antiberta2-cssp", "igbert"]


def auto_predict_multi_lm(input_file, target="psr_filter",
                           lms=None, lm_tag="all",
                           model_type="transformer_lm",
                           db_path=None, test_target=None):
    if lms is None:
        lms = _ALL_PLM_LMS

    print(f"\n{'═'*62}")
    print(f"  MULTI-LM PREDICTION")
    print(f"  File   : {os.path.basename(input_file)}")
    print(f"  Target : {target.upper()}  |  Model: {model_type.upper()}")
    print(f"  LMs    : {lms}")
    if db_path:
        print(f"  DB     : {os.path.basename(db_path)}")
    print(f"{'─'*62}")

    if input_file.lower().endswith((".xlsx", ".xls")):
        base_data = pd.read_excel(input_file)
    elif input_file.lower().endswith(".csv"):
        base_data = pd.read_csv(input_file)
    else:
        raise ValueError(f"Unsupported file format: {input_file}")

    if "BARCODE" not in base_data.columns:
        base_data["BARCODE"] = range(len(base_data))
    base_data["BARCODE"] = base_data["BARCODE"].astype(str).str.strip()

    if "HSEQ" in base_data.columns:
        base_data = _ensure_cdr3(base_data, verbose=True)
        base_data = fix_cdr3_c_prefix(base_data, verbose=True)

    base_data = base_data.set_index("BARCODE")

    db_stem = Path(db_path).stem if db_path else "default"
    ext     = ".pt" if model_type in ["cnn", "transformer_onehot", "transformer_lm"] else ".pkl"

    results_summary = []
    failed_lms      = []

    for lm in lms:
        print(f"\n  ── LM: {lm} ──")
        try:
            emb_file = f"{input_file}.{lm}.emb.csv"
            if not os.path.exists(emb_file):
                print(f"  Generating {lm} embedding ...")
                generate_embedding(input_file, lm=lm)
            embedding = pd.read_csv(emb_file, index_col=0)
            X, _lm_data = _align_embedding(base_data, embedding, context=f"predict/{lm}")

            _base2     = f"{MODEL_DIR}/FINAL_{target}_{lm}_{model_type}_{db_stem}"
            if os.path.exists(f"{_base2}_regression{ext}"):
                model_path = f"{_base2}_regression{ext}"
            else:
                model_path = f"{_base2}{ext}"
            if not os.path.exists(model_path):
                raise FileNotFoundError(f"Model not found: {model_path}")

            if model_type == "xgboost":
                model = XGBoostModel.load(model_path)
            elif model_type == "rf":
                model = RandomForestModel.load(model_path)
            elif model_type == "cnn":
                model = CNNModel.load(model_path, embedding_dim=X.shape[1])
            elif model_type == "transformer_lm":
                model = TransformerLMModel.load(model_path, embedding_dim=X.shape[1])

            # ── [FIX] RF uses (X_df, embeddings=array) ──────────────────────
            if model_type == "rf":
                _rf_X_df_m = base_data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in base_data.columns for c in ['HSEQ','LSEQ','CDR3']) else base_data.copy()
                _rf_X_df_m = _rf_X_df_m.loc[_lm_data.index]
                _SEQ_ONLY_ML = {"biophysical", "kmer", "onehot", "onehot_vh",
                                  "onehot_cdr3", "onehot_hcdr3", "none", "seq"}
                _rf_emb_m  = None if lm in _SEQ_ONLY_ML                              else (X.values if hasattr(X, 'values') else X)
                scores = model.predict_proba(_rf_X_df_m, embeddings=_rf_emb_m)
            elif model_type == "xgboost":
                _SEQ_ONLY_M = {"biophysical","kmer","onehot","onehot_vh",
                               "onehot_cdr3","onehot_hcdr3","none","seq"}
                _xgb_X_df_m = _lm_data[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in _lm_data.columns for c in ['HSEQ','LSEQ','CDR3']) else _lm_data.copy()
                _xgb_emb_m  = None if lm in _SEQ_ONLY_M else (
                    X.values if hasattr(X, 'values') else X)
                scores = model.predict_proba(_xgb_X_df_m, embeddings=_xgb_emb_m)
            else:
                scores = model.predict_proba(X)

            _thresh = getattr(model, "recommended_threshold", 0.5)
            labels  = (scores >= _thresh).astype(int)

            if model_type in ("transformer_lm", "cnn") and hasattr(model, 'save_hidden_states'):
                try:
                    _hidden_csv = str(Path(input_file).with_name(
                        f"{Path(input_file).stem}_pred_{target}_{lm}_{model_type}_{db_stem}_hidden.emb.csv"))
                    _X_hidden = X if hasattr(X, 'index') else pd.DataFrame(
                        X, index=base_data.loc[_lm_data.index].index)
                    model.save_hidden_states(_X_hidden, _hidden_csv)
                    print(f"  [hidden] {lm} → {Path(_hidden_csv).name}")
                except Exception as _he:
                    print(f"  [hidden] WARNING {lm}: {_he}")

            score_col = f"{model_type}_{lm}_{db_stem}_score"
            label_col = f"{model_type}_{lm}_{db_stem}_label"
            base_data[score_col] = np.nan
            base_data[label_col] = np.nan
            base_data.loc[_lm_data.index, score_col] = scores
            base_data.loc[_lm_data.index, label_col] = labels

            pos_rate = labels.mean()
            print(f"  threshold={_thresh:.4f}  n={len(scores):,}  pos_rate={pos_rate:.1%}")

            _opt_thresh   = _thresh
            opt_label_col = f"{model_type}_{lm}_{db_stem}_optimallabel"
            opt_labels    = (scores >= _opt_thresh).astype(int)
            base_data[opt_label_col] = np.nan
            base_data.loc[_lm_data.index, opt_label_col] = opt_labels
            print(f"  opt_threshold={_opt_thresh:.4f}  opt_pos_rate={opt_labels.mean():.1%}  → {opt_label_col}")
            results_summary.append((lm, len(scores), pos_rate, _thresh, _opt_thresh))

        except Exception as e:
            print(f"  [ERROR] {lm} failed: {e}")
            failed_lms.append((lm, str(e)))

    if not results_summary:
        print("\n[multi-lm] ERROR: All LMs failed — no output written.")
        return

    path        = Path(input_file)
    output_file = path.with_name(
        f"{path.stem}_pred_{target}_{lm_tag}_{model_type}_{db_stem}{path.suffix}")
    out_data = base_data.reset_index()
    if path.suffix.lower() in [".xlsx", ".xls"]:
        out_data.to_excel(output_file, index=False)
    else:
        out_data.to_csv(output_file, index=False)

    print(f"\n{'═'*62}")
    print(f"  MULTI-LM PREDICTION SUMMARY")
    print(f"{'─'*62}")
    print(f"  {'LM':25s}  {'n':>7}  {'pos_rate':>9}  {'threshold':>10}  {'opt_thresh':>10}")
    print(f"  {'─'*25}  {'─'*7}  {'─'*9}  {'─'*10}  {'─'*10}")
    for row in results_summary:
        lm_name, n, pos, thresh = row[0], row[1], row[2], row[3]
        opt_t = row[4] if len(row) > 4 else thresh
        print(f"  {lm_name:25s}  {n:>7,}  {pos:>9.1%}  {thresh:>10.4f}  {opt_t:>10.4f}")
    if failed_lms:
        print(f"{'─'*62}")
        for lm_name, err in failed_lms:
            print(f"  [FAILED] {lm_name}: {err}")
    print(f"{'─'*62}")
    print(f"  Output → {output_file}")
    print(f"{'═'*62}\n")

    _eval_label_col = test_target if test_target else target
    if _eval_label_col in out_data.columns and results_summary:
        print(f"\n[eval] Ground-truth '{_eval_label_col}' found — running per-LM eval ...")
        _eval_ok = _biophys_ok = False
        _base_cols = [c for c in out_data.columns
                      if not any(c.endswith(s) for s in ('_score', '_label', '_optimallabel'))]

        for row in results_summary:
            lm_name, _n, _pos, _thresh = row[0], row[1], row[2], row[3]
            _score_col    = f"{model_type}_{lm_name}_{db_stem}_score"
            _label_col    = f"{model_type}_{lm_name}_{db_stem}_label"
            _optlabel_col = f"{model_type}_{lm_name}_{db_stem}_optimallabel"
            _lm_stem = str(path.with_name(
                f"{path.stem}_pred_{target}_{lm_name}_{model_type}_{db_stem}"))
            print(f"\n  ── LM: {lm_name}  score_col={_score_col}")

            _lm_cols = _base_cols + [c for c in [_score_col, _label_col, _optlabel_col]
                                      if c in out_data.columns]
            _lm_df   = out_data[_lm_cols].copy()
            _lm_file = str(path.with_name(
                f"{path.stem}_pred_{target}_{lm_name}_{model_type}_{db_stem}{path.suffix}"))
            if path.suffix.lower() in [".xlsx", ".xls"]:
                _lm_df.to_excel(_lm_file, index=False)
            else:
                _lm_df.to_csv(_lm_file, index=False)
            print(f"  Per-LM file ({len(_lm_df.columns)} cols) → {Path(_lm_file).name}")

            try:
                from utils.evaluate_model import evaluate
                _eval_ok = True
                evaluate(file=_lm_file, target=target, score_col=_score_col,
                         cost_fp=1.0, cost_fn=3.0, out=_lm_stem,
                         test_target=test_target, model_type=model_type,
                         lm=lm_name, db_stem=db_stem, dataset_name=Path(input_file).stem)
            except ImportError:
                if not _eval_ok: print("[eval] utils/evaluate_model.py not found — skipping.")
                _eval_ok = False
            except Exception as _e:
                print(f"  [eval] WARNING: {_e}")

            try:
                from utils.plot_biophysical import plot_biophysical_report
                _biophys_ok = True
                plot_biophysical_report(file=_lm_file, target=target,
                                        test_target=test_target, out=_lm_stem,
                                        dataset_name=Path(input_file).stem)
            except ImportError:
                if not _biophys_ok: print("[biophys] utils/plot_biophysical.py not found — skipping.")
                _biophys_ok = False
            except Exception as _e:
                print(f"  [biophys] WARNING: {_e}")

        print(f"\n[eval] Per-LM outputs saved alongside:\n  {output_file}")


# ===========================================================================
# MAIN
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(description="IPI Antibody Developability Prediction Platform")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--predict",         type=str,            help="Predict on new file")
    group.add_argument("--build-embedding", type=str,            help="Generate embeddings only")
    group.add_argument("--kfold",           type=int,            help="Run k-fold CV")
    group.add_argument("--train",           action="store_true", help="Train final model")
    group.add_argument("--split-dataset",   action="store_true", dest="split_dataset",
                       help="Split --db into train + val files.")

    parser.add_argument("--target",
                        type=str,
                        default="psr_filter",
                        help=(
                            "Target column — any column in your database. "
                            "Binary (classification): psr_filter, sec_filter, purif_filter, spr_anno. "
                            "Continuous (regression): psr_norm_mean, psr_norm_dna, psr_norm_avidin, "
                            "psr_norm_smp, sec_monomer_pct, hic_retention_time, fab_tm, "
                            "viscosity_cP, hek_titer, or any numeric column."
                        ))
    parser.add_argument("--lm", default="antiberta2")
    parser.add_argument("--model", default="xgboost",
                        choices=["xgboost","rf","cnn","transformer_onehot","transformer_lm"])
    parser.add_argument("--db",      type=str)
    parser.add_argument("--cluster", type=float, default=0.8, metavar="THRESHOLD")
    parser.add_argument("--split",   type=float, default=0.0, metavar="TRAIN_FRAC")
    parser.add_argument("--val",     type=str,   default=None, metavar="VAL_FILE")
    parser.add_argument("--cluster_col", type=str, default="CDR3",
                        choices=["CDR3", "VH", "VHVL"])
    parser.add_argument("--no-aug",  dest="no_aug", action="store_true", default=True)
    parser.add_argument("--test_target", type=str, default=None)
    parser.add_argument("--mutagenesis", type=int, nargs="?",
        const=50, default=None, metavar="N",
        help="Run in-silico CDR3 mutagenesis on the first N antibodies. "
             "--mutagenesis       → uses N from mutagenesis.max_samples in YAML (default 50). "
             "--mutagenesis 200   → scans first 200 antibodies. "
             "--mutagenesis 0     → scans ALL antibodies (no limit). "
             "Saves TIFF heatmaps + PPT in {stem}_mutagenesis/ folder.")
    parser.add_argument("--threshold", type=float, default=None,
        metavar="T",
        help="Override classification threshold (0–1) for --predict. "
             "Default: use threshold embedded in model checkpoint from kfold, "
             "else 0.5.  Example: --threshold 0.42")
    parser.add_argument("--model_path", type=str, default=None,
        metavar="PATH",
        help="Explicit path to a pretrained checkpoint for --predict. "
             "Overrides the automatic FINAL_{target}_{lm}_{model}_{db}.pt lookup. "
             "Use this to select between Mode 1/2/3 or fine-tuned checkpoints. "
             "Example: --model_path models/FINAL_psr_filter_ablang_transformer_lm_ipi_lora8.pt")
    # ── Fine-tuning arguments ──────────────────────────────────────────────────
    parser.add_argument("--finetune", action="store_true", default=False,
        help="Fine-tune a pretrained model on --finetune_db (Level 2). "
             "Requires --pretrained and --finetune_db.")
    parser.add_argument("--finetune_plm", action="store_true", default=False,
        help="Mode 2: end-to-end fine-tuning of PLM + classifier together during --train. "
             "Slower but no .emb.csv needed. Requires GPU for practical speed.")
    parser.add_argument("--pretrained", type=str, default=None,
        metavar="PATH",
        help="Path to pretrained .pt checkpoint for --finetune.")
    parser.add_argument("--finetune_db", type=str, default=None,
        metavar="FILE",
        help="New dataset (Excel/CSV) for fine-tuning. Used with --finetune.")
    parser.add_argument("--freeze_layers", type=str, default="1",
        metavar="N",
        help="Classifier layer freezing for --finetune. "
             "'all' = head only, '0'/'none' = full, N = freeze first N layers. "
             "Default: 1 (freeze first transformer layer).")
    parser.add_argument("--freeze_plm_layers", type=int, default=10,
        metavar="N",
        help="PLM layers to freeze from bottom for --finetune_plm. "
             "Default: 10 (freeze first 10 of 12 layers, update top 2).")
    parser.add_argument("--finetune_lr", type=float, default=1e-6,
        metavar="LR",
        help="Learning rate for fine-tuning. Should be 10-100x lower than "
             "original training lr. Default: 1e-6.")
    parser.add_argument("--finetune_epochs", type=int, default=10,
        metavar="N",
        help="Number of fine-tuning epochs. Default: 10.")
    parser.add_argument("--lr_plm", type=float, default=1e-6,
        metavar="LR",
        help="Learning rate for PLM layers in --finetune_plm mode. Default: 1e-6.")
    parser.add_argument("--lr_classifier", type=float, default=1e-4,
        metavar="LR",
        help="Learning rate for classifier in --finetune_plm mode. Default: 1e-4.")
    parser.add_argument("--peft", type=str, default="none",
        choices=["none", "lora"],
        help="PEFT method for --finetune_plm. "
             "'none' = Mode 2 (standard layer unfreezing, default). "
             "'lora' = Mode 3 (LoRA low-rank adaptation — recommended). "
             "LoRA trains only ~400k params vs 20M for Mode 2. "
             "Safer for small datasets, works on CPU.")
    parser.add_argument("--lora_r", type=int, default=8,
        metavar="R",
        help="LoRA rank for --peft lora. "
             "4=small dataset (<1k), 8=medium (default), 16=large (>50k).")
    parser.add_argument("--lora_alpha", type=float, default=16.0,
        metavar="A",
        help="LoRA scaling alpha. Usually 2×lora_r. Default: 16.0.")
    parser.add_argument("--lora_layers", type=int, nargs="+", default=None,
        metavar="N",
        help="Which PLM encoder layers get LoRA. "
             "Default: last 2 layers (safe for small datasets). "
             "Example: --lora_layers 10 11  or  --lora_layers 8 9 10 11.")
    parser.add_argument("--cost_fn", type=float, default=3.0,
        metavar="W",
        help="Cost of a False Negative (missing a FAIL antibody) relative to FP. "
             "Higher → lower threshold → higher Recall(FAIL). "
             "Default=3.0 (missing a bad antibody costs 3× a false alarm). "
             "Use --cost_fn 1.0 for balanced (Youden J). "
             "Use --cost_fn 5.0 to aggressively catch polyreactive antibodies.")
    parser.add_argument("--cost_fp", type=float, default=1.0,
        metavar="W",
        help="Cost of a False Positive (flagging a good antibody). Default=1.0.")

    args    = parser.parse_args()

    # ── Smart --db detection ──────────────────────────────────────────────────
    # If --db points to a .pt or .pkl file → treat it as --model_path directly.
    # This allows the clean syntax:
    #   --db FINAL_psr_filter_ablang_transformer_lm_ipi_lora8.pt
    # instead of requiring a separate --model_path flag.
    _raw_db = args.db or get_default_db_path()
    if _raw_db and str(_raw_db).lower().endswith(('.pt', '.pkl')):
        # --db is a model checkpoint → use as model_path, clear db for lookup
        if not getattr(args, 'model_path', None):
            args.model_path = _raw_db
            print(f"[db] --db is a checkpoint file → using as --model_path")
            print(f"[db] model_path = {_raw_db}")
        # db_path is no longer a dataset — derive db_stem from checkpoint name
        # e.g. FINAL_psr_filter_ablang_transformer_lm_ipi_psr_trainset_lora8.pt
        #       → db_stem = ipi_psr_trainset  (everything between model_type and suffix)
        import re as _re
        _stem = Path(_raw_db).stem                # FINAL_psr_filter_ablang_..._lora8
        # Extract the db portion: after _{model_type}_ and before any _ft_/_lora_/_plmft
        _m = _re.search(
            r'(?:transformer_lm|transformer_onehot|rf|xgboost|cnn)_(.+?)(?:_ft_|_lora|_plmft|$)',
            _stem)
        _inferred_db = _m.group(1) if _m else _stem
        args.db  = None            # no dataset file — embeddings not needed for predict
        db_path  = None
        print(f"[db] Inferred db_stem from checkpoint: '{_inferred_db}'")
        print(f"[db] To load embeddings from a different file, "
              f"pass --db dataset.xlsx alongside --model_path")
    else:
        db_path = _raw_db

    _log_path = _setup_logging(args, db_path)

    _cluster_thresh = args.cluster
    # Column name and sequence depend on --cluster_col
    _cluster_col_src = getattr(args, 'cluster_col', 'CDR3').upper()
    _cluster_col_map = {
        'CDR3':  f"HCDR3_CLUSTER_{_cluster_thresh}",
        'HSEQ':  f"VH_CLUSTER_{_cluster_thresh}",
        'VH':    f"VH_CLUSTER_{_cluster_thresh}",
        'VHVL':  f"VHVL_CLUSTER_{_cluster_thresh}",
        'LSEQ':  f"VL_CLUSTER_{_cluster_thresh}",
    }
    _cluster_col = _cluster_col_map.get(_cluster_col_src,
                                         f"HCDR3_CLUSTER_{_cluster_thresh}")

    if args.split_dataset:
        if not db_path:
            parser.error("--db required for --split-dataset")
        if not (0.0 < args.split < 1.0):
            parser.error("--split must be in (0,1), e.g. 0.8 for 80% train")

        import sys as _sys
        _SKIP_LMS    = {"onehot", "onehot_vh"}
        _lm_raw      = args.lm.strip()
        _lm_explicit = "--lm" in _sys.argv

        if not _lm_explicit:
            _lms_to_check = []
            print("[split-dataset] --lm not specified — skipping embedding generation.")
        elif _lm_raw == "all":
            _lms_to_check = list(_ALL_PLM_LMS)
        elif "," in _lm_raw:
            _lms_to_check = [x.strip() for x in _lm_raw.split(",")
                             if x.strip() and x.strip() not in _SKIP_LMS]
        else:
            _lms_to_check = [_lm_raw] if _lm_raw not in _SKIP_LMS else []

        if _lms_to_check:
            print(f"\n[split-dataset] LMs to check: {_lms_to_check}")
            for _lm in _lms_to_check:
                _emb_path = f"{db_path}.{_lm}.emb.csv"
                if os.path.exists(_emb_path):
                    print(f"  [emb] {_lm}: ✓ exists — will split")
                else:
                    print(f"  [emb] {_lm}: not found — generating ...")
                    generate_embedding(db_path, lm=_lm)
                    if os.path.exists(_emb_path):
                        print(f"  [emb] {_lm}: ✓ generated")
                    else:
                        print(f"  [emb] {_lm}: ✗ generation failed — skipping")

        split_and_save(db_path=db_path, split=args.split,
                       cluster_thresh=_cluster_thresh,
                       cluster_col=getattr(args, 'cluster_col', 'CDR3'),
                       label_col=args.target)
        return

    if args.build_embedding:
        lms = ["ablang","antiberty","antiberta2","antiberta2-cssp"] if args.lm == "all" else [args.lm]
        if args.lm in ["onehot","onehot_vh"]:
            print("One-hot encoding does not require pre-generation.")
        else:
            for lm in lms:
                generate_embedding(args.build_embedding, lm=lm)
        return

    if args.kfold:
        if not db_path:
            parser.error("--db required for k-fold")
        X, data, y = load_data(db_path, lm=args.lm, label_col=args.target)
        title = f"{args.target.upper()}_{args.model}"

        # ── Auto-generate HCDR3 clustering if missing ─────────────────────────
        if _cluster_col not in data.columns:
            if "CDR3" in data.columns:
                print(f"\n[kfold] '{_cluster_col}' not found — computing automatically ...")
                print(f"[kfold] Clustering {len(data):,} CDR3 sequences "
                      f"(Levenshtein threshold={_cluster_thresh}) ...")
                try:
                    from utils.clustering import greedy_clustering_by_levenshtein
                    # Use the correct sequence column based on --cluster_col
                    _seq_col_map = {'CDR3': 'CDR3', 'HSEQ': 'HSEQ', 'VH': 'HSEQ',
                                    'VHVL': 'HSEQ', 'LSEQ': 'LSEQ'}
                    _seq_col = _seq_col_map.get(_cluster_col_src, 'CDR3')
                    if _seq_col not in data.columns:
                        print(f"[kfold] WARNING: sequence column '{_seq_col}' not found "
                              f"for clustering — trying CDR3")
                        _seq_col = 'CDR3' if 'CDR3' in data.columns else None
                    if _seq_col:
                        _seqs = data[_seq_col].fillna('').tolist()
                        data[_cluster_col] = greedy_clustering_by_levenshtein(
                            _seqs, _cluster_thresh)
                        n_clust = data[_cluster_col].nunique()
                        print(f"[kfold] {n_clust:,} clusters  "
                              f"(mean {len(data)/n_clust:.1f} sequences/cluster)")
                        try:
                            _df_save = read_dataframe(db_path)
                            _df_save[_cluster_col] = data[_cluster_col].values
                            save_dataframe(_df_save, db_path)
                            print(f"[kfold] ✓ Saved {_cluster_col} → {Path(db_path).name}")
                            print(f"[kfold]   (future runs will skip clustering step)")
                        except Exception as _ce:
                            print(f"[kfold] NOTE: could not save clustering — {_ce}")
                except ImportError:
                    print(f"[kfold] WARNING: utils/clustering.py not found — "
                          f"falling back to StratifiedKFold (no CDR3 leakage protection)")
                except Exception as _clust_e:
                    import traceback
                    print(f"[kfold] WARNING: clustering failed — {_clust_e}")
                    print(traceback.format_exc())
            else:
                print(f"[kfold] WARNING: '{_cluster_col}' not found and no CDR3/HSEQ column — "
                      f"StratifiedKFold will be used (sequence leakage not prevented)")
        else:
            n_clust = data[_cluster_col].nunique()
            print(f"[kfold] Using existing '{_cluster_col}': {n_clust:,} clusters")

        if args.model == "xgboost":
            db_stem   = Path(db_path).stem if db_path else ""
            _SEQ_ONLY = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3",
                         "biophysical","kmer","seq","none"}
            _xgb_emb  = None if args.lm in _SEQ_ONLY else (
                X.values if hasattr(X, 'values') else X)
            _xgb_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()

            # Build feature override dict from --lm (same logic as RF)
            _lm_k = args.lm.lower()
            if _lm_k == "biophysical":
                _xgb_feat = {'embedding': False, 'biophysical': True,
                             'kmer': False, 'onehot': False}
            elif _lm_k == "kmer":
                _xgb_feat = {'embedding': False, 'biophysical': False,
                             'kmer': True, 'onehot': False}
            elif _lm_k in ("onehot","onehot_vh","onehot_cdr3","onehot_hcdr3"):
                _oh_seq = {"onehot":"VHVL","onehot_vh":"VH",
                           "onehot_cdr3":"HCDR3","onehot_hcdr3":"HCDR3"}.get(_lm_k,"VHVL")
                _xgb_feat = {'embedding': False, 'biophysical': False,
                             'kmer': False, 'onehot': True,
                             '_onehot_sequence': _oh_seq}
            elif _lm_k in ("none","seq"):
                _xgb_feat = {'embedding': False}
            else:
                _xgb_feat = {'embedding': True}   # PLM mode

            XGBoostModel.kfold_validation(
                data, _xgb_X_df, y,
                embeddings        = _xgb_emb,
                embedding_lm      = args.lm,
                title             = title,
                kfold             = args.kfold,
                target            = args.target,
                cluster_col       = _cluster_col,
                db_stem           = db_stem,
                override_features = _xgb_feat,
                cost_fn           = getattr(args, 'cost_fn', 3.0),
                cost_fp           = getattr(args, 'cost_fp', 1.0),
            )
        elif args.model == "rf":
            # [FIX] RF kfold: pass X_df + embeddings separately
            db_stem  = Path(db_path).stem if db_path else ""
            _rf_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            # seq-only modes have no embedding array
            _SEQ_ONLY = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}
            _rf_emb = None if args.lm in _SEQ_ONLY else (X.values if hasattr(X, 'values') else X)
            # Auto-set feature flags from --lm (same logic as --train)
            _tmp_rf = RandomForestModel()
            _lm_k = args.lm.lower()
            if _lm_k == "biophysical":
                _tmp_rf.config['features'].update({'embedding': False, 'biophysical': True, 'kmer': False, 'onehot': False})
            elif _lm_k == "kmer":
                _tmp_rf.config['features'].update({'embedding': False, 'biophysical': False, 'kmer': True, 'onehot': False})
            elif _lm_k in ("onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3"):
                # Map --lm flag to onehot.sequence setting
                _oh_seq = {"onehot": "VHVL", "onehot_vh": "VH",
                           "onehot_cdr3": "HCDR3", "onehot_hcdr3": "HCDR3"}.get(_lm_k, "VHVL")
                _tmp_rf.config['features'].update({'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True})
                _tmp_rf.config.setdefault('onehot', {})['sequence'] = _oh_seq
            elif _lm_k in ("none", "seq"):
                _tmp_rf.config['features']['embedding'] = False
            else:
                _tmp_rf.config['features']['embedding'] = True
            # Build feature override dict from --lm
            _lm_k = args.lm.lower()
            if _lm_k == "biophysical":
                _rf_feat_override = {'embedding': False, 'biophysical': True, 'kmer': False}
            elif _lm_k == "kmer":
                _rf_feat_override = {'embedding': False, 'biophysical': False, 'kmer': True}
            elif _lm_k in ("none", "seq"):
                _rf_feat_override = {'embedding': False}
            elif _lm_k in ("onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3"):
                _oh_seq = {"onehot": "VHVL", "onehot_vh": "VH",
                           "onehot_cdr3": "HCDR3", "onehot_hcdr3": "HCDR3"}.get(_lm_k, "VHVL")
                _rf_feat_override = {'embedding': False, 'biophysical': False,
                                     'kmer': False, 'onehot': True,
                                     '_onehot_sequence': _oh_seq}   # picked up in kfold
            else:
                _rf_feat_override = {'embedding': True}   # PLM mode

            RandomForestModel.kfold_validation(
                data, _rf_X_df, y,
                embeddings        = _rf_emb,
                embedding_lm      = args.lm,
                title             = title,
                kfold             = args.kfold,
                target            = args.target,
                cluster_col       = _cluster_col,
                db_stem           = db_stem,
                override_features = _rf_feat_override,
                cost_fn           = getattr(args, 'cost_fn', 3.0),
                cost_fp           = getattr(args, 'cost_fp', 1.0),
            )
        elif args.model == "cnn":
            db_stem = Path(db_path).stem if db_path else ""
            CNNModel().kfold_validation(db_stem, data, X, y, embedding_lm=args.lm,
                                         title=title, kfold=args.kfold, target=args.target)
        elif args.model == "transformer_onehot":
            if _cluster_col not in data.columns:
                if "CDR3" in data.columns:
                    from utils.clustering import greedy_clustering_by_levenshtein
                    print(f"[kfold] Computing {_cluster_col} (threshold={_cluster_thresh}) ...")
                    data[_cluster_col] = greedy_clustering_by_levenshtein(
                        data["CDR3"].tolist(), _cluster_thresh)
                    n_clust = data[_cluster_col].nunique()
                    print(f"[kfold] {n_clust:,} clusters  (mean {len(data)/n_clust:.1f} sequences/cluster)")
                    try:
                        _df_save = read_dataframe(db_path)
                        _df_save[_cluster_col] = data[_cluster_col].values
                        save_dataframe(_df_save, db_path)
                        print(f"[kfold] Saved {_cluster_col} → {os.path.basename(db_path)}")
                    except Exception as _e:
                        print(f"[kfold] WARNING: could not save {_cluster_col} to db: {_e}")
                else:
                    print(f"[WARN] No CDR3 — falling back to StratifiedKFold.")
            # transformer_onehot: onehot=VH+VL+CDR3, onehot_vh=VH+CDR3
            # onehot_hcdr3/onehot_cdr3 → not supported, fall back to 'onehot'
            _kf_lm = args.lm if args.lm in ("onehot", "onehot_vh") else "onehot"
            TransformerOneHotModel().kfold_validation(
                data, X, y, embedding_lm=_kf_lm, title=title,
                kfold=args.kfold, target=args.target,
                cluster_col=_cluster_col,
                db_stem=Path(db_path).stem if db_path else "")
        elif args.model == "transformer_lm":
            db_stem = Path(db_path).stem if db_path else ""
            TransformerLMModel().kfold_validation(
                db_stem, data, X, y, embedding_lm=args.lm, title=title,
                kfold=args.kfold, target=args.target, cluster_col=_cluster_col)
        return

    # ── --finetune: Level 2 fine-tuning from pretrained checkpoint ───────────
    if getattr(args, 'finetune', False):
        # Validate required args
        if not getattr(args, 'pretrained', None):
            parser.error("--finetune requires --pretrained <path.pt>")
        if not getattr(args, 'finetune_db', None):
            parser.error("--finetune requires --finetune_db <dataset>")
        if args.model not in ('transformer_lm', 'transformer_onehot'):
            parser.error(f"--finetune only supported for transformer_lm and "
                         f"transformer_onehot (got --model {args.model})")

        _pt_path  = args.pretrained
        _ft_db    = args.finetune_db
        _ft_stem  = Path(_ft_db).stem
        _db_stem  = Path(db_path).stem if db_path else "pretrained"

        # Parse freeze_layers — accepts 'all', 'none', or int
        _fl = getattr(args, 'freeze_layers', '1')
        if str(_fl).lower() in ('none', '0'):
            _freeze = 0
        elif str(_fl).lower() == 'all':
            _freeze = 'all'
        else:
            try:    _freeze = int(_fl)
            except: _freeze = 1

        print(f"\n{'='*62}")
        print(f"  FINE-TUNE  (Level 2)")
        print(f"{'='*62}")
        print(f"  Pretrained  : {_pt_path}")
        print(f"  New dataset : {_ft_db}")
        print(f"  Target      : {args.target}")
        print(f"  LM          : {args.lm}")
        print(f"  Freeze      : {_freeze}")
        print(f"  LR          : {args.finetune_lr:.2e}")
        print(f"  Epochs      : {args.finetune_epochs}")
        print(f"{'='*62}\n")

        if args.model == 'transformer_lm':
            # Load pretrained checkpoint
            # Need embedding_dim — get from emb file or checkpoint
            _emb_path_ft = f"{_ft_db}.{args.lm}.emb.csv"
            if not os.path.exists(_emb_path_ft):
                print(f"[finetune] Generating {args.lm} embeddings for {_ft_db} ...")
                _emb_path_ft = generate_embedding(_ft_db, lm=args.lm)
            import pandas as _pd_ft
            _emb_dim_ft = _pd_ft.read_csv(_emb_path_ft, index_col=0).shape[1]

            print(f"[finetune] Loading pretrained: {_pt_path}  (emb_dim={_emb_dim_ft})")
            ft_model = TransformerLMModel.load(_pt_path, embedding_dim=_emb_dim_ft)

            ft_model.fine_tune(
                finetune_db    = _ft_db,
                target         = args.target,
                lm             = args.lm,
                freeze_layers  = _freeze,
                finetune_lr    = args.finetune_lr,
                finetune_epochs= args.finetune_epochs,
                db_stem        = _db_stem,
            )
            # Save fine-tuned model
            _ft_out = (f"{MODEL_DIR}/FINAL_{args.target}_{args.lm}"
                       f"_transformer_lm_{_db_stem}_ft_{_ft_stem}.pt")
            ft_model.save(_ft_out)
            print(f"\n[finetune] Fine-tuned model saved: {_ft_out}")

        elif args.model == 'transformer_onehot':
            ft_model = TransformerOneHotModel.load(_pt_path)
            ft_model.set_lm_mode(args.lm)
            # Load sequence data
            _ext_ft = Path(_ft_db).suffix.lower()
            import pandas as _pd_ft2
            _df_ft  = (_pd_ft2.read_excel(_ft_db) if _ext_ft in ('.xlsx','.xls')
                       else _pd_ft2.read_csv(_ft_db))
            if 'BARCODE' in _df_ft.columns:
                _df_ft = _df_ft.set_index('BARCODE')
            _df_ft = _df_ft.dropna(subset=[args.target])
            _y_ft  = _df_ft[args.target].astype(int).values
            _X_ft  = _df_ft[['HSEQ','LSEQ','CDR3']].copy()

            # Apply layer freezing then train
            import re as _re
            if _freeze == 'all':
                for n_, p in ft_model.model.named_parameters():
                    p.requires_grad = ('head' in n_)
            elif _freeze == 0:
                for p in ft_model.model.parameters():
                    p.requires_grad = True
            else:
                for n_, p in ft_model.model.named_parameters():
                    m_ = _re.search(r'(?:transformer|cdr3_transformer)\.layers\.(\d+)', n_)
                    p.requires_grad = (not m_) or (int(m_.group(1)) >= int(_freeze))

            ft_model.train(_X_ft, _y_ft,
                           target   = args.target,
                           db_stem  = _ft_stem,
                           epochs   = args.finetune_epochs)
            _ft_out = (f"{MODEL_DIR}/FINAL_{args.target}_{args.lm}"
                       f"_transformer_onehot_{_db_stem}_ft_{_ft_stem}.pt")
            ft_model.save(_ft_out)
            print(f"\n[finetune] Fine-tuned model saved: {_ft_out}")

        return

    # ── --finetune_plm: Mode 2 end-to-end PLM + classifier training ───────────
    if getattr(args, 'finetune_plm', False) and args.train:
        if args.model != 'transformer_lm':
            parser.error("--finetune_plm only supported for --model transformer_lm")
        if not db_path:
            parser.error("--finetune_plm requires --db")

        _db_stem_plm = Path(db_path).stem
        _ext_plm     = Path(db_path).suffix.lower()
        import pandas as _pd_plm
        _df_plm = (_pd_plm.read_excel(db_path) if _ext_plm in ('.xlsx','.xls')
                   else _pd_plm.read_csv(db_path))
        if 'BARCODE' in _df_plm.columns:
            _df_plm = _df_plm.set_index('BARCODE')
        _df_plm = _df_plm.dropna(subset=[args.target])
        _y_plm  = _df_plm[args.target].astype(int).values
        print(f"[plmft] Dataset: {db_path}  n={len(_y_plm):,}  "
              f"pos={_y_plm.mean():.1%}")

        plmft_model = TransformerLMModel()
        _peft_mode = getattr(args, 'peft', 'none')
        print(f"[plmft] Mode: {'3 (LoRA)' if _peft_mode=='lora' else '2 (full unfreeze)'}")
        plmft_model.train_with_plm_finetune(
            data_df           = _df_plm,
            target            = args.target,
            lm                = args.lm,
            freeze_plm_layers = getattr(args, 'freeze_plm_layers', 10),
            epochs            = getattr(args, 'finetune_epochs', 20),
            batch_size        = 16,
            lr_classifier     = getattr(args, 'lr_classifier', 1e-4),
            lr_plm            = getattr(args, 'lr_plm', 1e-6),
            db_stem           = _db_stem_plm,
            peft              = _peft_mode,
            lora_r            = getattr(args, 'lora_r',     8),
            lora_alpha        = getattr(args, 'lora_alpha', 16.0),
            lora_layers       = getattr(args, 'lora_layers', None),
        )
        # Filename tags the mode
        _plmft_suffix = f"_lora{args.lora_r}" if _peft_mode == 'lora' else "_plmft"
        _plmft_out = (f"{MODEL_DIR}/FINAL_{args.target}_{args.lm}"
                      f"_transformer_lm_{_db_stem_plm}{_plmft_suffix}.pt")
        plmft_model.save(_plmft_out)
        print(f"\n[plmft] Model saved: {_plmft_out}")
        return

    if args.train:
        if not db_path:
            parser.error("--db required for training")

        db_stem           = Path(db_path).stem
        _has_explicit_val = bool(getattr(args, 'val', None))
        _has_split        = 0.0 < args.split < 1.0
        val_X = val_y = val_db_path = None

        if _has_explicit_val:
            print(f"\n[train] Validation: explicit file → {args.val}")
            X, data, y             = load_data(db_path,  lm=args.lm, label_col=args.target)
            X_val, data_val, y_val = load_data(args.val, lm=args.lm, label_col=args.target)
            val_X, val_y, val_db_path = X_val, y_val, args.val
            print(f"[train] Train n={len(y):,}  mean={y.mean():.4f}  "
                  f"Val n={len(y_val):,}  mean={y_val.mean():.4f}")

        elif _has_split:
            _cc = getattr(args, 'cluster_col', 'CDR3')
            print(f"\n[train] Auto-split: train={args.split:.0%}/val={1-args.split:.0%}  cluster_col={_cc}")
            train_path, val_path, _, _ = split_and_save(
                db_path=db_path, split=args.split,
                cluster_thresh=_cluster_thresh, cluster_col=_cc, label_col=args.target)
            X, data, y             = load_data(train_path, lm=args.lm, label_col=args.target)
            X_val, data_val, y_val = load_data(val_path,   lm=args.lm, label_col=args.target)
            val_X, val_y, val_db_path = X_val, y_val, val_path
        else:
            X, data, y = load_data(db_path, lm=args.lm, label_col=args.target)

        # ── Auto-generate HCDR3 clustering if missing ─────────────────────────
        # Needed for train_test_split_group_stratified and logging purposes.
        # If already present in file → skip. Otherwise compute + save.
        if _cluster_col not in data.columns and "CDR3" in data.columns:
            print(f"\n[train] '{_cluster_col}' not found — computing automatically ...")
            print(f"[train] Clustering {len(data):,} CDR3 sequences "
                  f"(Levenshtein threshold={_cluster_thresh}) ...")
            try:
                from utils.clustering import greedy_clustering_by_levenshtein
                _seq_col_map2 = {'CDR3': 'CDR3', 'HSEQ': 'HSEQ', 'VH': 'HSEQ',
                                 'VHVL': 'HSEQ', 'LSEQ': 'LSEQ'}
                _seq_col2 = _seq_col_map2.get(_cluster_col_src, 'CDR3')
                if _seq_col2 not in data.columns:
                    _seq_col2 = 'CDR3' if 'CDR3' in data.columns else None
                if _seq_col2:
                    _seqs2 = data[_seq_col2].fillna('').tolist()
                    data[_cluster_col] = greedy_clustering_by_levenshtein(
                        _seqs2, _cluster_thresh)
                    n_clust = data[_cluster_col].nunique()
                    print(f"[train] {n_clust:,} clusters  "
                          f"(mean {len(data)/n_clust:.1f} sequences/cluster)")
                    try:
                        _df_save = read_dataframe(db_path)
                        _df_save[_cluster_col] = data[_cluster_col].values
                        save_dataframe(_df_save, db_path)
                        print(f"[train] ✓ Saved {_cluster_col} → {Path(db_path).name}")
                    except Exception as _ce:
                        print(f"[train] NOTE: could not save clustering — {_ce}")
            except ImportError:
                print(f"[train] WARNING: utils/clustering.py not found — "
                      f"group-stratified split unavailable")
            except Exception as _ce:
                print(f"[train] WARNING: clustering failed — {_ce}")

        if args.model == "xgboost":
            model = XGBoostModel(verbose=False)
            # Auto-set feature flags from --lm (same logic as RF)
            _lm_xgb = args.lm.lower()
            if _lm_xgb == "biophysical":
                model.config['features'].update(
                    {'embedding': False, 'biophysical': True, 'kmer': False, 'onehot': False})
                print("[XGB] --lm biophysical → features: embedding=False  biophysical=True")
            elif _lm_xgb == "kmer":
                model.config['features'].update(
                    {'embedding': False, 'biophysical': False, 'kmer': True, 'onehot': False})
                print("[XGB] --lm kmer → features: embedding=False  kmer=True")
            elif _lm_xgb in ("onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3"):
                _oh_seq_xgb = {"onehot": "VHVL", "onehot_vh": "VH",
                               "onehot_cdr3": "HCDR3", "onehot_hcdr3": "HCDR3"}.get(_lm_xgb, "VHVL")
                model.config['features'].update(
                    {'embedding': False, 'biophysical': False, 'kmer': False, 'onehot': True})
                model.config.setdefault('onehot', {})['sequence'] = _oh_seq_xgb
                print(f"[XGB] --lm {_lm_xgb} → features: onehot=True  sequence={_oh_seq_xgb}")
            elif _lm_xgb in ("none", "seq"):
                model.config['features']['embedding'] = False
                print(f"[XGB] --lm {_lm_xgb} → features: embedding=False")
            else:
                # PLM mode
                model.config['features']['embedding'] = True
                print(f"[XGB] --lm {_lm_xgb} → features: embedding=True")
            # Print config banner after overrides applied
            XGBoostModel.print_config_report(model.config)
        elif args.model == "rf":
            model = RandomForestModel(verbose=False)   # suppress banner — reprint after override
            # Auto-set feature flags from --lm so YAML edits are not needed.
            #   --lm biophysical  → embedding=False  biophysical=True  kmer=False
            #   --lm kmer         → embedding=False  biophysical=False kmer=True
            #   --lm none/seq     → embedding=False  (keep kmer/bio from YAML)
            #   --lm ablang etc.  → embedding=True   (keep kmer/bio from YAML)
            _lm = args.lm.lower()
            if _lm == "biophysical":
                model.config['features']['embedding']   = False
                model.config['features']['biophysical'] = True
                model.config['features']['kmer']        = False
                print("[RF] --lm biophysical → features: embedding=False  biophysical=True  kmer=False")
            elif _lm == "kmer":
                model.config['features']['embedding']   = False
                model.config['features']['biophysical'] = False
                model.config['features']['kmer']        = True
                print("[RF] --lm kmer → features: embedding=False  biophysical=False  kmer=True")
            elif _lm in ("none", "seq"):
                model.config['features']['embedding']   = False
                print(f"[RF] --lm {_lm} → features: embedding=False  (kmer/bio from YAML)")
            elif _lm in ("onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3"):
                # One-hot modes — no PLM, no biophysical/kmer unless explicitly added
                _oh_seq = {"onehot": "VHVL", "onehot_vh": "VH",
                           "onehot_cdr3": "HCDR3", "onehot_hcdr3": "HCDR3"}.get(_lm, "VHVL")
                # Ensure features dict has onehot key (may be absent from old YAML)
                model.config.setdefault('features', {}).update(
                    {'embedding': False, 'biophysical': False,
                     'kmer': False, 'onehot': True})
                # Ensure onehot section exists with correct sequence
                model.config.setdefault('onehot', {})['sequence'] = _oh_seq
                print(f"[RF] --lm {_lm} → features.onehot=True  sequence={_oh_seq}  "
                      f"embedding=False  kmer=False  biophysical=False")
            else:
                # PLM mode (ablang, antiberta2, …) — force embedding=True
                model.config['features']['embedding']   = True
                model.config['features'].setdefault('onehot', False)
                print(f"[RF] --lm {_lm} → features: embedding=True  "
                      f"(kmer={model.config['features'].get('kmer',False)}  "
                      f"bio={model.config['features'].get('biophysical',False)})")
            # Reprint config banner now that overrides are applied
            from models.randomforest import RandomForestModel as _RFM
            _RFM.print_config_report(model.config)
        elif args.model == "cnn":
            model = CNNModel()
        elif args.model == "transformer_onehot":
            model = TransformerOneHotModel()
            # transformer_onehot: onehot=VH+VL+CDR3, onehot_vh=VH+CDR3
            # onehot_hcdr3/onehot_cdr3 not applicable → fall back to onehot
            _t_lm = args.lm if args.lm in ("onehot", "onehot_vh") else "onehot"
            model.set_lm_mode(_t_lm)
        elif args.model == "transformer_lm":
            model = TransformerLMModel()

        _vkw = {}
        if val_X is not None and args.model in ("transformer_lm", "transformer_onehot", "cnn"):
            _vkw = {"val_X": val_X, "val_y": val_y}

        if args.model == "transformer_onehot":
            model.train(X, y, target=args.target, db_stem=db_stem,
                        cluster_col=_cluster_col, no_aug=args.no_aug, **_vkw)
        elif args.model == "transformer_lm":
            model.train(X, y, target=args.target, db_stem=db_stem,
                        embedding_lm=args.lm, cluster_col=_cluster_col, **_vkw)
        elif args.model == "rf":
            # [FIX] RF train: pass X_df + embeddings + optional val data
            _rf_train_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            # Seq-only modes: no embedding array — FeatureBuilder handles features internally
            _SEQ_ONLY_TRAIN = {"biophysical", "kmer", "onehot", "onehot_vh",
                               "onehot_cdr3", "onehot_hcdr3", "none", "seq"}
            _rf_train_emb  = None if args.lm in _SEQ_ONLY_TRAIN                              else (X.values if hasattr(X, 'values') else X)
            _rf_val_X_df   = None
            _rf_val_emb    = None
            if val_X is not None and val_y is not None and data_val is not None:
                _rf_val_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data_val.columns for c in ['HSEQ','LSEQ','CDR3']) else data_val.copy()
                _rf_val_emb  = None if args.lm in _SEQ_ONLY_TRAIN                                else (val_X.values if hasattr(val_X, 'values') else val_X)
            # full_stem used as output prefix for all RF files
            # e.g.  psr_norm_dna_ablang_rf_ipi_psr_train_elisa
            _y_uniq_rf   = len(set(y.tolist() if hasattr(y, 'tolist') else list(y)))
            _rf_task_tag = "_regression" if _y_uniq_rf > 2 else ""
            _rf_full_stem = f"{args.target}_{args.lm}_rf_{db_stem}{_rf_task_tag}"

            # Seq-only modes (biophysical, kmer, none): X is a DataFrame,
            # no embedding array exists — pass embeddings=None to FeatureBuilder
            _SEQ_ONLY = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}
            if args.lm in _SEQ_ONLY:
                _rf_train_emb = None
                _rf_val_emb   = None
                # Map --lm to feature flags — definitive override
                _oh_map = {"onehot": "VHVL", "onehot_vh": "VH",
                           "onehot_hcdr3": "HCDR3", "onehot_cdr3": "HCDR3"}
                if args.lm in _oh_map:
                    _oh_mode = _oh_map[args.lm]
                    model.config.setdefault('features', {}).update(
                        {'embedding': False, 'biophysical': False,
                         'kmer': False, 'onehot': True})
                    model.config.setdefault('onehot', {})['sequence'] = _oh_mode
                    print(f"[RF] --lm {args.lm} → onehot=True  sequence={_oh_mode}  embedding=False")
                elif args.lm == "biophysical":
                    model.config.setdefault('features', {}).update(
                        {'embedding': False, 'biophysical': True, 'kmer': False, 'onehot': False})
                    print(f"[RF] --lm biophysical → biophysical=True  embedding=False")
                elif args.lm == "kmer":
                    model.config.setdefault('features', {}).update(
                        {'embedding': False, 'biophysical': False, 'kmer': True, 'onehot': False})
                    print(f"[RF] --lm kmer → kmer=True  embedding=False")
                else:
                    model.config.setdefault('features', {})['embedding'] = False
                    print(f"[RF] --lm {args.lm} → embedding=False  (kmer/bio from YAML)")

            model.train(_rf_train_X_df, y,
                        embeddings     = _rf_train_emb,
                        val_X          = _rf_val_X_df,
                        val_y          = val_y if val_y is not None else None,
                        val_embeddings = _rf_val_emb,
                        target         = _rf_full_stem,
                        target_col     = args.target,
                        embedding_lm   = args.lm)
        elif args.model == "xgboost":
            # Mirror RF: pass X_df + embeddings separately
            _xgb_train_X_df = data[['HSEQ','LSEQ','CDR3']].copy() if all(
                c in data.columns for c in ['HSEQ','LSEQ','CDR3']) else data.copy()
            _SEQ_ONLY_XGB = {"biophysical", "kmer", "onehot", "onehot_vh",
                             "onehot_cdr3", "onehot_hcdr3", "none", "seq"}
            _xgb_train_emb = None if args.lm in _SEQ_ONLY_XGB else (
                X.values if hasattr(X, 'values') else X)
            _xgb_val_X_df = _xgb_val_emb = None
            if val_X is not None and val_y is not None and data_val is not None:
                _xgb_val_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data_val.columns for c in ['HSEQ','LSEQ','CDR3']) else data_val.copy()
                _xgb_val_emb = None if args.lm in _SEQ_ONLY_XGB else (
                    val_X.values if hasattr(val_X, 'values') else val_X)
            # Detect task from y values (same logic as auto-detection in model)
            _y_uniq_xgb   = len(set(y.tolist() if hasattr(y, 'tolist') else list(y)))
            _is_bin_xgb   = (_y_uniq_xgb <= 2)
            _xgb_task_tag = "_regression" if not _is_bin_xgb else ""
            _xgb_full_stem = f"{args.target}_{args.lm}_xgboost_{db_stem}{_xgb_task_tag}"
            model.train(_xgb_train_X_df, y,
                        embeddings     = _xgb_train_emb,
                        val_X          = _xgb_val_X_df,
                        val_y          = val_y if val_y is not None else None,
                        val_embeddings = _xgb_val_emb,
                        target         = _xgb_full_stem,
                        target_col     = args.target,
                        embedding_lm   = args.lm)
        else:
            model.train(X, y, embedding_lm=args.lm)

        ext  = ".pt" if args.model in ["cnn","transformer_onehot","transformer_lm"] else ".pkl"
        # Add _regression suffix so filename is self-documenting
        _task_suffix = ""
        if hasattr(model, 'task') and model.task == 'regression':
            _task_suffix = "_regression"
        elif args.model in ('xgboost', 'rf') and hasattr(model, 'task'):
            pass   # classification has no suffix (default)
        path = f"{MODEL_DIR}/FINAL_{args.target}_{args.lm}_{args.model}_{db_stem}{_task_suffix}{ext}"
        model.save(path)
        print(f"FINAL MODEL SAVED: {path}")

        if val_X is not None and val_db_path is not None:
            print(f"\n[train] Predicting on validation set ...")
            _vp = Path(val_db_path)
            # [FIX] RF val predict_proba
            if args.model == "rf":
                _val_rf_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if all(
                    c in data_val.columns for c in ['HSEQ','LSEQ','CDR3']) else data_val.copy()
                _SEQ_ONLY_VAL = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3",
                                  "biophysical","kmer","seq","none"}
                _val_rf_emb  = None if args.lm in _SEQ_ONLY_VAL else (
                    val_X.values if hasattr(val_X, 'values') else val_X)
                _vscores = model.predict_proba(_val_rf_X_df, embeddings=_val_rf_emb)
            elif args.model == "xgboost":
                _SEQ_ONLY_VAL_XGB = {"onehot","onehot_vh","onehot_cdr3","onehot_hcdr3",
                                     "biophysical","kmer","seq","none"}
                _val_xgb_X_df = data_val[['HSEQ','LSEQ','CDR3']].copy() if (
                    data_val is not None and
                    all(c in data_val.columns for c in ['HSEQ','LSEQ','CDR3'])
                ) else (data_val.copy() if data_val is not None else pd.DataFrame())
                _val_xgb_emb  = None if args.lm in _SEQ_ONLY_VAL_XGB else (
                    val_X.values if hasattr(val_X, 'values') else val_X)
                _vscores = model.predict_proba(_val_xgb_X_df, embeddings=_val_xgb_emb)
            else:
                _vscores = model.predict_proba(val_X)

            _vthresh = getattr(model, "recommended_threshold", 0.5)
            _vlabels = (_vscores >= _vthresh).astype(int)
            print(f"[train] Val threshold={_vthresh:.4f}  pos={_vlabels.mean():.1%}")

            val_out_df = read_dataframe(val_db_path)
            if "BARCODE" in val_out_df.columns:
                val_out_df = val_out_df.set_index("BARCODE")
            _score_col = f"{args.model}_{args.lm}_{db_stem}_score"
            _label_col = f"{args.model}_{args.lm}_{db_stem}_label"
            if len(val_out_df) == len(_vscores):
                val_out_df[_score_col] = _vscores
                val_out_df[_label_col] = _vlabels
            else:
                print(f"[train] WARNING: row count mismatch — skipping score columns")

            val_pred_path = str(_vp.parent /
                f"{_vp.stem}_pred_{args.target}_{args.lm}_{args.model}_{db_stem}{_vp.suffix}")
            save_dataframe(val_out_df.reset_index(), val_pred_path)
            print(f"[train] Val predictions → {val_pred_path}")

            _val_test_target = args.test_target if args.test_target else args.target
            _val_eval_stem   = str(_vp.parent /
                f"{_vp.stem}_pred_{args.target}_{args.lm}_{args.model}_{db_stem}")
            _val_flat = val_out_df.reset_index()

            if _val_test_target in _val_flat.columns:
                print(f"\n[train-eval] Ground-truth '{_val_test_target}' found — running evaluation ...")
                try:
                    from utils.evaluate_model import evaluate
                    evaluate(file=val_pred_path, target=args.target, score_col=_score_col,
                             cost_fp=1.0, cost_fn=3.0, out=_val_eval_stem,
                             test_target=_val_test_target, model_type=args.model,
                             lm=args.lm, db_stem=db_stem, dataset_name=_vp.stem)
                except ImportError:
                    print("[train-eval] utils/evaluate_model.py not found — skipping.")
                except Exception as _e:
                    print(f"[train-eval] WARNING: evaluation failed — {_e}")

                try:
                    from utils.plot_biophysical import plot_biophysical_report
                    plot_biophysical_report(file=val_pred_path, target=args.target,
                                            test_target=_val_test_target, out=_val_eval_stem,
                                            dataset_name=_vp.stem)
                except ImportError:
                    print("[train-eval] utils/plot_biophysical.py not found — skipping.")
                except Exception as _e:
                    print(f"[train-eval] WARNING: biophysical plot failed — {_e}")
            else:
                print(f"\n[train-eval] NOTE: '{_val_test_target}' not found in val set — skipping plots.")

    if args.predict:
        _VALID_PLM_LMS = {"ablang", "antiberty", "antiberta2", "antiberta2-cssp", "igbert"}
        _SEQ_ONLY_LMS  = {"onehot", "onehot_vh", "onehot_cdr3", "onehot_hcdr3", "biophysical", "kmer", "seq", "none"}
        _VALID_ALL_LMS = _VALID_PLM_LMS | _SEQ_ONLY_LMS

        _lm_raw = args.lm.strip()
        if _lm_raw == "all":
            _lm_list = _ALL_PLM_LMS
        elif "," in _lm_raw:
            _lm_list = [x.strip() for x in _lm_raw.split(",") if x.strip()]
            _invalid = [x for x in _lm_list if x not in _VALID_ALL_LMS]
            if _invalid:
                parser.error(f"Invalid LM(s): {_invalid}. Valid choices: {sorted(_VALID_ALL_LMS)} or 'all'")
        else:
            if _lm_raw not in _VALID_ALL_LMS:
                parser.error(f"Invalid LM: '{_lm_raw}'. Valid choices: {sorted(_VALID_ALL_LMS)} or 'all'")
            _lm_list = None

        if _lm_list is not None:
            _lm_tag = "all" if set(_lm_list) == set(_ALL_PLM_LMS) else "_".join(_lm_list)
            auto_predict_multi_lm(args.predict, target=args.target, lms=_lm_list,
                                   lm_tag=_lm_tag, model_type=args.model,
                                   db_path=args.db, test_target=args.test_target)
        else:
            auto_predict(args.predict, target=args.target, lm=_lm_raw,
                         model_type=args.model, db_path=args.db,
                         test_target=args.test_target,
                         run_mutagenesis=args.mutagenesis is not None,
                         mutagenesis_n=args.mutagenesis,
                         threshold=args.threshold,
                         model_path=getattr(args, 'model_path', None))

        # Mutagenesis is handled inside auto_predict() via _run_cdr3_mutagenesis()
        # No additional processing needed here.


if __name__ == "__main__":
    try:
        main()
    finally:
        if isinstance(sys.stdout, _Tee):
            ts_end = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            print(f"\n[log] Finished: {ts_end}")
            sys.stdout.flush()
            sys.stdout.close()
            sys.stdout = sys.__stdout__
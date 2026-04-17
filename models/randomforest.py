# models/randomforest.py
# Random Forest — Classification & Regression
# MLAbDev · IPI Antibody Developability Prediction Platform
# Supports: embedding | k-mer | biophysical features (combinable)
# Interpretation: SHAP (TreeExplainer) for k-mer + biophysical
# Updated: full rewrite — DEC-2025

import os, sys, copy, logging, warnings, datetime
warnings.filterwarnings('ignore')

project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(project_root)

from config import MODEL_DIR

try:
    from utils.threshold_optimizer import run_full_threshold_pipeline
    _THRESHOLD_OPT_AVAILABLE = True
except ImportError:
    _THRESHOLD_OPT_AVAILABLE = False

try:
    import shap as _shap_lib
    _SHAP_AVAILABLE = True
except ImportError:
    _SHAP_AVAILABLE = False

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import joblib, yaml
from itertools import product
from collections import Counter
from scipy.stats import randint, pearsonr, spearmanr

from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.model_selection import (
    StratifiedKFold, StratifiedGroupKFold, RandomizedSearchCV
)
from sklearn.metrics import (
    roc_auc_score, roc_curve, accuracy_score, f1_score,
    precision_score, recall_score,
    mean_squared_error, mean_absolute_error, r2_score,
)


# ══════════════════════════════════════════════════════════════════════════════
# 1.  AMINO ACID CONSTANTS & BIOPHYSICAL SCALES
# ══════════════════════════════════════════════════════════════════════════════

AMINO_ACIDS = 'ACDEFGHIKLMNPQRSTVWY'

_PKA = {'D': 3.9, 'E': 4.1, 'C': 8.3, 'Y': 10.1,
        'H': 6.0, 'K': 10.5, 'R': 12.5,
        'N_term': 8.0, 'C_term': 3.1}

_KD = {'A': 1.8, 'R':-4.5, 'N':-3.5, 'D':-3.5, 'C': 2.5,
       'Q':-3.5, 'E':-3.5, 'G':-0.4, 'H':-3.2, 'I': 4.5,
       'L': 3.8, 'K':-3.9, 'M': 1.9, 'F': 2.8, 'P':-1.6,
       'S':-0.8, 'T':-0.7, 'W':-0.9, 'Y':-1.3, 'V': 4.2}


# ══════════════════════════════════════════════════════════════════════════════
# 2.  BIOPHYSICAL & K-MER FEATURE COMPUTATION
# ══════════════════════════════════════════════════════════════════════════════

def _charge_at_ph7(seq):
    c = 0.0
    for aa in seq.upper():
        if   aa == 'D': c -= 1/(1+10**(_PKA['D']-7))
        elif aa == 'E': c -= 1/(1+10**(_PKA['E']-7))
        elif aa == 'H': c += 1/(1+10**(7-_PKA['H']))
        elif aa == 'K': c += 1/(1+10**(7-_PKA['K']))
        elif aa == 'R': c += 1/(1+10**(7-_PKA['R']))
    return round(c, 4)

def _gravy(seq):
    vals = [_KD.get(aa, 0.0) for aa in seq.upper()]
    return round(sum(vals)/len(vals), 4) if vals else 0.0

def _pi(seq):
    def _c(ph):
        c = 1/(1+10**(ph-_PKA['N_term'])) - 1/(1+10**(_PKA['C_term']-ph))
        for aa in seq.upper():
            if aa=='D': c -= 1/(1+10**(_PKA['D']-ph))
            elif aa=='E': c -= 1/(1+10**(_PKA['E']-ph))
            elif aa=='H': c += 1/(1+10**(ph-_PKA['H']))
            elif aa=='K': c += 1/(1+10**(ph-_PKA['K']))
            elif aa=='R': c += 1/(1+10**(ph-_PKA['R']))
        return c
    lo, hi = 0.0, 14.0
    for _ in range(200):
        mid = (lo+hi)/2
        if _c(mid) > 0: lo = mid
        else: hi = mid
    return round((lo+hi)/2, 3)

def _instability(seq):
    # Guruprasad dipeptide instability weights (abbreviated key set)
    DIWV = {'WC':1,'WM':24.68,'CM':33.6,'CH':33.6,'CD':20.26,'CT':33.6,
            'CL':20.26,'CP':20.26,'QD':20.26,'RD':20.26,'RH':20.26,
            'DG':20.26,'DS':20.26,'FD':54.96,'GD':20.26,'GH':33.6,
            'HN':20.26,'LR':20.26,'ND':20.26,'NL':20.26,'NP':20.26,
            'QR':20.26,'RR':58.28,'SR':20.26,'TD':20.26,'TP':20.26,
            'VN':1,'WR':1}
    seq = seq.upper()
    if len(seq) < 2: return 0.0
    total = sum(DIWV.get(seq[i:i+2], 1.0) for i in range(len(seq)-1))
    return round(10.0/len(seq)*total, 3)

def _strip_cdr3_loop(cdr3: str) -> str:
    """
    Strip conserved framework anchor residues from CDR3:
      N-terminal : remove leading 'CAR' or 'C' (conserved cysteine + Ala-Arg)
      C-terminal : remove trailing 'FDY', 'FDW', 'MDY', 'LDY', 'MDW' (J-gene motif)
    Returns only the hypervariable loop residues.

    Example:  CARGFDYW  →  strip CAR + FDY  →  G
              CARDRGYYY →  strip CAR        →  DRGYYY
              CARRGFDYW →  strip CAR + FDY  →  R
    """
    s = cdr3.upper().replace('-', '')
    # Strip N-terminal conserved residues: leading C or CAR
    if s.startswith('CAR'):
        s = s[3:]
    elif s.startswith('C'):
        s = s[1:]
    # Strip C-terminal J-gene motif: xDY or xDW (last 3 residues)
    # Common patterns: FDY, FDW, MDY, MDW, LDY, WDY, YDY
    if len(s) >= 3 and s[-2] == 'D' and s[-1] in 'YW':
        s = s[:-3]
    return s


def compute_biophysical_features(seq: str, feature_list: list,
                                  vh_seq: str = '') -> dict:
    """
    Compute biophysical features from CDR3 loop sequence.
    seq     : CDR3 with conserved flanks already stripped by _strip_cdr3_loop()
    vh_seq  : full VH sequence (for vh_length and ratio features)

    Feature naming convention:
      cdr3_{property}  e.g. cdr3_charge_ph7, cdr3_R, cdr3_length
      vh_length        length of full VH sequence
      cdr3_vh_length_ratio  CDR3 loop / VH length
    """
    seq = seq.upper().replace('-', '')
    if not seq:
        return {f: 0.0 for f in feature_list}
    counts = Counter(seq); n = len(seq)
    vh_len = len(vh_seq.upper().replace('-', '')) if vh_seq else 0
    out = {}
    for f in feature_list:
        if   f == 'length':               out[f] = float(n)
        elif f == 'pi':                   out[f] = _pi(seq)
        elif f in ('charge_ph7', 'charge'): out[f] = _charge_at_ph7(seq)
        elif f == 'hydrophobicity':       out[f] = _gravy(seq)
        elif f == 'aromaticity':          out[f] = round(sum(counts.get(a,0) for a in 'FYW')/n, 4)
        elif f == 'instability':          out[f] = _instability(seq)
        elif f == 'net_charge_sq':
            c = _charge_at_ph7(seq);      out[f] = round(c*c, 4)
        elif f == 'frac_charged':         out[f] = round(sum(counts.get(a,0) for a in 'DERKH')/n, 4)
        elif f == 'frac_hydrophobic':     out[f] = round(sum(counts.get(a,0) for a in 'AVILMFWP')/n, 4)
        elif f == 'cdr3_length':          out[f] = float(n)
        elif f == 'vh_length':            out[f] = float(vh_len)
        elif f == 'vh_cdr3_length_ratio': out[f] = round(vh_len / n, 4) if n > 0 else 0.0
        elif f in ('vh_charge_ph7', 'vh_charge'): out[f] = _charge_at_ph7(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f == 'vh_hydrophobicity':    out[f] = _gravy(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f.startswith('count_'):      out[f] = float(counts.get(f.split('_')[1], 0))
        # Single amino acid counts — e.g. 'R', 'K', 'W', 'D'
        elif len(f) == 1 and f in 'ACDEFGHIKLMNPQRSTVWY':
                                          out[f] = float(counts.get(f, 0))
        else:                             out[f] = 0.0
    return out

def compute_kmer_features(seq: str, k_list: list, normalize: bool=True) -> dict:
    seq = seq.upper().replace('-','')
    out = {}
    for k in k_list:
        all_kmers = [''.join(p) for p in product(AMINO_ACIDS, repeat=k)]
        cnts = Counter(seq[i:i+k] for i in range(max(0, len(seq)-k+1))) if len(seq)>=k else {}
        total = sum(cnts.values()) or 1
        for km in all_kmers:
            v = cnts.get(km, 0)
            out[f'{k}mer_{km}'] = v/total if normalize else float(v)
    return out



# ── ASCII lookup for fast one-hot (same as transformer_onehot.py) ─────────────
_AA_LOOKUP_RF = np.full(128, 20, dtype=np.int32)
for _aa_rf, _idx_rf in {aa: i for i, aa in enumerate(AMINO_ACIDS)}.items():
    _AA_LOOKUP_RF[ord(_aa_rf)] = _idx_rf


def compute_onehot_features(seq: str, max_len: int) -> np.ndarray:
    """
    One-hot encode a sequence into a flattened vector of shape (max_len * 20,).
    Unknown / padding positions → all-zero row.
    """
    seq = seq.upper().replace('-', '')[:max_len]
    out = np.zeros((max_len, len(AMINO_ACIDS)), dtype=np.float32)
    if seq:
        chars  = np.frombuffer(seq.encode('ascii'), dtype=np.uint8)
        aa_idx = _AA_LOOKUP_RF[np.clip(chars, 0, 127)]
        pos    = np.arange(len(chars))
        valid  = aa_idx < 20
        out[pos[valid], aa_idx[valid]] = 1.0
    return out.flatten()   # (max_len * 20,)


def onehot_feature_names(prefix: str, max_len: int) -> list:
    """
    Generate feature names for a one-hot encoded sequence.
    Format: {prefix}_pos{N}_{AA}  e.g. H_pos001_A, H_pos001_C, ...
    """
    names = []
    for pos in range(1, max_len + 1):
        for aa in AMINO_ACIDS:
            names.append(f"{prefix}_pos{pos:03d}_{aa}")
    return names


# ══════════════════════════════════════════════════════════════════════════════
# 3.  FEATURE BUILDER
# ══════════════════════════════════════════════════════════════════════════════

class FeatureBuilder:
    """
    Builds combined feature matrix for RandomForestModel.

    Supported feature types (any combination):
      • embedding   – PLM vector (ablang 480d, antiberta2 1024d, ...)
      • kmer        – k-mer frequencies from VH | CDR3 | VHVL
      • biophysical – pI, charge, hydrophobicity, AA counts, VH-level
      • onehot      – position × AA one-hot from HCDR3 | VH | VHVL

    YAML keys:
      features.embedding / kmer / biophysical / onehot  : bool
      kmer.sequence     : CDR3 | VH | VHVL
      kmer.k            : [1, 2]
      kmer.normalize    : true
      biophysical.features : [length, pi, charge_ph7, ...]
      onehot.sequence   : HCDR3 | VH | VHVL  (default HCDR3)
    """

    def __init__(self, config: dict):
        self.cfg        = config
        self.feat_cfg   = config.get('features', {})
        self.km_cfg     = config.get('kmer',     {})
        self.bp_cfg     = config.get('biophysical', {})
        self.oh_cfg     = config.get('onehot',   {})
        # Ensure onehot key exists in feat_cfg — may be absent from old YAML
        self.feat_cfg.setdefault('onehot', False)

        self.feature_names_ : list = None
        self._emb_dim               = None

        # Max sequence lengths — fitted from training data
        self._oh_max_vh   = 0
        self._oh_max_vl   = 0
        self._oh_max_cdr3 = 0

    # ── Sequence accessors ────────────────────────────────────────────────────

    def _get_kmer_seq(self, row) -> str:
        """Sequence used for k-mer features (kmer.sequence in YAML)."""
        src = self.km_cfg.get('sequence', 'CDR3').upper()
        vh  = str(row.get('HSEQ', '') or '')
        vl  = str(row.get('LSEQ', '') or '')
        cdr = str(row.get('CDR3', '') or '')
        if src == 'VH':   return vh
        if src == 'VHVL': return vh + vl
        return cdr  # default CDR3

    def _get_cdr3(self, row) -> str:
        return str(row.get('CDR3', '') or '')

    def _get_vh(self, row) -> str:
        return str(row.get('HSEQ', '') or '')

    def _get_vl(self, row) -> str:
        return str(row.get('LSEQ', '') or '')

    def _oh_segs(self, row) -> list:
        """
        Return list of (seq, max_len, prefix) for one-hot encoding.
        Controlled by onehot.sequence in YAML: HCDR3 | VH | VHVL
        """
        mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
        vh   = self._get_vh(row)
        vl   = self._get_vl(row)
        cdr  = self._get_cdr3(row)
        if mode == 'VH':
            return [(vh, self._oh_max_vh, 'oh_vh')]
        elif mode == 'VHVL':
            return [(vh, self._oh_max_vh, 'oh_vh'),
                    (vl, self._oh_max_vl, 'oh_vl')]
        else:  # HCDR3 (default)
            return [(cdr, self._oh_max_cdr3, 'oh_hcdr3')]

    # ── Chain tag for output file naming ─────────────────────────────────────

    @property
    def chain_tag(self) -> str:
        """
        Short tag added to output filenames for kmer/onehot modes.
        e.g. 'HCDR3', 'VH', 'VHVL'
        """
        if self.feat_cfg.get('kmer') and not self.feat_cfg.get('onehot'):
            return self.km_cfg.get('sequence', 'CDR3').upper()
        if self.feat_cfg.get('onehot'):
            return self.oh_cfg.get('sequence', 'HCDR3').upper()
        return ''

    # ── Fit ───────────────────────────────────────────────────────────────────

    def fit(self, X_df: pd.DataFrame, embeddings: np.ndarray = None):
        names = []

        # Embedding
        if self.feat_cfg.get('embedding'):
            if embeddings is None:
                raise ValueError(
                    "features.embedding=True but embeddings=None.\n"
                    "  Pass embeddings=X.values to train()/kfold_validation().\n"
                    "  Or set features.embedding: false in YAML.")
            self._emb_dim = embeddings.shape[1]
            names += [f'emb_{i}' for i in range(self._emb_dim)]

        # K-mer
        if self.feat_cfg.get('kmer'):
            kf    = compute_kmer_features(
                        self._get_kmer_seq(X_df.iloc[0]),
                        self.km_cfg.get('k', [1, 2]),
                        self.km_cfg.get('normalize', True))
            names += list(kf.keys())

        # Biophysical
        if self.feat_cfg.get('biophysical'):
            bf = compute_biophysical_features(
                     self._get_cdr3(X_df.iloc[0]),
                     self.bp_cfg.get('features', []),
                     vh_seq=self._get_vh(X_df.iloc[0]))
            def _rename(k):
                k = k.replace('count_', '').replace('_ph7', '')
                if k.startswith('vh_') or k.startswith('cdr3_'):
                    return k
                return f'cdr3_{k}'
            names += [_rename(k) for k in bf.keys()]

        # One-hot — fit max lengths from training data, capped by YAML max_lengths.
        # Priority:
        #   1. YAML onehot.max_lengths.{vh|vl|hcdr3} — hard cap (safety ceiling)
        #   2. Actual max observed in training data    — fitted value (stored in pkl)
        #   3. Fallback defaults (135/135/25)
        #
        # At predict time the saved _oh_max_* values from the pkl are used directly,
        # so predict sequences longer than training max are truncated with a warning.
        if self.feat_cfg.get('onehot'):
            mode     = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _max_cfg = self.oh_cfg.get('max_lengths', {})
            # YAML caps (0 or missing → no cap applied)
            _cap_vh   = int(_max_cfg.get('vh',   0) or 0)
            _cap_vl   = int(_max_cfg.get('vl',   0) or 0)
            _cap_cdr3 = int(_max_cfg.get('hcdr3',0) or 0)

            if mode in ('VH', 'VHVL'):
                _data_max_vh   = max((len(str(r.get('HSEQ','') or ''))
                                      for _, r in X_df.iterrows()), default=135)
                self._oh_max_vh  = _cap_vh  if _cap_vh  > 0 else _data_max_vh
                print(f"  [onehot] VH  max_len={self._oh_max_vh:3d}  "
                      f"(data_max={_data_max_vh}"
                      + (f"  capped at {_cap_vh}" if _cap_vh > 0 else "") + ")")
            if mode == 'VHVL':
                _data_max_vl   = max((len(str(r.get('LSEQ','') or ''))
                                      for _, r in X_df.iterrows()), default=135)
                self._oh_max_vl  = _cap_vl  if _cap_vl  > 0 else _data_max_vl
                print(f"  [onehot] VL  max_len={self._oh_max_vl:3d}  "
                      f"(data_max={_data_max_vl}"
                      + (f"  capped at {_cap_vl}" if _cap_vl > 0 else "") + ")")
            if mode == 'HCDR3':
                _data_max_cdr3 = max((len(str(r.get('CDR3','') or ''))
                                      for _, r in X_df.iterrows()), default=25)
                self._oh_max_cdr3 = _cap_cdr3 if _cap_cdr3 > 0 else _data_max_cdr3
                print(f"  [onehot] CDR3 max_len={self._oh_max_cdr3:3d}  "
                      f"(data_max={_data_max_cdr3}"
                      + (f"  capped at {_cap_cdr3}" if _cap_cdr3 > 0 else "") + ")")

            # Feature names
            row0 = X_df.iloc[0]
            for _seq, max_len, pfx in self._oh_segs(row0):
                for pos in range(1, max_len + 1):
                    for aa in AMINO_ACIDS:
                        names.append(f"{pfx}_{pos}_{aa}")

        if not names:
            raise ValueError(
                "No features selected — enable at least one of: "
                "embedding | kmer | biophysical | onehot")
        self.feature_names_ = names

        # ── Print feature dimension report ────────────────────────────────────
        _AA = len(AMINO_ACIDS)   # 20
        _SEP = '─' * 52
        print(f"[FeatureBuilder] Feature dimensions:")
        print(f"  {_SEP}")
        _total = 0

        if self.feat_cfg.get('embedding') and self._emb_dim:
            _d = self._emb_dim
            print(f"  embedding       : {_d:>6,} d  (PLM vector)")
            _total += _d

        if self.feat_cfg.get('kmer'):
            _k   = self.km_cfg.get('k', [1, 2])
            _seq = self.km_cfg.get('sequence', 'CDR3')
            _d   = len([n for n in names if 'mer_' in n])
            print(f"  kmer            : {_d:>6,} d  "
                  f"(k={_k}  seq={_seq})")
            _total += _d

        if self.feat_cfg.get('biophysical'):
            _d = len([n for n in names
                      if n.startswith('cdr3_') or n.startswith('vh_')])
            print(f"  biophysical     : {_d:>6,} d  "
                  f"(CDR3 + VH properties)")
            _total += _d

        if self.feat_cfg.get('onehot'):
            _mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            if _mode == 'VHVL':
                _vh = self._oh_max_vh
                _vl = self._oh_max_vl
                _d  = (_vh + _vl) * _AA
                print(f"  onehot (VHVL)   : {_d:>6,} d  "
                      f"= 2 chains × pos × 20 AA")
                print(f"    VH  : {_vh:>3} pos × {_AA} AA = {_vh*_AA:,}")
                print(f"    VL  : {_vl:>3} pos × {_AA} AA = {_vl*_AA:,}")
            elif _mode == 'VH':
                _vh = self._oh_max_vh
                _d  = _vh * _AA
                print(f"  onehot (VH)     : {_d:>6,} d  "
                      f"= {_vh} pos × {_AA} AA")
            else:   # HCDR3
                _cdr3 = self._oh_max_cdr3
                _d    = _cdr3 * _AA
                print(f"  onehot (HCDR3)  : {_d:>6,} d  "
                      f"= {_cdr3} pos × {_AA} AA")
            _total += _d

        print(f"  {_SEP}")
        print(f"  TOTAL           : {len(names):>6,} d")
        print(f"  {_SEP}")
        return self

    # ── Transform ─────────────────────────────────────────────────────────────

    def transform(self, X_df: pd.DataFrame,
                  embeddings: np.ndarray = None) -> np.ndarray:
        parts = []
        n     = len(X_df)

        if self.feat_cfg.get('embedding') and embeddings is not None:
            parts.append(embeddings.astype(np.float32))

        if self.feat_cfg.get('kmer'):
            k_list = self.km_cfg.get('k', [1, 2])
            norm   = self.km_cfg.get('normalize', True)
            rows   = [list(compute_kmer_features(
                               self._get_kmer_seq(X_df.iloc[i]),
                               k_list, norm).values())
                      for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))

        if self.feat_cfg.get('biophysical'):
            fl   = self.bp_cfg.get('features', [])
            rows = [list(compute_biophysical_features(
                            self._get_cdr3(X_df.iloc[i]),
                            fl,
                            vh_seq=self._get_vh(X_df.iloc[i])).values())
                    for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))

        if self.feat_cfg.get('onehot'):
            # Warn if any predict sequence exceeds the training max length
            # (sequences will be silently truncated — could shift residue positions)
            mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _warned = False
            for i in range(min(n, 5)):   # check first 5 rows only for speed
                row = X_df.iloc[i]
                if mode in ('VH', 'VHVL') and self._oh_max_vh > 0:
                    _l = len(str(row.get('HSEQ','') or ''))
                    if _l > self._oh_max_vh and not _warned:
                        print(f"  [onehot] WARNING: VH sequence length {_l} > "
                              f"training max {self._oh_max_vh} — truncated. "
                              f"Consider increasing onehot.max_lengths.vh in YAML.")
                        _warned = True
                if mode == 'VHVL' and self._oh_max_vl > 0:
                    _l = len(str(row.get('LSEQ','') or ''))
                    if _l > self._oh_max_vl and not _warned:
                        print(f"  [onehot] WARNING: VL sequence length {_l} > "
                              f"training max {self._oh_max_vl} — truncated. "
                              f"Consider increasing onehot.max_lengths.vl in YAML.")
                        _warned = True
                if mode == 'HCDR3' and self._oh_max_cdr3 > 0:
                    _l = len(str(row.get('CDR3','') or ''))
                    if _l > self._oh_max_cdr3 and not _warned:
                        print(f"  [onehot] WARNING: CDR3 length {_l} > "
                              f"training max {self._oh_max_cdr3} — truncated. "
                              f"Re-train with this data or increase max_lengths.hcdr3.")
                        _warned = True
            oh_rows = []
            for i in range(n):
                row  = X_df.iloc[i]
                vecs = [compute_onehot_features(seq, max_len)
                        for seq, max_len, _ in self._oh_segs(row)]
                oh_rows.append(np.concatenate(vecs))
            parts.append(np.array(oh_rows, dtype=np.float32))

        return np.hstack(parts) if parts else np.zeros((n, 0), dtype=np.float32)

    def fit_transform(self, X_df, embeddings=None):
        return self.fit(X_df, embeddings).transform(X_df, embeddings)

    # ── Properties ────────────────────────────────────────────────────────────

    @property
    def n_features(self):
        return len(self.feature_names_ or [])

    @property
    def non_embedding_feature_names(self):
        """kmer + biophysical + onehot names — interpretable, used for SHAP."""
        return [n for n in (self.feature_names_ or [])
                if not n.startswith('emb_')]

    @property
    def non_embedding_indices(self):
        return [i for i, n in enumerate(self.feature_names_ or [])
                if not n.startswith('emb_')]

    @property
    def onehot_feature_names(self):
        """One-hot feature names (oh_hcdr3_1_A, oh_vh_1_A, ...)."""
        return [n for n in (self.feature_names_ or [])
                if n.startswith('oh_')]

    @property
    def onehot_indices(self):
        return [i for i, n in enumerate(self.feature_names_ or [])
                if n.startswith('oh_')]


# ══════════════════════════════════════════════════════════════════════════════
# 4.  HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _deep_merge(base: dict, override: dict) -> dict:
    for k, v in override.items():
        if isinstance(v, dict) and isinstance(base.get(k), dict):
            _deep_merge(base[k], v)
        else:
            base[k] = v
    return base

def _setup_logger(log_path: str) -> logging.Logger:
    logger = logging.getLogger(f'RF_{os.path.basename(log_path)}')
    logger.setLevel(logging.DEBUG)
    if not logger.handlers:
        fh = logging.FileHandler(log_path, mode='a', encoding='utf-8')
        fh.setFormatter(logging.Formatter('%(message)s'))
        logger.addHandler(fh)
    return logger

def _log(logger, msg: str):
    print(msg)
    if logger: logger.info(msg)


# ══════════════════════════════════════════════════════════════════════════════
# 5.  MAIN CLASS
# ══════════════════════════════════════════════════════════════════════════════

class RandomForestModel:
    """
    Random Forest for antibody developability — classification & regression.

    Feature modes (combinable via YAML):
      embedding   – PLM embedding vector
      kmer        – 1/2/3-mer frequencies from VH / CDR3
      biophysical – pI, charge, CDR3 length, AA counts, hydrophobicity …

    Interpretation:
      SHAP TreeExplainer for kmer + biophysical features.
      (For pure embedding, use IG from transformer_lm instead.)

    Examples
    --------
    # Classification
    model = RandomForestModel()
    RandomForestModel.kfold_validation(data, X, y, emb, target='psr_filter')

    # Regression
    cfg = RandomForestModel.auto_detect_config(n=11000, pos_rate=0.5,
                                               task='regression')
    model = RandomForestModel(config=cfg)
    RandomForestModel.kfold_validation(data, X, y, emb, target='psr_norm_mean')

    # k-mer only (no PLM needed)
    cfg['features'] = {'embedding': False, 'kmer': True, 'biophysical': True}
    model = RandomForestModel(config=cfg)
    """

    _DEFAULT_CONFIG = {
        'mode': 'manual',
        'task': 'classification',
        'features': {'embedding': True, 'kmer': False, 'biophysical': False},
        'kmer': {'k': [1,2,3], 'sequence': 'CDR3', 'normalize': True},
        'biophysical': {
            'sequence': 'CDR3',
            'features': [
                # CDR3 loop only (CAR/FDY flanks stripped)
                'length', 'pi', 'charge_ph7', 'hydrophobicity',
                'aromaticity', 'instability',
                'frac_charged', 'frac_hydrophobic',
                # Single AA — anionic (protective)
                'D', 'E',
                # Single AA — cationic (polyreactivity drivers)
                'R', 'K', 'H',
                # Single AA — aromatic
                'W', 'Y', 'F',
                # Single AA — structural
                'G', 'P', 'C',
                # VH-level features
                'vh_length', 'vh_cdr3_length_ratio',
                'vh_charge', 'vh_hydrophobicity',
            ],
        },
        'model': {
            'n_estimators': 1000, 'max_depth': 15,
            'min_samples_leaf': 5, 'min_samples_split': 10,
            'max_features': 0.3, 'criterion': 'entropy',
            'bootstrap': True, 'oob_score': True,
            'n_jobs': -1, 'random_state': 42,
        },
        'training': {
            'class_weight': 'balanced',
            'hyperparam_search': False,
            'search_n_iter': 20, 'search_cv': 3,
        },
        'shap': {
            'enabled': True, 'max_samples': 500,
            'top_features': 30, 'plot_types': ['bar','beeswarm'],
        },
    }

    # ── auto-detect ───────────────────────────────────────────────────────────

    @staticmethod
    def auto_detect_config(n: int, pos_rate: float,
                           task: str = 'classification') -> dict:
        """
        Derive RF config from dataset size, class balance and task.

        Size tiers:
          xs  n <  5k  →  500 trees, depth 10, min_leaf 3,  max_feat sqrt
          sm  5–20k    → 1000 trees, depth 15, min_leaf 5,  max_feat 0.3
          md  20–80k   → 2000 trees, depth 20, min_leaf 8,  max_feat 0.2
          lg  80–200k  → 2000 trees, depth 25, min_leaf 10, max_feat 0.15
          xl  >200k    → 3000 trees, depth 30, min_leaf 15, max_feat 0.1
        """
        if   n <  5_000: tier='xs'; n_e,d,ml,mf = 500, 10, 3,  'sqrt'
        elif n < 20_000: tier='sm'; n_e,d,ml,mf = 1000,15, 5,  0.3
        elif n < 80_000: tier='md'; n_e,d,ml,mf = 2000,20, 8,  0.2
        elif n <200_000: tier='lg'; n_e,d,ml,mf = 2000,25,10,  0.15
        else:            tier='xl'; n_e,d,ml,mf = 3000,30,15,  0.1

        criterion = 'squared_error' if task=='regression' else 'entropy'
        if task == 'regression':
            cw = None
        else:
            min_rate = min(pos_rate, 1-pos_rate)
            if   min_rate >= 0.40: cw = 'balanced'
            elif min_rate >= 0.10: cw = 'balanced_subsample'
            else:
                ratio = (1-pos_rate)/(pos_rate+1e-8)
                cw    = {0:1, 1:ratio}
                print(f"  [WARN] Severe imbalance → explicit weights {cw}")

        cfg = copy.deepcopy(RandomForestModel._DEFAULT_CONFIG)
        cfg['mode'] = 'auto'; cfg['task'] = task
        cfg['model'].update({'n_estimators':n_e,'max_depth':d,
                             'min_samples_leaf':ml,'max_features':mf,
                             'criterion':criterion})
        cfg['training']['class_weight'] = cw
        cfg['_auto'] = {'n':n,'pos_rate':round(pos_rate,4),'size_tier':tier,'task':task}
        return cfg

    @staticmethod
    def print_config_report(cfg: dict) -> None:
        a  = cfg.get('_auto',{}); m = cfg['model']
        t  = cfg['training'];     ft= cfg.get('features',{})
        km = cfg.get('kmer',{});  bp= cfg.get('biophysical',{})
        sh = cfg.get('shap',{});  task = cfg.get('task','classification')
        W=62; sep='═'*W; sep2='─'*W
        print(f"\n{sep}")
        if a:
            print(f"  RandomForest  ·  AUTO  ({task})  tier={a.get('size_tier')}  n={a.get('n'):,}")
        else:
            print(f"  RandomForest  ·  MANUAL  ({task})")
        print(sep2)
        print(f"  MODEL          n_est={m['n_estimators']}  depth={m['max_depth']}  "
              f"min_leaf={m['min_samples_leaf']}  max_feat={m['max_features']}")
        print(f"                 criterion={m['criterion']}  oob={m['oob_score']}")
        print(sep2)
        print(f"  FEATURES       embedding={ft.get('embedding')}  "
              f"kmer={ft.get('kmer')}  biophysical={ft.get('biophysical')}  "
              f"onehot={ft.get('onehot', False)}")
        if ft.get('kmer'):
            print(f"                 k-mer k={km.get('k')}  seq={km.get('sequence')}")
        if ft.get('biophysical'):
            print(f"                 bio features={len(bp.get('features',[]))}  "
                  f"seq={bp.get('sequence')}")
        print(sep2)
        if task == 'regression':
            print(f"  TRAINING       search={t.get('hyperparam_search')}")
        else:
            print(f"  TRAINING       class_weight={t.get('class_weight')}  "
                  f"search={t.get('hyperparam_search')}")
        print(f"  SHAP           enabled={sh.get('enabled')}  "
              f"top={sh.get('top_features')}  plots={sh.get('plot_types')}")
        print(f"{sep}\n")

    # ── init ──────────────────────────────────────────────────────────────────

    def __init__(self, config_path: str = "config/random_forest.yaml",
                 config: dict = None, verbose: bool = True):
        if config is not None:
            self.config = _deep_merge(copy.deepcopy(self._DEFAULT_CONFIG), config)
            mode = "auto" if '_auto' in config else "dict"
            if verbose:
                print(f"[RandomForestModel] config ← {mode}")
        elif os.path.exists(config_path):
            with open(config_path) as f:
                user_cfg = yaml.safe_load(f) or {}
            self.config = _deep_merge(copy.deepcopy(self._DEFAULT_CONFIG), user_cfg)
            if verbose:
                print(f"[RandomForestModel] config ← {config_path}")
        else:
            self.config = copy.deepcopy(self._DEFAULT_CONFIG)
            if verbose:
                print(f"[RandomForestModel] {config_path} not found — using defaults")
        self.model = None
        self.fb_   = None
        self.task  = self.config.get('task', 'classification')
        if verbose:
            RandomForestModel.print_config_report(self.config)

    # ── internal helpers ──────────────────────────────────────────────────────

    def _resolve_config(self, y: np.ndarray) -> None:
        if self.config.get('mode','manual') != 'auto': return
        n = len(y)
        pos_rate = float(np.nanmean(y)) if self.task=='classification' \
                   else float(np.mean(y > np.median(y)))
        resolved = RandomForestModel.auto_detect_config(n, pos_rate, self.task)
        for s in ('features','kmer','biophysical','shap'):
            if s in self.config: resolved[s] = copy.deepcopy(self.config[s])
        self.config = resolved
        self.task   = self.config.get('task', self.task)   # keep task in sync
        RandomForestModel.print_config_report(self.config)

    def apply_lm_profile(self, lm: str, logger=None) -> None:
        """
        Apply LM-specific hyperparameter overrides from lm_profiles in YAML.
        Called before _build_sklearn_model() with the --lm name.

        Priority: lm_profiles.{lm} > model defaults
        Only keys present in the profile are overridden.
        """
        profiles = self.config.get('lm_profiles', {})
        profile  = profiles.get(lm, {})
        if not profile:
            return
        m = self.config['model']
        for k, v in profile.items():
            m[k] = v
        _log(logger, f"[RF] lm_profile '{lm}' applied: "
                     + "  ".join(f"{k}={v}" for k, v in profile.items()))

    def _build_sklearn_model(self):
        m  = self.config['model']
        t  = self.config['training']
        kw = dict(n_estimators=m['n_estimators'], max_depth=m['max_depth'],
                  min_samples_leaf=m['min_samples_leaf'],
                  min_samples_split=m['min_samples_split'],
                  max_features=m['max_features'], bootstrap=m['bootstrap'],
                  oob_score=m['oob_score'], n_jobs=m['n_jobs'],
                  random_state=m['random_state'])
        if self.task == 'classification':
            kw['criterion']    = m.get('criterion','entropy')
            kw['class_weight'] = t.get('class_weight','balanced')
            return RandomForestClassifier(**kw)
        else:
            kw['criterion'] = m.get('criterion','squared_error')
            return RandomForestRegressor(**kw)

    def _hyperparam_search(self, rf, X, y, logger=None):
        t = self.config['training']
        param_dist = {
            'n_estimators':      [500,1000,1500,2000],
            'max_depth':         [8,10,12,15,20,None],
            'max_features':      ['sqrt',0.2,0.3,0.4,0.5],
            'min_samples_leaf':  randint(2,15),
            'min_samples_split': randint(5,30),
        }
        scoring = 'r2' if self.task=='regression' else 'roc_auc'
        cv = (t['search_cv'] if self.task=='regression'
              else StratifiedKFold(n_splits=t['search_cv'],shuffle=True,random_state=42))
        search = RandomizedSearchCV(rf, param_dist, n_iter=t['search_n_iter'],
                                    cv=cv, scoring=scoring, n_jobs=-1,
                                    random_state=42, verbose=0)
        search.fit(X, y)
        _log(logger, f"  [search] best_params={search.best_params_}")
        _log(logger, f"  [search] CV {scoring}={search.best_score_:.4f}")
        return search.best_estimator_

    # ── _auto_fix_features ───────────────────────────────────────────────────

    def _auto_fix_features(self, embeddings, X_df=None):
        """
        Auto-correct features config so users do not need to edit YAML.

        Rules (applied in order):
          1. embedding=True  but embeddings=None
             → seq-only --lm was used without overriding YAML
             → detect which seq features are available and enable them
               • CDR3 column present  → biophysical=True (safest default)
               • all features off     → biophysical=True as fallback
             Prints a clear warning so the user knows what happened.

          2. embedding=True  and embeddings provided
             → PLM mode — leave as-is.

          3. onehot=True and CDR3/HSEQ column absent from X_df
             → warn, disable onehot.

        Called automatically at the start of train() and kfold_validation().
        """
        f = self.config.setdefault('features', {})

        # Rule 1: embedding=True but no embeddings array
        if f.get('embedding') and embeddings is None:
            # Figure out the best seq-only fallback
            _has_cdr3 = (X_df is not None and 'CDR3' in X_df.columns)
            _any_seq  = f.get('kmer') or f.get('biophysical') or f.get('onehot')

            f['embedding'] = False          # must disable — embeddings=None

            if not _any_seq:
                # Nothing else enabled — default to biophysical
                f['biophysical'] = True
                print("[RF] WARNING: embedding=True but no embeddings provided.")
                print("  Auto-corrected: embedding=False  biophysical=True")
                print("  → To use PLM: provide embeddings via --lm ablang/antiberta2/...")
                print("  → To use seq-only: set --lm biophysical or --lm kmer in YAML")
            else:
                print(f"[RF] WARNING: embedding=True but no embeddings provided → embedding=False")
                print(f"  Seq features active: kmer={f.get('kmer')}  "
                      f"bio={f.get('biophysical')}  onehot={f.get('onehot')}")

    # ── train ─────────────────────────────────────────────────────────────────

    def train(self, X_df: pd.DataFrame, y,
              embeddings: np.ndarray = None,
              val_X: pd.DataFrame = None,
              val_y = None,
              val_embeddings: np.ndarray = None,
              target: str = "model",
              target_col: str = "",
              label_col: str = "",
              embedding_lm: str = "",
              logger=None):
        """
        Train on full dataset, with optional post-fit val evaluation.

        X_df           - DataFrame with HSEQ / CDR3 columns
        y              - int labels (classification) or float scores (regression)
        embeddings     - (n, emb_dim) PLM array; None if embedding disabled
        val_X          - validation DataFrame (same columns as X_df)
        val_y          - validation labels / scores
        val_embeddings - (n_val, emb_dim) validation embeddings
        NOTE: RF has no early stopping. val data is used only for post-training
              evaluation and reporting. Use kfold_validation() for model selection.
        """
        self._resolve_config(np.asarray(y, dtype=float))
        # Auto-detect regression if YAML says classification but y is continuous
        _y_tmp_rf   = np.asarray(y, dtype=float)
        _unique_rf  = len(set(_y_tmp_rf.tolist()))
        _is_bin_rf  = (_unique_rf <= 2 and
                       set(_y_tmp_rf.tolist()).issubset({0, 1, 0.0, 1.0}))
        if self.task == 'classification' and not _is_bin_rf:
            self.task = 'regression'
            self.config['task'] = 'regression'
            print(f"[train] Auto-detected task=regression "
                  f"({_unique_rf} unique values — not binary)")
        y = np.asarray(y, dtype=float if self.task=='regression' else int)

        # ── Safety: auto-correct feature flags from available data ───────────
        # If embedding=True but no embeddings provided → likely a seq-only --lm
        # that wasn't overridden in the calling script. Auto-detect and fix.
        self._auto_fix_features(embeddings, X_df)

        self.fb_ = FeatureBuilder(self.config)
        self.fb_.fit(X_df, embeddings)
        # Apply LM-specific hyperparameter profile BEFORE building the model
        self.apply_lm_profile(embedding_lm, logger=logger)
        X_feat   = self.fb_.transform(X_df, embeddings)

        _log(logger, f"[train] n={len(y):,}  features={X_feat.shape[1]}  task={self.task}")
        _log(logger, f"[train] features: emb={self.config['features'].get('embedding')}  "
                     f"kmer={self.config['features'].get('kmer')}  "
                     f"bio={self.config['features'].get('biophysical')}")

        rf = self._build_sklearn_model()
        if self.config['training'].get('hyperparam_search'):
            _log(logger, "  [search] Running RandomizedSearchCV...")
            self.model = self._hyperparam_search(rf, X_feat, y, logger)
        else:
            self.model = rf
            self.model.fit(X_feat, y)

        if getattr(self.model, 'oob_score_', None) is not None:
            _log(logger, f"  OOB score = {self.model.oob_score_:.4f}")

        # Post-training val evaluation (no early stopping — one-shot after fit)
        if val_X is not None and val_y is not None:
            val_y_arr  = np.asarray(val_y, dtype=float if self.task=='regression' else int)
            X_val_feat = self.fb_.transform(val_X, val_embeddings)
            _log(logger, f"\n[train] Val evaluation (n={len(val_y_arr):,}  "
                         + (f"pos={val_y_arr.mean():.1%}" if self.task=='classification' else f"mean={val_y_arr.mean():.3f}") + "):")

            if self.task == 'classification':
                from sklearn.metrics import roc_auc_score, accuracy_score, f1_score, recall_score
                val_probs = self.model.predict_proba(X_val_feat)[:, 1]
                val_preds = (val_probs >= 0.5).astype(int)
                try:
                    val_auc = roc_auc_score(val_y_arr, val_probs)
                except Exception:
                    val_auc = 0.5
                val_acc  = accuracy_score(val_y_arr, val_preds)
                val_f1   = f1_score(val_y_arr, val_preds, zero_division=0)
                val_rcf  = recall_score(val_y_arr, val_preds, pos_label=0, zero_division=0)
                _log(logger, f"  val_auc={val_auc:.4f}  val_acc={val_acc:.4f}  "
                             f"val_f1={val_f1:.4f}  val_rec_fail={val_rcf:.4f}")
                _log(logger, f"  OOB score is the RF equivalent of val_auc during training.")
            else:
                from sklearn.metrics import r2_score, mean_absolute_error
                from scipy.stats import pearsonr, spearmanr
                val_preds_r = self.model.predict(X_val_feat)
                val_r2  = r2_score(val_y_arr, val_preds_r)
                val_mae = mean_absolute_error(val_y_arr, val_preds_r)
                val_rp  = pearsonr(val_y_arr, val_preds_r)[0]
                val_rs  = spearmanr(val_y_arr, val_preds_r)[0]
                _log(logger, f"  val_r2={val_r2:.4f}  val_pearson={val_rp:.4f}  "
                             f"val_spearman={val_rs:.4f}  val_mae={val_mae:.4f}")

        _log(logger, "[train] completed.")

        # ── Auto SHAP after training ──────────────────────────────────────────
        if self.config.get('shap', {}).get('enabled', False) and _SHAP_AVAILABLE:
            sh_top  = self.config.get('shap', {}).get('top_features', 50)
            _prefix = target  # full stem e.g. "psr_filter_biophysical_rf_ipi_psr_trainset"
            # Parse lm and db from prefix: {target_col}_{lm}_{model}_{db}
            _parts  = str(target).split('_')
            # lm is typically at index 1, but safest to use config
            _lm_tag = self.config.get('features', {}).get('_lm_name', '')
            # If not stored in config, try to parse from prefix
            if not _lm_tag and len(_parts) > 1:
                _lm_tag = _parts[1]
            _db_tag = '_'.join(_parts[3:]) if len(_parts) > 3 else _prefix

            _log(logger, f"\n[SHAP] Running on TRAIN set (top {sh_top}) ...")
            # Parse lm_name and db_name from the full target stem
            # stem format: {target}_{lm}_{model}_{db_stem}
            # e.g. "psr_filter_biophysical_rf_ipi_psr_trainset_train"
            # lm_name is the embedding/feature mode passed by caller
            _stem_parts = str(_prefix).split('_rf_')  # split at _rf_
            _shap_lm    = _stem_parts[0].split('_')[-1] if len(_stem_parts) > 1                           else self.config.get('_lm_name', '')
            _shap_db    = _stem_parts[1] if len(_stem_parts) > 1 else _prefix
            try:
                self.shap_analysis(X_df, embeddings,
                                   output_prefix   = _prefix,
                                   split_tag       = "train",
                                   top_n           = sh_top,
                                   barcodes        = list(X_df.index.astype(str)) if hasattr(X_df, 'index') else None,
                                   actual_labels   = list(y),
                                   actual_col_name = str(target).split('_')[0] if '_' in str(target) else target,
                                   lm_name         = _shap_lm,
                                   db_name         = _shap_db,
                                   logger          = logger)
            except Exception as _se:
                import traceback as _tb
                _log(logger, f"[SHAP] train SHAP failed — {_se}\n{_tb.format_exc()}")

            if val_X is not None and val_y is not None:
                _log(logger, f"\n[SHAP] Running on VAL set (top {sh_top}) ...")
                try:
                    self.shap_analysis(val_X, val_embeddings,
                                       output_prefix   = _prefix,
                                       split_tag       = "val",
                                       top_n           = sh_top,
                                       barcodes        = list(val_X.index.astype(str)) if hasattr(val_X, 'index') else None,
                                       actual_labels   = list(val_y) if val_y is not None else None,
                                       actual_col_name = str(target).split('_')[0] if '_' in str(target) else target,
                                       lm_name         = _shap_lm,
                                       db_name         = _shap_db,
                                       logger          = logger)
                except Exception as _se:
                    import traceback as _tb
                    _log(logger, f"[SHAP] val SHAP failed — {_se}\n{_tb.format_exc()}")
        elif not _SHAP_AVAILABLE:
            _log(logger, "[SHAP] skipped — pip install shap to enable")

        return self

    # ── predict ───────────────────────────────────────────────────────────────

    def predict_proba(self, X_df, embeddings=None) -> np.ndarray:
        if self.task == 'regression':
            raise RuntimeError("predict_proba not available for regression. Use predict().")
        return self.model.predict_proba(self.fb_.transform(X_df, embeddings))[:,1]

    def predict(self, X_df, embeddings=None, threshold: float = None) -> np.ndarray:
        X_feat = self.fb_.transform(X_df, embeddings)
        if self.task == 'regression':
            return self.model.predict(X_feat)
        t = threshold if threshold is not None else getattr(self,'recommended_threshold',0.5)
        return (self.model.predict_proba(X_feat)[:,1] >= t).astype(int)

    # ── save / load ───────────────────────────────────────────────────────────

    def save(self, path: str):
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        rt = getattr(self,'recommended_threshold',None)
        joblib.dump({'model':self.model,'feature_builder':self.fb_,
                     'config':self.config,'task':self.task,
                     'recommended_threshold':rt}, path)
        note = f"  threshold={rt:.4f}" if rt else "  threshold=None"
        print(f"[save] → {path}{note}")

    @classmethod
    def load(cls, path, config_path="config/random_forest.yaml"):
        inst = cls(config_path)
        pl   = joblib.load(path)
        if isinstance(pl, dict) and 'model' in pl:
            inst.model  = pl['model']
            inst.fb_    = pl.get('feature_builder')
            inst.config = _deep_merge(inst.config, pl.get('config',{}))
            inst.task   = pl.get('task','classification')
            rt = pl.get('recommended_threshold')
            inst.recommended_threshold = float(rt) if rt else 0.5
            flag = f"={rt:.4f}" if rt and rt!=0.5 else "=0.5 (default)"
            print(f"[load] recommended_threshold{flag}")
        else:
            inst.model = pl
            inst.recommended_threshold = 0.5
        print(f"[load] ← {path}")
        return inst

    # ── CDR3 in-silico mutagenesis ──────────────────────────────────────────────

    def cdr3_mutagenesis(self, X_df, embeddings=None,
                          output_dir=".", split_tag="predict",
                          lm_name="", db_name="",
                          barcodes=None, actual_labels=None, label_col="",
                          pub_dpi=300, logger=None):
        """
        CDR3 in-silico mutagenesis — RF version.
        Mutates every CDR3 position to all 20 AAs, predicts score,
        plots 20 x L heatmap per antibody (Nature Biotech style).
        Saves TIFFs + PPT.
        """
        if self.model is None or self.fb_ is None:
            _log(logger, "[MUTAGENESIS] model not trained"); return
        os.makedirs(output_dir, exist_ok=True)
        if barcodes is None:
            barcodes = list(X_df.index.astype(str))
        task_str = "CL" if self.task == "classification" else "REG"
        lm_str   = lm_name if lm_name else "RF"
        db_str   = db_name if db_name else ""
        thresh   = getattr(self, "recommended_threshold", 0.5)

        wt_feats = self.fb_.transform(X_df, embeddings)
        wt_scores = (self.model.predict_proba(wt_feats)[:, 1]
                     if self.task == "classification"
                     else self.model.predict(wt_feats))

        saved_paths = []

        for ab_idx in range(len(X_df)):
            bc       = barcodes[ab_idx] if ab_idx < len(barcodes) else str(ab_idx)
            act      = (actual_labels[ab_idx]
                        if actual_labels is not None and ab_idx < len(actual_labels)
                        else None)
            wt_score = float(wt_scores[ab_idx])
            row_wt   = X_df.iloc[ab_idx]
            cdr3_wt  = str(row_wt.get("CDR3", "") or "").upper().replace("-", "")
            cdr3_len = len(cdr3_wt)
            if cdr3_len == 0:
                _log(logger, f"[MUTAGENESIS] {bc}: empty CDR3 — skip"); continue
            _log(logger, f"[MUTAGENESIS] {ab_idx+1}/{len(X_df)}  {bc}  "
                         f"CDR3={cdr3_wt}  wt={wt_score:.4f}")

            score_mat = np.full((len(AMINO_ACIDS), cdr3_len), np.nan, dtype=np.float32)
            for pos in range(cdr3_len):
                for aa_i, mut_aa in enumerate(AMINO_ACIDS):
                    cdr3_mut        = cdr3_wt[:pos] + mut_aa + cdr3_wt[pos+1:]
                    row_m           = row_wt.copy()
                    row_m["CDR3"]   = cdr3_mut
                    emb_i           = (embeddings[[ab_idx]]
                                       if embeddings is not None else None)
                    feat_m          = self.fb_.transform(pd.DataFrame([row_m]), emb_i)
                    score_mat[aa_i, pos] = float(
                        self.model.predict_proba(feat_m)[0, 1]
                        if self.task == "classification"
                        else self.model.predict(feat_m)[0])

            # Heatmap
            fig, ax = plt.subplots(figsize=(max(6.0, cdr3_len * 0.42), 5.5))
            vmin, vmax = (0.0, 1.0) if self.task == "classification" else (
                float(np.nanpercentile(score_mat, 2)),
                float(np.nanpercentile(score_mat, 98)))
            cbar_lbl = "P(PASS)" if self.task == "classification" else "Predicted score"
            import matplotlib.colors as _mc
            _cmap = 'RdBu'
            _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0) \
                    if self.task == 'classification' else \
                    _mc.TwoSlopeNorm(vmin=float(score_mat.min()),
                                     vcenter=float((score_mat.min()+score_mat.max())/2),
                                     vmax=float(score_mat.max()))
            im   = ax.imshow(score_mat, cmap=_cmap, norm=_norm,
                             aspect='auto', interpolation='nearest')
            cbar = plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02)
            cbar.set_label(cbar_lbl, fontsize=7.5); cbar.ax.tick_params(labelsize=7)
            ax.set_yticks(range(len(AMINO_ACIDS)))
            ax.set_yticklabels(list(AMINO_ACIDS), fontsize=7.5)
            ax.set_xticks(range(cdr3_len))
            ax.set_xticklabels([f"{cdr3_wt[i]}\n{i+1}" for i in range(cdr3_len)],
                               fontsize=7.5)
            ax.set_xlabel("CDR3 position  (WT residue / number)", fontsize=8)
            ax.set_ylabel("Mutant amino acid", fontsize=8)
            # WT diagonal box
            for pos in range(cdr3_len):
                r = AMINO_ACIDS.find(cdr3_wt[pos])
                if r >= 0:
                    ax.add_patch(plt.Rectangle((pos-0.5, r-0.5), 1, 1,
                        fill=False, edgecolor="black", lw=1.5, zorder=5))
            # Title
            wt_lbl     = ("PASS" if wt_score >= thresh else "FAIL") if self.task == "classification" else ""
            score_line = (f"WT: P(PASS)={wt_score:.4f} → {wt_lbl}"
                          if self.task == "classification"
                          else f"WT predicted={wt_score:.4f}")
            col_name   = label_col if label_col else "label"
            act_part   = (f"  |  Actual ({col_name})="
                          + ("PASS" if int(act)==1 else "FAIL")
                          if act is not None else "")
            ax.set_title(
                f"IPI MLAbDev Platform: CDR3 in-silico mutagenesis — {bc}{act_part}\n"
                f"{score_line}\n"
                f"RF | {lm_str} | {task_str} | {db_str}",
                fontsize=8, loc="center", pad=8, linespacing=1.4)
            plt.tight_layout()
            bc_safe  = str(bc).replace("/","_").replace(" ","_")
            img_path = os.path.join(output_dir,
                f"{ab_idx+1:04d}_{bc_safe}_mutagenesis_{split_tag}.tiff")
            plt.savefig(img_path, dpi=pub_dpi, format="tiff", bbox_inches="tight")
            plt.close()
            saved_paths.append((img_path, bc, score_line))

        _log(logger, f"[MUTAGENESIS] {len(saved_paths)} heatmaps → {output_dir}/")

        if saved_paths:
            try:
                from pptx import Presentation as _Prs
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
                prs = _Prs()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]
                for img_path, bc, score_line in saved_paths:
                    slide = prs.slides.add_slide(blank)
                    img_w = Inches(10.5)
                    slide.shapes.add_picture(img_path,
                        (prs.slide_width-img_w)/2, Inches(0.3),
                        width=img_w, height=Inches(6.8))
                    txb = slide.shapes.add_textbox(
                        Inches(0.1), Inches(7.1), Inches(13.0), Inches(0.35))
                    tf  = txb.text_frame
                    tf.text = f"{bc}  |  {score_line}  |  RF | {lm_str} | {task_str} | {db_str}"
                    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]; run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x88, 0x87, 0x80)
                ppt = os.path.join(output_dir,
                    f"mutagenesis_{split_tag}_{lm_str}_{db_str}.pptx")
                prs.save(ppt)
                _log(logger, f"[MUTAGENESIS] PPT ({len(saved_paths)} slides) → {ppt}")
            except ImportError:
                _log(logger, "[MUTAGENESIS] pip install python-pptx for PPT")
            except Exception as _me:
                import traceback as _tb
                _log(logger, f"[MUTAGENESIS] PPT failed — {_me}")

        # ── SHAP ──────────────────────────────────────────────────────────────────

    def cdr3_mutagenesis_heatmap(
        self,
        X_df: pd.DataFrame,
        embeddings: np.ndarray = None,
        output_dir: str = None,
        pub_dpi: int = 300,
        fig_size: tuple = (10, 4),
        logger = None,
    ) -> dict:
        """
        In-silico CDR3 single-point mutagenesis for every antibody in X_df.
        For each position × amino acid substitution, predicts score/probability
        using the trained RF model and generates a heatmap per antibody.

        Saves:
          {output_dir}/{barcode}_cdr3_mutagenesis.tiff   (Nature 300 DPI)
          {output_dir}/cdr3_mutagenesis_all.pptx         (one slide per antibody)

        Returns dict: {barcode: heatmap_matrix (20 × cdr3_len)}
        """
        if self.model is None or self.fb_ is None:
            _log(logger, "[MUTA] model not trained/loaded"); return {}

        out_dir = output_dir or MODEL_DIR
        os.makedirs(out_dir, exist_ok=True)
        results = {}
        saved_paths = []

        # Get prediction scores for WT sequences first
        for i in range(len(X_df)):
            row = X_df.iloc[i]
            bc  = str(X_df.index[i]) if hasattr(X_df, 'index') else str(i)
            cdr3_seq = str(row.get('CDR3', '') or '').upper().replace('-', '')
            if not cdr3_seq:
                _log(logger, f"  [MUTA] {bc}: empty CDR3 — skipped")
                continue

            n_pos  = len(cdr3_seq)
            n_aa   = len(AMINO_ACIDS)
            matrix = np.zeros((n_aa, n_pos), dtype=np.float32)

            # Build mutation DataFrame for all substitutions at once (vectorised)
            mut_rows = []
            mut_info = []   # (aa_idx, pos_idx)
            for pos in range(n_pos):
                for aa_idx, aa in enumerate(AMINO_ACIDS):
                    mutant_cdr3 = cdr3_seq[:pos] + aa + cdr3_seq[pos+1:]
                    mut_row = dict(row)
                    mut_row['CDR3'] = mutant_cdr3
                    # Also mutate HSEQ if available (replace CDR3 region)
                    if 'HSEQ' in mut_row and mut_row['HSEQ']:
                        hseq = str(mut_row['HSEQ'])
                        cdr3_start = hseq.find(cdr3_seq)
                        if cdr3_start >= 0:
                            mut_row['HSEQ'] = (hseq[:cdr3_start] + mutant_cdr3 +
                                               hseq[cdr3_start + n_pos:])
                    mut_rows.append(mut_row)
                    mut_info.append((aa_idx, pos))

            # Predict all mutants in one batch
            mut_df = pd.DataFrame(mut_rows)
            mut_df.index = [f"{bc}_mut_{j}" for j in range(len(mut_df))]

            # Build embeddings for mutants if needed
            mut_emb = None
            if embeddings is not None and self.feat_cfg.get('embedding', False):
                _log(logger,
                     f"  [MUTA] {bc}: embedding mode — mutant embeddings not available.\n"
                     "         Mutagenesis requires kmer/biophysical/onehot features.\n"
                     "         Set features.embedding=false in YAML.")
                continue

            try:
                X_mut = self.fb_.transform(mut_df, mut_emb)
                if self.task == 'classification':
                    scores = self.model.predict_proba(X_mut)[:, 1]
                else:
                    scores = self.model.predict(X_mut)
            except Exception as _e:
                _log(logger, f"  [MUTA] {bc}: prediction failed — {_e}"); continue

            for j, (aa_idx, pos_idx) in enumerate(mut_info):
                matrix[aa_idx, pos_idx] = scores[j]

            results[bc] = matrix

            # ── WT score (diagonal reference line) ────────────────────────────
            wt_row = pd.DataFrame([dict(row)])
            wt_row.index = [bc]
            try:
                X_wt = self.fb_.transform(wt_row, None)
                wt_score = (self.model.predict_proba(X_wt)[:, 1][0]
                            if self.task == 'classification'
                            else self.model.predict(X_wt)[0])
            except Exception:
                wt_score = None

            # ── Plot heatmap ─────────────────────────────────────────────────
            fig, ax = plt.subplots(figsize=(max(fig_size[0], n_pos * 0.45), fig_size[1]))

            score_label = "P(PASS)" if self.task == 'classification' else "Predicted score"
            cmap = 'RdBu'
            import matplotlib.colors as _mc
            if self.task == 'classification':
                _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0)
            else:
                _vmin = float(matrix.min()); _vmax = float(matrix.max())
                _norm = _mc.TwoSlopeNorm(vmin=_vmin, vcenter=(_vmin+_vmax)/2, vmax=_vmax)

            im = ax.imshow(matrix, aspect='auto', cmap=cmap, norm=_norm)

            # ── Cell annotations ─────────────────────────────────────────────
            _cell_fsz = max(4.0, min(7.0, 120.0 / max(n_pos, 1)))
            for _ai in range(n_aa):
                for _pi in range(n_pos):
                    _v = matrix[_ai, _pi]
                    if np.isnan(_v): continue
                    _tc = 'white' if (_v < 0.35 or _v > 0.65) else '#333333'
                    ax.text(_pi, _ai, f"{_v:.2f}", ha='center', va='center',
                            fontsize=_cell_fsz, color=_tc,
                            fontweight='bold' if abs(_v - 0.5) > 0.3 else 'normal')

            plt.colorbar(im, ax=ax, label=score_label, fraction=0.03, pad=0.02)

            ax.set_xticks(range(n_pos))
            ax.set_xticklabels([f"{cdr3_seq[p]}\n{p+1}" for p in range(n_pos)],
                               fontsize=8)
            ax.set_yticks(range(n_aa))
            ax.set_yticklabels(list(AMINO_ACIDS), fontsize=7)
            ax.set_xlabel('CDR3 position  (WT residue shown above position number)',
                          fontsize=8)
            ax.set_ylabel('Substituted AA', fontsize=8)

            wt_str = f"  WT {score_label}={wt_score:.4f}" if wt_score is not None else ""
            ax.set_title(
                f"IPI MLAbDev · CDR3 Mutagenesis Heatmap\n"
                f"ID: {bc}{wt_str}\n"
                f"RF | {self.task.upper()} | CDR3={cdr3_seq}",
                fontsize=8, loc='center', pad=6
            )

            # Mark WT residues on diagonal
            for pos in range(n_pos):
                wt_aa_idx = AMINO_ACIDS.find(cdr3_seq[pos])
                if wt_aa_idx >= 0:
                    ax.add_patch(plt.Rectangle(
                        (pos - 0.5, wt_aa_idx - 0.5), 1, 1,
                        fill=False, edgecolor='black', lw=1.5, zorder=5
                    ))

            plt.tight_layout()
            bc_safe  = bc.replace('/', '_').replace(' ', '_')
            img_path = os.path.join(out_dir, f"{bc_safe}_cdr3_mutagenesis.tiff")
            plt.savefig(img_path, dpi=pub_dpi, format='tiff', bbox_inches='tight')
            plt.close()
            saved_paths.append((bc, img_path, wt_score))
            _log(logger, f"  [MUTA] {bc} → {os.path.basename(img_path)}")

        # ── Assemble PPT ──────────────────────────────────────────────────────
        if saved_paths:
            try:
                from pptx import Presentation as _Prs
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor

                prs = _Prs()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                for bc, img_path, wt_score in saved_paths:
                    slide = prs.slides.add_slide(blank)
                    img_w = Inches(fig_size[0])
                    img_h = Inches(fig_size[1])
                    left  = (prs.slide_width - img_w) / 2
                    slide.shapes.add_picture(img_path, left, Inches(0.6),
                                             width=img_w, height=img_h)
                    txb = slide.shapes.add_textbox(
                        Inches(0.15), Inches(7.1), Inches(13.0), Inches(0.35))
                    tf  = txb.text_frame
                    wt_str = f"  WT score={wt_score:.4f}" if wt_score is not None else ""
                    tf.text = f"CDR3 Mutagenesis  |  {bc}{wt_str}"
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    p.runs[0].font.size = Pt(7)
                    p.runs[0].font.color.rgb = RGBColor(0x88, 0x87, 0x80)

                ppt_path = os.path.join(out_dir, "cdr3_mutagenesis_all.pptx")
                prs.save(ppt_path)
                _log(logger, f"[MUTA] PPT ({len(saved_paths)} slides) → {ppt_path}")
            except ImportError:
                _log(logger, "[MUTA] python-pptx not installed — PPT skipped")
            except Exception as _pe:
                _log(logger, f"[MUTA] PPT failed — {_pe}")

        _log(logger, f"[MUTA] complete  ({len(results)} antibodies)")
        return results

    def shap_analysis(self, X_df, embeddings=None,
                      output_prefix="rf_shap",
                      split_tag="train",
                      top_n=50,
                      barcodes=None,
                      actual_labels=None,
                      actual_col_name="label",
                      lm_name="",
                      db_name="",
                      logger=None):
        """
        SHAP TreeExplainer — works for ALL feature modes.
        Generates beeswarm/bar/heatmap for all samples, plus a multi-page
        PDF waterfall (one page per antibody) with full title info.

        Parameters
        ----------
        X_df          : DataFrame with HSEQ/CDR3 columns (index = BARCODE)
        embeddings    : (n, emb_dim) PLM array or None
        output_prefix : filename stem
        split_tag     : 'train' | 'val' | 'predict'
        top_n         : top features shown per plot
        barcodes      : list of antibody IDs — if None, inferred from X_df.index
        actual_labels : list/array of ground-truth labels — shown in title if provided
        actual_col_name: target column name (e.g. 'psr_filter') — shown in title
        lm_name       : LM name for title (e.g. 'ablang', 'biophysical')
        db_name       : database stem for title (e.g. 'ipi_psr_trainset_train')

        YAML shap config keys used here:
          top_features          : features per waterfall panel  (default 30)
          waterfall_top_features: override top_n for waterfall only (default = top_features)
          max_waterfall_samples : max antibodies in PDF          (default 50)
          waterfall_per_page    : 1 | 2 | 4 panels per page     (default 1)
          plot_dpi              : DPI for PNG outputs            (default 150)
          pub_dpi               : DPI for publication PDF        (default 300)
        """
        if not _SHAP_AVAILABLE:
            _log(logger, "[SHAP] shap not installed — pip install shap"); return
        if self.model is None or self.fb_ is None:
            _log(logger, "[SHAP] model not trained"); return

        # Infer barcodes from X_df index or provided list
        if barcodes is None:
            barcodes = list(X_df.index.astype(str)) if hasattr(X_df, 'index') else [str(i) for i in range(len(X_df))]

        sh_cfg     = self.config.get('shap', {})
        max_s      = sh_cfg.get('max_samples', 500)
        plots      = sh_cfg.get('plot_types', ['bar', 'beeswarm'])
        top_n      = sh_cfg.get('top_features', top_n)        # YAML overrides arg
        wf_top_n   = sh_cfg.get('waterfall_top_features', top_n)  # waterfall-specific
        max_wf     = sh_cfg.get('max_waterfall_samples', 50)  # max antibodies in PDF
        wf_per_pg  = sh_cfg.get('waterfall_per_page', 1)      # 1 | 2 | 4 per page
        plot_dpi   = sh_cfg.get('plot_dpi', 150)               # PNG dpi
        pub_dpi    = sh_cfg.get('pub_dpi',  300)               # PDF dpi (publication)
        # beeswarm is always generated regardless of plots list
        if 'beeswarm' not in plots:
            plots = list(plots) + ['beeswarm']

        # Build feature matrix and names
        X_feat   = self.fb_.transform(X_df, embeddings)
        n_feat   = X_feat.shape[1]
        ne_names = self.fb_.non_embedding_feature_names   # kmer + biophysical
        ne_idx   = self.fb_.non_embedding_indices

        # ── Embedding-only mode: SHAP is not meaningful ───────────────────────
        # PLM dimensions (emb_0 … emb_N) are entangled latent features —
        # no single dimension maps to a biological property.
        # Use Integrated Gradients from transformer_lm for sequence-level
        # attribution on PLM embeddings.
        if not ne_idx:
            _log(logger,
                 "[SHAP] Skipped — embedding-only mode.\n"
                 "  PLM dimensions (emb_0..emb_N) are latent and not interpretable.\n"
                 "  For sequence-level attribution use transformer_lm IG instead:\n"
                 "    model.global_ig_analysis(dataset)  in transformer_lm.py")
            return

        # Slice to interpretable features only (kmer + biophysical).
        # Embedding dims are excluded — they have no biological name.
        # SHAP runs only on ne_idx columns so no empty labels appear.
        ne_idx_arr = np.array(ne_idx)
        X_ne   = X_feat[:, ne_idx_arr]          # (n, n_interpretable)
        all_names = ne_names                     # names align 1:1 with X_ne columns

        # Subsample for speed
        # Keep full (unsubsampled) matrix for predict_proba — correct scores
        if len(X_ne) > max_s:
            idx_s   = np.random.choice(len(X_ne), max_s, replace=False)
            X_shap  = X_ne[idx_s]
            # For waterfall predictions: use same subsample but from full X_feat
            _X_full_wf = X_feat[idx_s]
        else:
            X_shap     = X_ne
            _X_full_wf = X_feat

        _log(logger, f"\n[SHAP] {split_tag.upper()} — TreeExplainer  "
                     f"n={len(X_shap)}  interpretable features={len(ne_idx)}  "
                     f"(embedding dims excluded)  top={top_n}")

        explainer = _shap_lib.TreeExplainer(self.model)
        try:
            shap_values = explainer.shap_values(X_shap, check_additivity=True)
        except Exception:
            shap_values = explainer.shap_values(X_shap, check_additivity=False)

        # Robustly extract (n_samples, n_features) array for class-1 or regression.
        # SHAP returns different shapes across versions:
        #   regression           : ndarray (n, f)
        #   binary classification:
        #     shap < 0.41        : list[ (n,f), (n,f) ]           → take [1]
        #     shap >= 0.42 binary: ndarray (n, f)                  → direct
        #     shap >= 0.42 multi : ndarray (n, f, n_classes)       → take [:,:,1]
        #     some versions      : ndarray (n_classes, n, f)       → take [1]
        if isinstance(shap_values, list):
            # list of arrays — take class 1
            sv = np.array(shap_values[1], dtype=np.float64)
        elif isinstance(shap_values, np.ndarray):
            if shap_values.ndim == 3:
                # could be (n, f, classes) or (classes, n, f)
                if shap_values.shape[0] == len(X_shap):
                    sv = shap_values[:, :, 1]        # (n, f, classes) → class 1
                else:
                    sv = shap_values[1]              # (classes, n, f) → class 1
            else:
                sv = shap_values                     # (n, f) — regression or binary
        else:
            sv = np.array(shap_values, dtype=np.float64)

        # Ensure sv is always a plain float64 numpy array (n_samples, n_features)
        sv = np.asarray(sv, dtype=np.float64)
        if sv.ndim != 2:
            raise ValueError(f"[SHAP] Unexpected sv shape after extraction: {sv.shape}")
        _log(logger, f"  shap_values shape={sv.shape}  dtype={sv.dtype}")

        mean_abs = np.mean(np.abs(sv), axis=0)
        # Convert to plain Python ints to avoid numpy scalar indexing errors
        top_idx  = [int(i) for i in np.argsort(mean_abs)[::-1][:top_n]]
        top_vals = [float(mean_abs[i]) for i in top_idx]
        top_lbls = [all_names[i] for i in top_idx]

        os.makedirs(MODEL_DIR, exist_ok=True)
        prefix = f"{output_prefix}_{split_tag}"

        # ── Print top-N table ─────────────────────────────────────────────────
        _log(logger, f"  {'Rank':>4}  {'Feature':35s}  {'Mean |SHAP|':>11}")
        _log(logger, f"  {'─'*4}  {'─'*35}  {'─'*11}")
        for rank, (lbl, val) in enumerate(zip(top_lbls, top_vals), 1):
            _log(logger, f"  {rank:4d}  {lbl:35s}  {val:11.6f}")

        # ── CSV export ────────────────────────────────────────────────────────
        import csv
        csv_path = os.path.join(MODEL_DIR, f"{prefix}_shap_top{top_n}.csv")
        with open(csv_path, 'w', newline='') as f:
            w = csv.writer(f)
            w.writerow(['rank', 'feature', 'mean_abs_shap'])
            for rank, (lbl, val) in enumerate(zip(top_lbls, top_vals), 1):
                w.writerow([rank, lbl, f"{val:.8f}"])
        _log(logger, f"[SHAP] csv  → {csv_path}")

        # ── Bar chart ─────────────────────────────────────────────────────────
        if 'bar' in plots:
            colours = []
            for l in top_lbls:
                if l.startswith('emb_'):              colours.append('#378ADD')  # PLM embedding
                elif l.startswith('1mer_'):            colours.append('#534AB7')  # 1-mer
                elif l.startswith('2mer_'):            colours.append('#7F77DD')  # 2-mer
                elif l.startswith('3mer_'):            colours.append('#AFA9EC')  # 3-mer
                elif l.startswith('oh_vh'):            colours.append('#E07B39')  # one-hot VH
                elif l.startswith('oh_vl'):            colours.append('#F5B942')  # one-hot VL
                elif l.startswith('oh_hcdr3') or l.startswith('oh_cdr3'):
                                                       colours.append('#C94A8C')  # one-hot HCDR3
                elif l.startswith('oh_'):              colours.append('#E07B39')  # one-hot other
                elif l.startswith('cdr3_') or l.startswith('vh_'):
                                                       colours.append('#1D9E75')  # biophysical
                else:                                  colours.append('#1D9E75')  # biophysical

            fig, ax = plt.subplots(figsize=(8.27, min(11.69, max(5, top_n * 0.20))))
            ax.barh(range(len(top_idx)), top_vals, color=colours)
            ax.set_yticks(range(len(top_idx)))
            ax.set_yticklabels(top_lbls, fontsize=7)
            ax.invert_yaxis()
            ax.set_xlabel('Mean |SHAP value|')
            ax.set_title(f'SHAP top {top_n} — {output_prefix} [{split_tag}]')
            patches = [mpatches.Patch(color='#378ADD', label='embedding'),
                       mpatches.Patch(color='#534AB7', label='1-mer'),
                       mpatches.Patch(color='#7F77DD', label='2-mer'),
                       mpatches.Patch(color='#AFA9EC', label='3-mer'),
                       mpatches.Patch(color='#1D9E75', label='biophysical'),
                       mpatches.Patch(color='#E07B39', label='one-hot VH'),
                       mpatches.Patch(color='#F5B942', label='one-hot VL'),
                       mpatches.Patch(color='#C94A8C', label='one-hot HCDR3')]
            # Only show legend entries that are actually present
            _present = set(colours)
            _color_map = {'#378ADD': 'embedding', '#534AB7': '1-mer',
                          '#7F77DD': '2-mer', '#AFA9EC': '3-mer',
                          '#1D9E75': 'biophysical', '#E07B39': 'one-hot VH',
                          '#F5B942': 'one-hot VL', '#C94A8C': 'one-hot HCDR3'}
            patches = [mpatches.Patch(color=c, label=l)
                       for c, l in _color_map.items() if c in _present]
            ax.legend(handles=patches, loc='lower right', fontsize=8)
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_bar.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] bar  → {out}")

        # ── Beeswarm ──────────────────────────────────────────────────────────
        if 'beeswarm' in plots:
            _shap_lib.summary_plot(sv[:, top_idx],
                                   X_shap[:, top_idx] if hasattr(X_shap, '__getitem__') else X_shap,
                                   feature_names=top_lbls,
                                   show=False, max_display=top_n)
            plt.title(f'SHAP beeswarm — {output_prefix} [{split_tag}]')
            plt.gcf().set_size_inches(8.27, min(11.69, max(5, top_n * 0.22)))
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_beeswarm.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] beeswarm → {out}")

        # ── Heatmap ───────────────────────────────────────────────────────────
        if 'heatmap' in plots:
            n_heat = min(top_n, 20)
            sub    = sv[:min(100, len(sv)), :][:, top_idx[:n_heat]]   # top_idx is plain int list
            lbls   = top_lbls[:n_heat]
            fig, ax = plt.subplots(figsize=(11, 7))
            vmax    = np.abs(sub).max()
            im      = ax.imshow(sub.T, aspect='auto', cmap='RdBu_r',
                                vmin=-vmax, vmax=vmax)
            plt.colorbar(im, ax=ax, label='SHAP value')
            ax.set_yticks(range(len(lbls)))
            ax.set_yticklabels(lbls, fontsize=7)
            ax.set_xlabel('Sample index')
            ax.set_title(f'SHAP heatmap top {n_heat} — {output_prefix} [{split_tag}]')
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_heatmap.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] heatmap → {out}")

        # ── Waterfall — individual TIFF/JPEG + PPT slide per antibody ────────
        # YAML controls:
        #   max_waterfall_samples  : max antibodies          (default 50)
        #   waterfall_top_features : features per plot       (default 20)
        #   waterfall_fig_size     : [w, h] in inches        (default [8, 6])
        #   waterfall_format       : 'tiff' | 'jpeg' | 'png' (default 'tiff')
        #   pub_dpi                : image DPI               (default 300)
        #   waterfall_ppt          : also build PPT          (default true)
        _log(logger, f"\n[SHAP] Generating waterfall plots ...")
        _log(logger, f"  shap enabled={sh_cfg.get('enabled', False)}  "
                     f"max_wf={max_wf}  wf_top_n={wf_top_n}  "
                     f"fmt={sh_cfg.get('waterfall_format','tiff')}  "
                     f"ppt={sh_cfg.get('waterfall_ppt', True)}")
        try:
            import matplotlib.patches as _mp

            # ── Read YAML settings ────────────────────────────────────────────
            n_wf      = min(len(sv), max_wf)
            wf_size   = sh_cfg.get('waterfall_fig_size', [8, 6])
            wf_fmt    = sh_cfg.get('waterfall_format', 'tiff').lower().strip('.')
            make_ppt  = sh_cfg.get('waterfall_ppt', True)
            if wf_fmt not in ('tiff', 'jpeg', 'jpg', 'png'):
                wf_fmt = 'tiff'

            if n_wf < len(sv):
                _log(logger, f"  Waterfall limited to first {n_wf}/{len(sv)} "
                             f"(max_waterfall_samples={max_wf})")

            # ── Output folder: {MODEL_DIR}/{prefix}_waterfalls/ ───────────────
            wf_dir = os.path.join(MODEL_DIR, f"{prefix}_shap_waterfalls")
            os.makedirs(wf_dir, exist_ok=True)

            # ── Predictions — use full feature matrix for correct probabilities ────
            # X_shap is sliced to ne_idx only. For predict_proba, use _X_full_wf
            # (full feature matrix) so P(PASS) matches model's actual output.
            if self.task == 'classification':
                _thresh         = getattr(self, 'recommended_threshold', 0.5)
                all_probs       = self.model.predict_proba(_X_full_wf[:n_wf])[:, 1]
                all_labels_pred = (all_probs >= _thresh).astype(int)
            else:
                all_probs       = self.model.predict(_X_full_wf[:n_wf])
                all_labels_pred = None

            expected_val = float(
                explainer.expected_value[1]
                if isinstance(explainer.expected_value, (list, np.ndarray))
                else explainer.expected_value
            )
            task_str = "CL"  if self.task == 'classification' else "REG"
            lm_str   = lm_name if lm_name else "unknown_lm"
            db_str   = db_name if db_name else output_prefix

            saved_paths = []

            for s_idx in range(n_wf):
                bc  = barcodes[s_idx]                       if barcodes and s_idx < len(barcodes) else str(s_idx)
                act = actual_labels[s_idx]                       if actual_labels is not None                       and s_idx < len(actual_labels) else None

                # Features
                n_f       = min(wf_top_n, len(top_idx))
                shap_vals = [float(sv[s_idx][i])       for i in top_idx[:n_f]]
                feat_vals = [float(X_shap[s_idx, i])   for i in top_idx[:n_f]]
                lbls_s    = top_lbls[:n_f]

                triples = sorted(zip(shap_vals, lbls_s, feat_vals),
                                 key=lambda x: abs(x[0]), reverse=True)
                wf_s    = [x[0] for x in triples]
                wf_l    = [f"{x[1]} = {x[2]:.3g}" for x in triples]
                cols    = ['#D62728' if v < 0 else '#1F77B4' for v in wf_s]

                cumsum  = expected_val
                bottoms = []
                for v in wf_s:
                    bottoms.append(cumsum)
                    cumsum += v

                # ── Fixed figure size — no scaling ───────────────────────────
                fig, ax = plt.subplots(figsize=(wf_size[0], wf_size[1]))

                # Bar height and font size fixed to figure size / n_features
                bar_h = 0.70   # fixed — same solid bar height regardless of n features
                fsz   = 8.5    # fixed font size — consistent across all plots

                # Compute x range once for inside/outside text decision
                _all_ends = [b + v for b, v in zip(bottoms, wf_s)]
                _xmin  = min(bottoms + [expected_val] + _all_ends)
                _xmax  = max(bottoms + [expected_val] + _all_ends)
                _xrng  = max(abs(_xmax - _xmin), 1e-6)

                for i, (v, lbl, bot, col) in enumerate(
                        zip(wf_s, wf_l, bottoms, cols)):
                    ax.barh(i, v, left=bot, color=col,
                            height=bar_h, edgecolor='white', linewidth=0.4)
                    bw = abs(v)
                    # Put text inside bar only if bar is wide enough to read
                    inside = bw / _xrng > 0.08
                    if inside:
                        xt = bot + v / 2
                        ha = 'center'
                        tc = 'white'
                    else:
                        pad = _xrng * 0.012
                        xt  = bot + v + pad * (1 if v >= 0 else -1)
                        ha  = 'left' if v >= 0 else 'right'
                        tc  = col
                    # Only show text if value is meaningful (> 0.5% of range)
                    if bw / _xrng > 0.005:
                        ax.text(xt, i, f"{v:+.4f}", va='center', ha=ha,
                                fontsize=max(6.5, fsz - 0.5),
                                color=tc,
                                fontweight='bold' if bw / _xrng > 0.05 else 'normal')

                ax.set_yticks(range(n_f))
                ax.set_yticklabels(wf_l, fontsize=max(7.0, fsz))
                ax.invert_yaxis()
                # Pad left so outside-bar text never overlaps y-axis labels
                _xpad = _xrng * 0.18
                ax.set_xlim(_xmin - _xpad, _xmax + _xrng * 0.06)
                ax.axvline(expected_val, color='#888',    lw=0.8, ls='--')
                ax.axvline(cumsum,       color='#222299', lw=2.0, ls='-')
                ax.set_xlabel('SHAP value  (cumulative impact on model output)', fontsize=max(8.0, fsz))
                ax.tick_params(axis='x', labelsize=max(7.0, fsz - 0.5))

                # Title
                _col_name = actual_col_name if actual_col_name else "label"
                act_part  = (f"  |  Actual ({actual_col_name}) = "
                             f"{'PASS' if int(act)==1 else 'FAIL'}"
                             ) if act is not None else ""
                if self.task == 'classification':
                    prob    = float(all_probs[s_idx])
                    ml_lbl  = "PASS" if all_labels_pred[s_idx] == 1 else "FAIL"
                    sc_part = f"P(PASS)={prob:.4f}  →  {ml_lbl}"
                else:
                    sc_part = f"Predicted={float(all_probs[s_idx]):.4f}"

                # Line 1: "SHAP waterfall — {ID}"
                # Line 2: P(PASS)=X.XXXX → PASS/FAIL  |  Actual (col)=PASS/FAIL  |  RF|lm|CL|db
                _act_str = ""
                if act is not None:
                    _act_val  = "PASS" if int(act) == 1 else "FAIL"
                    _col_name = actual_col_name if actual_col_name else "label"
                    _act_str  = f"  |  Actual ({_col_name}) = {_act_val}"

                _id_part   = f"IPI MLAbDev Platform: Prediction interpretability — {bc}"
                _info_part = f"{sc_part}{_act_str}  |  RandomForest | {lm_str} | {task_str} | {db_str}"
                ax.set_title(
                    f"{_id_part}\n{_info_part}",
                    fontsize=8.5, loc='center', pad=10
                )
                ax.legend(handles=[
                    _mp.Patch(color='#1F77B4', label='toward PASS  (class 1 +)'),
                    _mp.Patch(color='#D62728', label='toward FAIL  (class 0 −)'),
                    plt.Line2D([0],[0], color='#999',    ls='--', lw=1.0,
                               label=f'baseline={expected_val:.3f}'),
                    plt.Line2D([0],[0], color='#1A237E', ls='-',  lw=2.0,
                               label=f'final={cumsum if abs(cumsum) > 1e-4 else 0.0:.4f}'),
                ], fontsize=max(5.0, fsz - 1.0),
                   loc='lower right', framealpha=0.85,
                   edgecolor='#ccc')

                plt.tight_layout()

                # Safe filename from barcode
                bc_safe = str(bc).replace('/', '_').replace('\\', '_').replace(' ', '_')
                img_path = os.path.join(
                    wf_dir,
                    f"{s_idx+1:04d}_{bc_safe}_waterfall.{wf_fmt}"
                )
                save_kw = dict(dpi=pub_dpi, bbox_inches='tight')
                if wf_fmt in ('jpeg', 'jpg'):
                    save_kw['format']  = 'jpeg'
                    save_kw['pil_kwargs'] = {'quality': 95}
                elif wf_fmt == 'tiff':
                    save_kw['format'] = 'tiff'
                plt.savefig(img_path, **save_kw)
                plt.close()
                saved_paths.append(img_path)

            _log(logger, f"[SHAP] {n_wf} waterfall images ({wf_fmt}, "
                         f"{pub_dpi} DPI) → {wf_dir}/")

            # ── PPT — one slide per antibody ──────────────────────────────────
            if make_ppt:
                try:
                    from pptx import Presentation as _Prs
                    from pptx.util import Inches, Pt
                    from pptx.enum.text import PP_ALIGN
                    import io as _io

                    prs   = _Prs()
                    # Widescreen 16:9 slide
                    prs.slide_width  = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    blank_layout = prs.slide_layouts[6]  # blank

                    for img_path, s_idx in zip(saved_paths,
                                               range(len(saved_paths))):
                        bc  = barcodes[s_idx]                               if barcodes and s_idx < len(barcodes) else str(s_idx)
                        slide = prs.slides.add_slide(blank_layout)

                        # Image — centered, leaving margin for title
                        img_w = Inches(wf_size[0])
                        img_h = Inches(wf_size[1])
                        left  = (prs.slide_width  - img_w) / 2
                        top   = Inches(0.6)
                        slide.shapes.add_picture(
                            img_path, left, top,
                            width=img_w, height=img_h
                        )

                        # Slide number + barcode footer
                        txb = slide.shapes.add_textbox(
                            Inches(0.15), Inches(7.1),
                            Inches(13.0), Inches(0.35)
                        )
                        tf  = txb.text_frame
                        tf.text = (f"{s_idx+1}/{n_wf}  |  {bc}  |  "
                                   f"RF | {lm_str} | {task_str} | {db_str}  |  "
                                   f"{split_tag}")
                        from pptx.dml.color import RGBColor
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0]
                        run.font.size = Pt(7)
                        run.font.color.rgb = RGBColor(0x88, 0x87, 0x80)

                    ppt_path = os.path.join(MODEL_DIR,
                                            f"{prefix}_shap_waterfall.pptx")
                    prs.save(ppt_path)
                    _log(logger,
                         f"[SHAP] waterfall PPT ({n_wf} slides) → {ppt_path}")

                except ImportError:
                    _log(logger,
                         "[SHAP] python-pptx not installed — PPT skipped.\n"
                         "  pip install python-pptx")
                except Exception as _pe:
                    import traceback as _tb
                    _log(logger, f"[SHAP] PPT failed — {_pe}")
                    _log(logger, _tb.format_exc())

        except Exception as _we:
            import traceback as _tb
            _log(logger, f"[SHAP] waterfall failed — {_we}")
            _log(logger, _tb.format_exc())
        _log(logger, f"[SHAP] {split_tag} complete.")

    # ── compare_waterfall ────────────────────────────────────────────────────
    @staticmethod
    def compare_waterfall(
        models:        list,
        model_names:   list,
        X_df:          'pd.DataFrame',
        embeddings_list: list = None,
        barcode:       str  = None,
        actual_label:  int  = None,
        label_col:     str  = "",
        top_n:         int  = 20,
        fig_size:      tuple = (16, 7),
        pub_dpi:       int  = 300,
        output_path:   str  = None,
        logger         = None,
    ):
        """
        Side-by-side SHAP waterfall comparison of multiple models for ONE antibody.

        Parameters
        ----------
        models        : list of trained RandomForestModel instances
        model_names   : list of display names  e.g. ['biophysical', 'kmer']
        X_df          : single-row DataFrame (BARCODE index, HSEQ/CDR3 columns)
        embeddings_list: list of embedding arrays (one per model) or None per model
        barcode       : antibody ID for title (inferred from X_df.index if None)
        actual_label  : ground-truth label (0/1) or None
        label_col     : column name shown in title  e.g. 'psr_filter'
        top_n         : features per panel
        fig_size      : total figure size (w, h) — split equally across panels
        pub_dpi       : DPI for saved image
        output_path   : save path (.tiff/.png/.jpeg) — if None, just show
        """
        if not _SHAP_AVAILABLE:
            print("[compare_waterfall] shap not installed — pip install shap")
            return

        n_models = len(models)
        if barcode is None:
            barcode = str(X_df.index[0]) if hasattr(X_df, 'index') else "unknown"
        if embeddings_list is None:
            embeddings_list = [None] * n_models

        import matplotlib.patches as _mp
        fig, axes = plt.subplots(
            1, n_models,
            figsize=fig_size,
            sharey=False
        )
        if n_models == 1:
            axes = [axes]

        for ax, model, mname, emb in zip(
                axes, models, model_names, embeddings_list):

            if model.model is None or model.fb_ is None:
                ax.text(0.5, 0.5, f"{mname}\nNot trained",
                        ha='center', va='center', transform=ax.transAxes)
                continue

            # Build feature matrix
            X_feat   = model.fb_.transform(X_df, emb)
            ne_idx   = model.fb_.non_embedding_indices
            ne_names = model.fb_.non_embedding_feature_names
            if not ne_idx:
                ax.text(0.5, 0.5, f"{mname}\nEmbedding-only\n(no SHAP)",
                        ha='center', va='center', transform=ax.transAxes)
                continue

            X_ne = X_feat[:, np.array(ne_idx)]

            # SHAP
            explainer = _shap_lib.TreeExplainer(model.model)
            try:
                sv_raw = explainer.shap_values(X_ne, check_additivity=False)
            except TypeError:
                sv_raw = explainer.shap_values(X_ne)

            if isinstance(sv_raw, list):
                sv = np.asarray(sv_raw[1], dtype=np.float64)
            elif isinstance(sv_raw, np.ndarray) and sv_raw.ndim == 3:
                sv = sv_raw[:, :, 1] if sv_raw.shape[0] == len(X_ne)                      else sv_raw[1]
            else:
                sv = np.asarray(sv_raw, dtype=np.float64)

            expected_val = float(
                explainer.expected_value[1]
                if isinstance(explainer.expected_value, (list, np.ndarray))
                else explainer.expected_value
            )

            # Prediction score
            if model.task == 'classification':
                _thresh   = getattr(model, 'recommended_threshold', 0.5)
                prob      = float(model.model.predict_proba(X_ne)[:, 1][0])
                ml_label  = "PASS" if prob >= _thresh else "FAIL"
                score_str = f"P(PASS)={prob:.4f} → {ml_label}"
            else:
                prob      = float(model.model.predict(X_ne)[0])
                score_str = f"Pred={prob:.4f}"

            # Top features
            sv0      = sv[0]
            all_vals = [(float(sv0[i]), ne_names[i],
                         float(X_ne[0, i]))
                        for i in range(len(ne_names))]
            top      = sorted(all_vals, key=lambda x: abs(x[0]),
                               reverse=True)[:top_n]

            wf_s    = [x[0] for x in top]
            wf_lbls = [f"{x[1]} = {x[2]:.3g}" for x in top]
            cols    = ['#E05C3A' if v < 0 else '#2EAA6B' for v in wf_s]

            cumsum  = expected_val
            bottoms = []
            for v in wf_s:
                bottoms.append(cumsum)
                cumsum += v

            # x range for text placement
            all_ends = [b + v for b, v in zip(bottoms, wf_s)]
            _xrng = max(abs(max(bottoms + all_ends) -
                             min(bottoms + all_ends)), 1e-6)

            bar_h = 0.70
            fsz   = max(6.5, min(8.5, 130.0 / max(top_n, 1)))

            for i, (v, lbl, bot, col) in enumerate(
                    zip(wf_s, wf_lbls, bottoms, cols)):
                ax.barh(i, v, left=bot, color=col,
                        height=bar_h, edgecolor='white', linewidth=0.4)
                bw = abs(v)
                inside = bw / _xrng > 0.08
                if inside:
                    xt, ha, tc = bot + v/2, 'center', 'white'
                else:
                    xt = bot + v + _xrng * 0.012 * (1 if v >= 0 else -1)
                    ha = 'left' if v >= 0 else 'right'
                    tc = col
                ax.text(xt, i, f"{v:+.4f}", va='center', ha=ha,
                        fontsize=max(5.5, fsz - 1.5),
                        color=tc,
                        fontweight='bold' if bw / _xrng > 0.05 else 'normal')

            ax.set_yticks(range(top_n))
            ax.set_yticklabels(wf_lbls, fontsize=fsz)
            ax.invert_yaxis()
            ax.axvline(expected_val, color='#999',    lw=1.0, ls='--')
            ax.axvline(cumsum,       color='#1A237E', lw=2.0, ls='-')
            ax.set_xlabel('SHAP value  (cumulative impact)',
                          fontsize=max(7.0, fsz))
            ax.tick_params(axis='x', labelsize=max(6.5, fsz - 1))
            ax.set_title(f"{mname}\n{score_str}",
                         fontsize=max(7.5, fsz), pad=6)
            ax.legend(handles=[
                _mp.Patch(color='#2EAA6B', label='toward PASS (+)'),
                _mp.Patch(color='#E05C3A', label='toward FAIL  (−)'),
                plt.Line2D([0],[0], color='#999',    ls='--', lw=1.0,
                           label=f'baseline={expected_val:.3f}'),
                plt.Line2D([0],[0], color='#1A237E', ls='-',  lw=2.0,
                           label=f'final={cumsum:.3f}'),
            ], fontsize=max(5.0, fsz - 2), loc='lower right',
               framealpha=0.75)

        # Shared super-title
        act_part = ""
        if actual_label is not None:
            _col = label_col if label_col else "label"
            act_part = f"  |  Actual ({_col}) = "                        f"{'PASS' if int(actual_label)==1 else 'FAIL'}"
        fig.suptitle(
            f"IPI MLAbDev Platform: Prediction interpretability — {barcode}"
            f"\nModel comparison: {' vs '.join(model_names)}{act_part}",
            fontsize=9, y=1.02
        )
        plt.tight_layout()

        if output_path:
            fmt = output_path.rsplit('.', 1)[-1].lower()
            plt.savefig(output_path, dpi=pub_dpi,
                        bbox_inches='tight', format=fmt)
            if logger:
                _log(logger, f"[compare_waterfall] → {output_path}")
            else:
                print(f"[compare_waterfall] → {output_path}")
        return fig


    # ── CDR3 In-silico Mutagenesis ────────────────────────────────────────────

    def cdr3_mutagenesis(self, X_df: pd.DataFrame,
                         embeddings=None,
                         output_dir: str = ".",
                         barcodes: list  = None,
                         lm_name: str    = "",
                         db_name: str    = "",
                         pub_dpi: int    = 300,
                         fmt: str        = "tiff",
                         make_ppt: bool  = True,
                         logger=None):
        """
        In-silico CDR3 mutagenesis scan.
        For each antibody and each CDR3 position, substitute all 20 AAs
        and score with the trained RF model.

        Output per antibody (Nature Biotech quality):
          • {barcode}_cdr3_mutagenesis.{fmt}  — 20×CDR3_len heatmap
        Output overall:
          • cdr3_mutagenesis_all.pptx          — all heatmaps in PPT

        Heatmap:
          rows    = 20 amino acids (ACDEFGHIKLMNPQRSTVWY)
          columns = CDR3 position
          color   = P(PASS) of mutant  (RdBu: red=FAIL, white=0.5, blue=PASS)
          WT residue boxed with black rectangle
        """
        if self.model is None or self.fb_ is None:
            _log(logger, "[Mutagenesis] Model not trained."); return

        os.makedirs(output_dir, exist_ok=True)
        if barcodes is None:
            barcodes = list(X_df.index.astype(str))

        task_str = "CL"  if self.task == 'classification' else "REG"
        lm_str   = lm_name if lm_name else "RF"
        db_str   = db_name if db_name else ""
        saved    = []

        for s_idx in range(len(X_df)):
            bc   = barcodes[s_idx] if s_idx < len(barcodes) else str(s_idx)
            row  = X_df.iloc[s_idx]
            cdr3 = str(row.get('CDR3', '') or '').upper().replace('-', '')
            if not cdr3:
                _log(logger, f"  [Mutagenesis] {bc}: CDR3 missing — skipped")
                continue

            n_pos  = len(cdr3)
            n_aa   = len(AMINO_ACIDS)

            # WT prediction
            wt_score = float(self.predict_proba(
                X_df.iloc[[s_idx]], embeddings=None)[0])

            _log(logger, f"  [Mutagenesis] {bc}  CDR3={cdr3}  "
                         f"WT P(PASS)={wt_score:.4f}  scanning {n_pos}×{n_aa}={n_pos*n_aa} mutants ...")

            # Score matrix: (n_aa, n_pos)
            score_mat = np.full((n_aa, n_pos), np.nan, dtype=np.float32)

            for pos_i in range(n_pos):
                wt_aa = cdr3[pos_i]
                for aa_i, mut_aa in enumerate(AMINO_ACIDS):
                    mut_cdr3 = cdr3[:pos_i] + mut_aa + cdr3[pos_i+1:]
                    mut_row  = row.to_dict()
                    mut_row['CDR3'] = mut_cdr3
                    # Splice into VH if possible
                    vh = str(row.get('HSEQ', '') or '')
                    if vh and wt_aa in vh:
                        cdr3_start = vh.find(cdr3)
                        if cdr3_start >= 0:
                            mut_row['HSEQ'] = (vh[:cdr3_start] +
                                               mut_cdr3 +
                                               vh[cdr3_start + n_pos:])
                    mut_df = pd.DataFrame([mut_row], index=[bc])
                    try:
                        sc = float(self.predict_proba(mut_df, embeddings=None)[0])
                    except Exception:
                        sc = wt_score
                    score_mat[aa_i, pos_i] = sc

            # ── Plot heatmap ──────────────────────────────────────────────────
            # Width: 0.72" per column (enough to show "0.00" at ~6pt)
            # Height: fixed 7" for 20 AA rows
            fig_w = max(10, n_pos * 0.72)
            fig, ax = plt.subplots(figsize=(fig_w, 7))
            import matplotlib.colors as _mc
            if self.task == 'classification':
                cmap  = 'RdBu'
                _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0)
            else:
                cmap  = 'coolwarm'
                _vmin = float(np.nanmin(score_mat))
                _vmax = float(np.nanmax(score_mat))
                _norm = _mc.TwoSlopeNorm(vmin=_vmin, vcenter=(_vmin+_vmax)/2, vmax=_vmax)

            im = ax.imshow(score_mat, cmap=cmap, norm=_norm, aspect='auto')

            # ── Annotate each cell with P(PASS) score ─────────────────────
            # Font size scales with number of cells so it always fits.
            # Text colour: white on dark cells (score < 0.35 or > 0.65),
            # dark on mid (white background near 0.5).
            cell_fsz = max(4.0, min(7.0, 120.0 / max(n_pos, 1)))
            for aa_i in range(n_aa):
                for pos_i in range(n_pos):
                    val = score_mat[aa_i, pos_i]
                    if np.isnan(val):
                        continue
                    # Choose text colour for contrast against RdBu background
                    # red (<0.35) and blue (>0.65) are dark → white text
                    # white middle (0.35–0.65) → dark text
                    txt_col = 'white' if (val < 0.35 or val > 0.65) else '#333333'
                    ax.text(pos_i, aa_i, f"{val:.2f}",
                            ha='center', va='center',
                            fontsize=cell_fsz, color=txt_col,
                            fontweight='bold' if abs(val - 0.5) > 0.3 else 'normal')

            # ── Box WT residue at each position ──────────────────────────────
            for pos_i, wt_aa in enumerate(cdr3):
                if wt_aa in AMINO_ACIDS:
                    aa_i = AMINO_ACIDS.index(wt_aa)
                    rect = plt.Rectangle(
                        (pos_i - 0.5, aa_i - 0.5), 1, 1,
                        linewidth=2.0, edgecolor='black', facecolor='none',
                        zorder=5
                    )
                    ax.add_patch(rect)

            # Axes
            ax.set_xticks(range(n_pos))
            ax.set_xticklabels([f"{cdr3[i]}\n{i+1}" for i in range(n_pos)],
                               fontsize=8)
            ax.set_yticks(range(n_aa))
            ax.set_yticklabels(list(AMINO_ACIDS), fontsize=8)
            ax.set_xlabel('CDR3 position  (WT residue shown)', fontsize=9)
            ax.set_ylabel('Mutant amino acid', fontsize=9)

            # Colorbar
            cbar = plt.colorbar(im, ax=ax, shrink=0.85, pad=0.02)
            cb_lbl = 'P(PASS)' if self.task == 'classification' else 'Predicted score'
            cbar.set_label(cb_lbl, fontsize=9, labelpad=8)
            if self.task == 'classification':
                cbar.set_ticks([0.0, 0.25, 0.5, 0.75, 1.0])
                cbar.set_ticklabels(['0.0\nFAIL', '0.25', '0.50\n(border)', '0.75', '1.0\nPASS'])
            cbar.ax.tick_params(labelsize=7.5)
            # Annotation inside colorbar
            cbar.ax.text(0.5, 0.05,  'FAIL', ha='center', va='bottom',
                         fontsize=7, color='white', transform=cbar.ax.transAxes,
                         fontweight='bold')
            cbar.ax.text(0.5, 0.95,  'PASS', ha='center', va='top',
                         fontsize=7, color='white', transform=cbar.ax.transAxes,
                         fontweight='bold')

            # Title
            ax.set_title(
                f"IPI MLAbDev · CDR3 Mutagenesis Heatmap\n"
                f"ID: {bc}   WT P(PASS)={wt_score:.4f}\n"
                f"RF | {lm_str} | {task_str}" +
                (f" | {db_str}" if db_str else ""),
                fontsize=9, loc='center', pad=8
            )
            plt.tight_layout()

            bc_safe  = bc.replace('/', '_').replace(' ', '_')
            img_path = os.path.join(output_dir,
                                    f"{bc_safe}_cdr3_mutagenesis.{fmt}")
            save_kw  = dict(dpi=pub_dpi, bbox_inches='tight')
            if fmt == 'tiff':  save_kw['format'] = 'tiff'
            elif fmt in ('jpeg','jpg'):
                save_kw['format'] = 'jpeg'
                save_kw['pil_kwargs'] = {'quality': 95}
            plt.savefig(img_path, **save_kw)
            plt.close()
            saved.append(img_path)
            _log(logger, f"    → {img_path}")

        _log(logger, f"[Mutagenesis] {len(saved)} heatmaps saved → {output_dir}/")

        # ── PPT ───────────────────────────────────────────────────────────────
        if make_ppt and saved:
            try:
                from pptx import Presentation as _Prs
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor

                prs = _Prs()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                for img_path, s_idx in zip(saved, range(len(saved))):
                    bc    = barcodes[s_idx] if s_idx < len(barcodes) else str(s_idx)
                    slide = prs.slides.add_slide(blank)
                    img_w = Inches(10)
                    img_h = Inches(6.5)
                    left  = (prs.slide_width  - img_w) / 2
                    top   = Inches(0.4)
                    slide.shapes.add_picture(img_path, left, top,
                                             width=img_w, height=img_h)
                    txb = slide.shapes.add_textbox(
                        Inches(0.15), Inches(7.1),
                        Inches(13.0), Inches(0.35))
                    tf  = txb.text_frame
                    tf.text = (f"{s_idx+1}/{len(saved)}  |  {bc}  |  "
                               f"RF | {lm_str} | {task_str} | {db_str}")
                    p   = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x88, 0x87, 0x80)

                ppt_path = os.path.join(output_dir, 
                            f"cdr3_mutagenesis_{target}_{lm}_{model_type}_{db_stem}_all.pptx")
                prs.save(ppt_path)
                _log(logger,
                     f"[Mutagenesis] PPT ({len(saved)} slides) → {ppt_path}")

            except ImportError:
                _log(logger, "[Mutagenesis] pip install python-pptx for PPT output")
            except Exception as _pe:
                import traceback as _tb
                _log(logger, f"[Mutagenesis] PPT failed — {_pe}\n{_tb.format_exc()}")


        # ── kfold_validation ──────────────────────────────────────────────────────

    @classmethod
    def kfold_validation(
        cls,
        data,
        X_df:         pd.DataFrame,
        y,
        embeddings:   np.ndarray   = None,
        embedding_lm: str          = "ablang",
        title:        str          = "RF",
        kfold:        int          = 10,
        target:       str          = "psr_filter",
        cluster_col:  str          = "HCDR3_CLUSTER_0.8",
        db_stem:      str          = "",
        override_features: dict    = None,
        cost_fn:          float   = 3.0,
        cost_fp:          float   = 1.0,
    ):
        # override_features: dict of feature flags to apply to every fold's config
        # e.g. {'embedding': False, 'biophysical': True, 'kmer': False}
        # Derived automatically from --lm in predict_developability.py
        """
        CDR3-stratified k-fold for classification AND regression.

        Parameters
        ----------
        data         – full DataFrame (BARCODE, cluster_col)
        X_df         – feature DataFrame (HSEQ, CDR3)
        y            – labels / continuous scores
        embeddings   – (n, emb_dim) PLM array or None
        embedding_lm – PLM name tag for filenames
        title        – display title
        kfold        – number of folds
        target       – prediction target column name
        cluster_col  – CDR3 cluster column for leakage-free splits
        db_stem      – database tag for filenames
        """
        os.makedirs(MODEL_DIR, exist_ok=True)

        # ── Logging ───────────────────────────────────────────────────────────
        ts           = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        _db_tag      = f"_{db_stem}" if db_stem else ""
        kfold_actual = kfold   # may be reduced later if fewer clusters than folds
        log_name     = f"kfold_{target}_{embedding_lm}_rf_{db_stem}_k{kfold_actual}_{ts}.log"
        log_path = os.path.join(MODEL_DIR, log_name)
        logger   = _setup_logger(log_path)

        _log(logger, f"[log] {'='*58}")
        _log(logger, f"[log] Started : {datetime.datetime.now()}")
        _log(logger, f"[log] Command : kfold_validation "
                     f"target={target} lm={embedding_lm} k={kfold}")
        _log(logger, f"[log] {'='*58}\n")

        # ── Config ────────────────────────────────────────────────────────────
        y_arr  = np.asarray(y, dtype=float)
        _holder = cls.__new__(cls)
        _holder.config = copy.deepcopy(cls._DEFAULT_CONFIG)
        _cfg_path = "config/random_forest.yaml"
        if os.path.exists(_cfg_path):
            with open(_cfg_path) as f:
                _deep_merge(_holder.config, yaml.safe_load(f) or {})
        _holder.task = _holder.config.get('task','classification')
        # Auto-detect regression if YAML says classification but y is continuous
        _y_unique_rf  = len(set(y_arr.tolist()))
        _is_binary_rf = (_y_unique_rf <= 2 and
                         set(y_arr.tolist()).issubset({0, 1, 0.0, 1.0}))
        if _holder.task == 'classification' and not _is_binary_rf:
            _holder.task = 'regression'
            _holder.config['task'] = 'regression'
            print(f"[kfold] Auto-detected task=regression "
                  f"({_y_unique_rf} unique values — not binary)")
        _holder._resolve_config(y_arr)
        resolved_cfg = _holder.config
        task = resolved_cfg.get('task','classification')
        del _holder

        data = data.copy()
        data['BARCODE'] = data.index.astype(str).tolist()
        n = len(y_arr)

        # ── Print config once at kfold start (not once per fold) ─────────────
        _log(logger, "")
        RandomForestModel.print_config_report(resolved_cfg)
        _log(logger, f"[kfold] n={n:,}  task={task}  lm={embedding_lm}  "
                     f"target={target}")

        # ── Compute cluster column if missing ─────────────────────────────────
        # Mirrors transformer_lm — CDR3-stratified splits prevent leakage.
        if cluster_col not in data.columns:
            if 'CDR3' in data.columns:
                _log(logger, f"[kfold] '{cluster_col}' not found — computing from CDR3 ...")
                try:
                    from utils.clustering import greedy_clustering_by_levenshtein
                    _thresh = 0.8
                    try:
                        _thresh = float(cluster_col.rsplit('_', 1)[-1])
                    except Exception:
                        pass
                    data[cluster_col] = greedy_clustering_by_levenshtein(
                        data['CDR3'].tolist(), _thresh)
                    _log(logger, f"[kfold] {data[cluster_col].nunique():,} clusters "
                                 f"(threshold={_thresh})")
                except Exception as _ce:
                    _log(logger, f"[kfold] WARNING: clustering failed ({_ce}) "
                                 f"— falling back to StratifiedKFold")
            else:
                _log(logger, f"[kfold] WARNING: '{cluster_col}' not in data and "
                             f"no CDR3 column — using StratifiedKFold (no leakage protection)")

        # ── Stratifier ────────────────────────────────────────────────────────
        # For regression: stratify on median-split (above/below median)
        y_strat = (y_arr > np.median(y_arr)).astype(int)
        kfold_actual = kfold

        if cluster_col in data.columns:
            groups = data[cluster_col].values
            n_grp  = len(np.unique(groups))
            if n_grp < kfold_actual:
                kfold_actual = n_grp
                _log(logger, f"[kfold] WARNING: {n_grp} clusters → "
                             f"reducing folds to {kfold_actual}")
            splitter   = StratifiedGroupKFold(n_splits=kfold_actual,
                                              shuffle=True, random_state=42)
            split_iter = splitter.split(np.arange(n), y_strat, groups)
            _log(logger, f"[kfold] StratifiedGroupKFold on '{cluster_col}' "
                         f"({n_grp} clusters, {kfold_actual} folds)")
        else:
            _log(logger, f"[kfold] WARNING: '{cluster_col}' not found → StratifiedKFold")
            splitter   = StratifiedKFold(n_splits=kfold_actual,
                                         shuffle=True, random_state=42)
            split_iter = splitter.split(np.arange(n), y_strat)

        # ── Tracking ──────────────────────────────────────────────────────────
        mean_fpr        = np.linspace(0,1,100)
        tprs            = []
        primary_metrics = []   # AUC (classif) or R² (regression)
        fold_metrics    = []
        all_records     = []
        best_metric     = -1.0
        best_fold_num   = -1
        best_fold_model = None
        last_va_idx     = None  # for SHAP on best fold

        plt.figure(figsize=(8,7))
        _log(logger, f"\n[kfold] {kfold_actual}-fold | {title} | {target.upper()}")

        for fold, (tr_idx, va_idx) in enumerate(split_iter, 1):
            last_va_idx = va_idx
            _log(logger, f"\n── Fold {fold}/{kfold_actual} ──")

            # Leakage check
            if cluster_col in data.columns:
                tr_g  = set(groups[tr_idx]); va_g = set(groups[va_idx])
                ok    = not bool(tr_g & va_g)
                _log(logger, f"  {'[OK]  No CDR3 leakage' if ok else f'[WARN] {len(tr_g & va_g)} leaked'} | "
                             f"train_clusters={len(tr_g)}  val_clusters={len(va_g)}")

            y_tr = y_arr[tr_idx]; y_va = y_arr[va_idx]
            if task == 'classification':
                _log(logger, f"  Train={len(tr_idx):,} pos={y_tr.mean():.1%}  "
                             f"Val={len(va_idx):,} pos={y_va.mean():.1%}")
            else:
                _log(logger, f"  Train={len(tr_idx):,} μ={y_tr.mean():.3f} σ={y_tr.std():.3f}  "
                             f"Val={len(va_idx):,} μ={y_va.mean():.3f} σ={y_va.std():.3f}")

            # Build fold model — verbose=False suppresses per-fold config table
            inst      = cls(config=resolved_cfg, verbose=False)
            # Apply feature flag overrides from --lm (e.g. biophysical, kmer)
            if override_features:
                # Strip internal keys before applying to features dict
                _oh_seq = override_features.pop('_onehot_sequence', None)
                inst.config['features'].update(override_features)
                if _oh_seq:
                    inst.config.setdefault('onehot', {})['sequence'] = _oh_seq
            inst.task = task
            inst.fb_  = FeatureBuilder(inst.config)   # use inst.config (with overrides)
            # Apply LM-specific hyperparameter profile per fold
            inst.apply_lm_profile(embedding_lm, logger=logger)

            X_df_tr = X_df.iloc[tr_idx]; X_df_va = X_df.iloc[va_idx]
            emb_tr  = embeddings[tr_idx] if embeddings is not None else None
            emb_va  = embeddings[va_idx] if embeddings is not None else None

            X_tr = inst.fb_.fit_transform(X_df_tr, emb_tr)
            X_va = inst.fb_.transform(X_df_va, emb_va)
            _log(logger, f"  features={X_tr.shape[1]}")

            rf_sk = inst._build_sklearn_model()
            if resolved_cfg['training'].get('hyperparam_search'):
                _log(logger, "  [search] RandomizedSearchCV...")
                inst.model = inst._hyperparam_search(rf_sk, X_tr, y_tr, logger)
            else:
                inst.model = rf_sk
                inst.model.fit(X_tr, y_tr.astype(float if task=='regression' else int))

            oob = getattr(inst.model,'oob_score_',None)
            if oob is not None:
                _log(logger, f"  OOB = {oob:.4f}")

            # ── Evaluate ──────────────────────────────────────────────────────
            barcodes = data.iloc[va_idx]['BARCODE'].tolist()

            if task == 'classification':
                if len(np.unique(y_va.astype(int))) < 2:
                    _log(logger, f"  Skipping fold {fold} — single class in val.")
                    continue

                prob = inst.model.predict_proba(X_va)[:,1]
                pred = (prob >= 0.5).astype(int)

                fold_auc  = roc_auc_score(y_va.astype(int), prob)
                fold_acc  = accuracy_score(y_va.astype(int), pred)
                fold_f1   = f1_score(y_va.astype(int), pred, zero_division=0)
                fold_prec = precision_score(y_va.astype(int), pred, zero_division=0)
                fold_rec  = recall_score(y_va.astype(int), pred, zero_division=0)
                fold_rcf  = recall_score(y_va.astype(int), pred,
                                         pos_label=0, zero_division=0)

                _log(logger,
                     f"  Fold {fold} → AUC={fold_auc:.4f}  Acc={fold_acc:.4f}  "
                     f"F1={fold_f1:.4f}  Prec={fold_prec:.4f}  "
                     f"Rec={fold_rec:.4f}  Rec(Fail)={fold_rcf:.4f}")

                fold_metrics.append({'fold':fold,'auc':fold_auc,'acc':fold_acc,
                                     'f1':fold_f1,'precision':fold_prec,
                                     'recall':fold_rec,'rec_fail':fold_rcf})
                primary_metrics.append(fold_auc)

                fpr,tpr,_ = roc_curve(y_va.astype(int), prob)
                tprs.append(np.interp(mean_fpr,fpr,tpr)); tprs[-1][0]=0.0
                plt.plot(fpr,tpr,alpha=0.3,lw=1,
                         label=f'Fold {fold} ({fold_auc:.3f})')

                for bc,tv,pv,prv in zip(barcodes,y_va,pred,prob):
                    all_records.append({'BARCODE':bc,'fold':fold,
                                        'true':tv,'pred':pv,'prob':prv})
                pm = fold_auc

            else:  # regression
                preds_r  = inst.model.predict(X_va)
                fold_r2  = r2_score(y_va, preds_r)
                fold_rp  = pearsonr(y_va, preds_r)[0]
                fold_rs  = spearmanr(y_va, preds_r)[0]
                fold_mae = mean_absolute_error(y_va, preds_r)
                fold_rmse= mean_squared_error(y_va, preds_r)**0.5

                _log(logger,
                     f"  Fold {fold} → R²={fold_r2:.4f}  Pearson={fold_rp:.4f}  "
                     f"Spearman={fold_rs:.4f}  MAE={fold_mae:.4f}  RMSE={fold_rmse:.4f}")

                fold_metrics.append({'fold':fold,'r2':fold_r2,'pearson':fold_rp,
                                     'spearman':fold_rs,'mae':fold_mae,'rmse':fold_rmse})
                primary_metrics.append(fold_r2)

                for bc,tv,pv in zip(barcodes,y_va,preds_r):
                    all_records.append({'BARCODE':bc,'fold':fold,
                                        'true':tv,'pred':pv,'prob':pv})
                pm = fold_r2

            # Save fold checkpoint
            fold_path = os.path.join(
                MODEL_DIR,
                f"{target}_{embedding_lm}_rf_{db_stem}_k{kfold_actual}_fold{fold}.pkl")
            inst.save(fold_path)

            if pm > best_metric:
                best_metric     = pm
                best_fold_num   = fold
                best_fold_model = inst

        if not primary_metrics:
            _log(logger,"[kfold] No valid folds."); return None,None,None,None,None,None

        # ── Save best model ───────────────────────────────────────────────────
        _task_tag_rf = "_regression" if task == 'regression' else ""
        best_path = None
        if best_fold_model is not None:
            best_path = os.path.join(
                MODEL_DIR,
                f"BEST_{target}_{embedding_lm}_rf_{db_stem}_k{kfold_actual}_fold{best_fold_num}{_task_tag_rf}.pkl")
            best_fold_model.save(best_path)
            _log(logger, f"\n[kfold] Best fold → {best_path}"
                         f"  (fold={best_fold_num}, metric={best_metric:.4f})")

        # ── Save fold predictions ─────────────────────────────────────────────
        pred_path = None
        if all_records:
            pred_path = os.path.join(
                MODEL_DIR,
                f"fold_preds_{target}_{embedding_lm}_rf_{db_stem}_k{kfold_actual}{_task_tag_rf}.csv")
            df_p = pd.DataFrame(all_records)
            df_p['best_fold'] = (df_p['fold']==best_fold_num).astype(int)
            df_p.to_csv(pred_path, index=False)
            _log(logger, f"[kfold] Fold predictions → {pred_path}")

        # ── Summary ───────────────────────────────────────────────────────────
        sep = '═'*62
        _log(logger, f"\n{sep}")
        _log(logger, f"  {kfold_actual}-FOLD CV — {target.upper()}  [{task}]")
        _log(logger, f"{'─'*62}")

        if task == 'classification':
            mean_auc = float(np.mean(primary_metrics))
            std_auc  = float(np.std(primary_metrics))
            means    = {k: float(np.mean([m[k] for m in fold_metrics]))
                        for k in ('acc','f1','precision','recall','rec_fail')}
            _log(logger,
                 f"  {'Fold':>5}  {'AUC':>7}  {'Acc':>7}  {'F1':>7}  "
                 f"{'Prec':>7}  {'Rec':>7}  {'Rec(F)':>7}")
            _log(logger, f"  {'─'*5}  "+"  ".join(["─"*7]*6))
            for m in fold_metrics:
                mk = " ←" if m['fold']==best_fold_num else ""
                _log(logger,
                     f"  {m['fold']:5d}  {m['auc']:7.4f}  {m['acc']:7.4f}  "
                     f"{m['f1']:7.4f}  {m['precision']:7.4f}  "
                     f"{m['recall']:7.4f}  {m['rec_fail']:7.4f}{mk}")
            _log(logger, f"{'─'*62}")
            _log(logger,
                 f"  {'Mean':>5}  {mean_auc:7.4f}  {means['acc']:7.4f}  "
                 f"{means['f1']:7.4f}  {means['precision']:7.4f}  "
                 f"{means['recall']:7.4f}  {means['rec_fail']:7.4f}")
            _log(logger, f"  {'±Std':>5}  {std_auc:7.4f}")
            _log(logger, f"  Best fold : {best_fold_num}  (AUC={best_metric:.4f})")
            _log(logger, f"  Rec(Fail) : {means['rec_fail']:.4f}  ← minority class recall")

            # ROC plot
            if tprs:
                mean_tpr = np.mean(tprs,axis=0); mean_tpr[-1]=1.0
                std_tpr  = np.std(tprs,axis=0)
                plt.plot(mean_fpr,mean_tpr,'b',lw=3,
                         label=f'Mean (AUC={mean_auc:.3f}±{std_auc:.3f})')
                plt.fill_between(mean_fpr,
                                 np.maximum(mean_tpr-std_tpr,0),
                                 np.minimum(mean_tpr+std_tpr,1),
                                 color='lightblue',alpha=0.3,label='±1 std')
                plt.plot([0,1],[0,1],'--',color='gray',lw=0.8)
                plt.xlim([0,1]); plt.ylim([0,1.05])
                plt.xlabel('False Positive Rate', fontsize=9)
                plt.ylabel('True Positive Rate',  fontsize=9)
                plt.title(
                    f'{title} — {target.upper()} | {embedding_lm}\n'
                    f'{kfold_actual}-Fold SGKF ROC\n'
                    f'Acc={means["acc"]:.3f}  F1={means["f1"]:.3f}  '
                    f'Prec={means["precision"]:.3f}  Rec={means["recall"]:.3f}  '
                    f'Rec(Fail)={means["rec_fail"]:.3f}',
                    fontsize=8)
                plt.legend(loc='lower right', fontsize=5)
                plt.grid(alpha=0.3); plt.tight_layout()
                roc_path = os.path.join(
                    MODEL_DIR,
                    f"CV_ROC_{target}_{embedding_lm}_rf_{db_stem}_k{kfold_actual}.png")
                plt.savefig(roc_path,dpi=150,bbox_inches='tight')
                plt.close()
                _log(logger, f"  ROC plot  : {roc_path}")

        else:
            # Regression summary
            means_r = {k: float(np.mean([m[k] for m in fold_metrics]))
                       for k in ('r2','pearson','spearman','mae','rmse')}
            std_r2  = float(np.std([m['r2'] for m in fold_metrics]))
            _log(logger,
                 f"  {'Fold':>5}  {'R²':>7}  {'Pearson':>8}  "
                 f"{'Spearman':>9}  {'MAE':>7}  {'RMSE':>7}")
            _log(logger, f"  {'─'*5}  "+"  ".join(["─"*7]*5))
            for m in fold_metrics:
                mk = " ←" if m['fold']==best_fold_num else ""
                _log(logger,
                     f"  {m['fold']:5d}  {m['r2']:7.4f}  {m['pearson']:8.4f}  "
                     f"{m['spearman']:9.4f}  {m['mae']:7.4f}  {m['rmse']:7.4f}{mk}")
            _log(logger, f"{'─'*62}")
            _log(logger,
                 f"  {'Mean':>5}  {means_r['r2']:7.4f}  {means_r['pearson']:8.4f}  "
                 f"{means_r['spearman']:9.4f}  {means_r['mae']:7.4f}  "
                 f"{means_r['rmse']:7.4f}")
            _log(logger, f"  {'±Std':>5}  {std_r2:7.4f}")
            _log(logger, f"  Best fold : {best_fold_num}  (R²={best_metric:.4f})")
            plt.close()

        _log(logger, f"{sep}\n")

        # ── Threshold optimisation (classification only) ───────────────────────
        if task=='classification' and _THRESHOLD_OPT_AVAILABLE and pred_path:
            _log(logger, f"\n[threshold] Starting optimisation  ···")
            try:
                stability = run_full_threshold_pipeline(
                    fold_preds_csv=pred_path, target=target, lm=embedding_lm,
                    model='rf', db_stem=db_stem or '',
                    best_ckpt_path=best_path, output_dir=MODEL_DIR,
                    cost_fp=cost_fp, cost_fn=cost_fn)
                rec_t = stability['pooled_threshold']
                cls.recommended_threshold = float(rec_t)
                _log(logger, f"  Pooled OOF threshold : {rec_t:.4f}")
            except Exception as e:
                _log(logger, f"[threshold] WARNING: {e} — defaulting to 0.5")

        # ── SHAP (kmer / biophysical features) ───────────────────────────────
        sh_cfg   = resolved_cfg.get('shap',{})
        has_ne   = (resolved_cfg['features'].get('kmer') or
                    resolved_cfg['features'].get('biophysical'))
        if sh_cfg.get('enabled') and has_ne and best_fold_model is not None:
            _log(logger, f"\n[SHAP] Running on best fold ({best_fold_num})...")
            # Use last val split as representative sample
            va_X  = X_df.iloc[last_va_idx] if last_va_idx is not None else X_df
            va_em = (embeddings[last_va_idx]
                     if (embeddings is not None and last_va_idx is not None)
                     else embeddings)
            best_fold_model.shap_analysis(
                va_X, va_em,
                output_prefix=f"{target}_{embedding_lm}_rf{_db_tag}",
                logger=logger)
        elif sh_cfg.get('enabled') and not has_ne:
            _log(logger, "\n[SHAP] Embedding-only — SHAP skipped. "
                         "Use IG from transformer_lm for embedding interpretation.")

        _log(logger, f"\n[log] Finished: {datetime.datetime.now()}")
        _log(logger, f"[log] → {log_path}")
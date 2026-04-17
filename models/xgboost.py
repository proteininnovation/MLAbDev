# models/xgboost.py
# XGBoost Antibody Developability Model
# IPI Antibody Developability Prediction Platform — Production Version DEC-2025
#
# Full parity with randomforest.py:
#   FeatureBuilder (biophysical / kmer / onehot / PLM embeddings)
#   SHAP (bar, beeswarm, heatmap, waterfall, PPT)
#   CDR3 mutagenesis heatmaps
#   Threshold optimisation pipeline
#   LM-specific hyperparameter profiles
#   CDR3-cluster stratified k-fold (no sequence leakage)
#   Log file mirroring
#
# XGBoost-specific differences vs RF:
#   scale_pos_weight instead of class_weight='balanced'
#   colsample_bytree instead of max_features
#   early_stopping_rounds + eval_set (optional)
#   No OOB score (use val AUC instead)
#   TreeExplainer works natively (same SHAP API as RF)
#
# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLES
# ══════════════════════════════════════════════════════════════════════════════
#
#   # K-fold validation
#   python predict_developability.py --kfold 10 \
#       --target psr_filter --lm biophysical --model xgboost \
#       --db data/ipi_psr_trainset.xlsx --cost_fn 3.0
#
#   # Train on full dataset
#   python predict_developability.py --train \
#       --target psr_filter --lm ablang --model xgboost \
#       --db data/ipi_psr_trainset.xlsx
#
#   # Predict + SHAP + CDR3 mutagenesis
#   python predict_developability.py --predict data/new_cohort.xlsx \
#       --target psr_filter --lm biophysical --model xgboost \
#       --db data/ipi_psr_trainset.xlsx --mutagenesis 50
#
# ══════════════════════════════════════════════════════════════════════════════


from config import MODEL_DIR

import os
import copy
import datetime
import logging
import joblib
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import xgboost as xgb
from pathlib import Path
from sklearn.model_selection import (
    RandomizedSearchCV, StratifiedKFold, StratifiedGroupKFold)
from sklearn.metrics import (
    roc_curve, auc, f1_score, accuracy_score, roc_auc_score,
    precision_score, recall_score, confusion_matrix)
from scipy.stats import uniform, randint
from collections import Counter
from itertools import product
from collections import Counter
from itertools import product
import warnings
import yaml

try:
    import shap as _shap_lib
    _SHAP_AVAILABLE = True
except ImportError:
    _SHAP_AVAILABLE = False

try:
    from utils.threshold_optimizer import run_full_threshold_pipeline
    _THRESHOLD_OPT_AVAILABLE = True
except ImportError:
    _THRESHOLD_OPT_AVAILABLE = False


# ══════════════════════════════════════════════════════════════════════════
# AMINO ACID CONSTANTS & BIOPHYSICAL SCALES
# ══════════════════════════════════════════════════════════════════════════

AMINO_ACIDS = 'ACDEFGHIKLMNPQRSTVWY'

_PKA = {'D': 3.9, 'E': 4.1, 'C': 8.3, 'Y': 10.1,
        'H': 6.0, 'K': 10.5, 'R': 12.5,
        'N_term': 8.0, 'C_term': 3.1}

_KD = {'A': 1.8, 'R':-4.5, 'N':-3.5, 'D':-3.5, 'C': 2.5,
       'Q':-3.5, 'E':-3.5, 'G':-0.4, 'H':-3.2, 'I': 4.5,
       'L': 3.8, 'K':-3.9, 'M': 1.9, 'F': 2.8, 'P':-1.6,
       'S':-0.8, 'T':-0.7, 'W':-0.9, 'Y':-1.3, 'V': 4.2}


# ══════════════════════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════════════════
# FEATURE HELPERS + FeatureBuilder (shared with RandomForestModel)
# ══════════════════════════════════════════════════════════════════════════
def _charge_at_ph7(seq):
    c = 0.0
    for aa in seq.upper():
        if   aa == 'D': c -= 1/(1+10**(_PKA['D']-7))
        elif aa == 'E': c -= 1/(1+10**(_PKA['E']-7))
        elif aa == 'H': c += 1/(1+10**(7-_PKA['H']))
        elif aa == 'K': c += 1/(1+10**(7-_PKA['K']))
        elif aa == 'R': c += 1/(1+10**(7-_PKA['R']))
    return round(c, 4)

def _gravy(seq):
    vals = [_KD.get(aa, 0.0) for aa in seq.upper()]
    return round(sum(vals)/len(vals), 4) if vals else 0.0

def _pi(seq):
    def _c(ph):
        c = 1/(1+10**(ph-_PKA['N_term'])) - 1/(1+10**(_PKA['C_term']-ph))
        for aa in seq.upper():
            if aa=='D': c -= 1/(1+10**(_PKA['D']-ph))
            elif aa=='E': c -= 1/(1+10**(_PKA['E']-ph))
            elif aa=='H': c += 1/(1+10**(ph-_PKA['H']))
            elif aa=='K': c += 1/(1+10**(ph-_PKA['K']))
            elif aa=='R': c += 1/(1+10**(ph-_PKA['R']))
        return c
    lo, hi = 0.0, 14.0
    for _ in range(200):
        mid = (lo+hi)/2
        if _c(mid) > 0: lo = mid
        else: hi = mid
    return round((lo+hi)/2, 3)

def _instability(seq):
    # Guruprasad dipeptide instability weights (abbreviated key set)
    DIWV = {'WC':1,'WM':24.68,'CM':33.6,'CH':33.6,'CD':20.26,'CT':33.6,
            'CL':20.26,'CP':20.26,'QD':20.26,'RD':20.26,'RH':20.26,
            'DG':20.26,'DS':20.26,'FD':54.96,'GD':20.26,'GH':33.6,
            'HN':20.26,'LR':20.26,'ND':20.26,'NL':20.26,'NP':20.26,
            'QR':20.26,'RR':58.28,'SR':20.26,'TD':20.26,'TP':20.26,
            'VN':1,'WR':1}
    seq = seq.upper()
    if len(seq) < 2: return 0.0
    total = sum(DIWV.get(seq[i:i+2], 1.0) for i in range(len(seq)-1))
    return round(10.0/len(seq)*total, 3)

def _strip_cdr3_loop(cdr3: str) -> str:
    """
    Strip conserved framework anchor residues from CDR3:
      N-terminal : remove leading 'CAR' or 'C' (conserved cysteine + Ala-Arg)
      C-terminal : remove trailing 'FDY', 'FDW', 'MDY', 'LDY', 'MDW' (J-gene motif)
    Returns only the hypervariable loop residues.

    Example:  CARGFDYW  →  strip CAR + FDY  →  G
              CARDRGYYY →  strip CAR        →  DRGYYY
              CARRGFDYW →  strip CAR + FDY  →  R
    """
    s = cdr3.upper().replace('-', '')
    # Strip N-terminal conserved residues: leading C or CAR
    if s.startswith('CAR'):
        s = s[3:]
    elif s.startswith('C'):
        s = s[1:]
    # Strip C-terminal J-gene motif: xDY or xDW (last 3 residues)
    # Common patterns: FDY, FDW, MDY, MDW, LDY, WDY, YDY
    if len(s) >= 3 and s[-2] == 'D' and s[-1] in 'YW':
        s = s[:-3]
    return s


def compute_biophysical_features(seq: str, feature_list: list,
                                  vh_seq: str = '') -> dict:
    """
    Compute biophysical features from CDR3 loop sequence.
    seq     : CDR3 with conserved flanks already stripped by _strip_cdr3_loop()
    vh_seq  : full VH sequence (for vh_length and ratio features)

    Feature naming convention:
      cdr3_{property}  e.g. cdr3_charge_ph7, cdr3_R, cdr3_length
      vh_length        length of full VH sequence
      cdr3_vh_length_ratio  CDR3 loop / VH length
    """
    seq = seq.upper().replace('-', '')
    if not seq:
        return {f: 0.0 for f in feature_list}
    counts = Counter(seq); n = len(seq)
    vh_len = len(vh_seq.upper().replace('-', '')) if vh_seq else 0
    out = {}
    for f in feature_list:
        if   f == 'length':               out[f] = float(n)
        elif f == 'pi':                   out[f] = _pi(seq)
        elif f in ('charge_ph7', 'charge'): out[f] = _charge_at_ph7(seq)
        elif f == 'hydrophobicity':       out[f] = _gravy(seq)
        elif f == 'aromaticity':          out[f] = round(sum(counts.get(a,0) for a in 'FYW')/n, 4)
        elif f == 'instability':          out[f] = _instability(seq)
        elif f == 'net_charge_sq':
            c = _charge_at_ph7(seq);      out[f] = round(c*c, 4)
        elif f == 'frac_charged':         out[f] = round(sum(counts.get(a,0) for a in 'DERKH')/n, 4)
        elif f == 'frac_hydrophobic':     out[f] = round(sum(counts.get(a,0) for a in 'AVILMFWP')/n, 4)
        elif f == 'cdr3_length':          out[f] = float(n)
        elif f == 'vh_length':            out[f] = float(vh_len)
        elif f == 'vh_cdr3_length_ratio': out[f] = round(vh_len / n, 4) if n > 0 else 0.0
        elif f in ('vh_charge_ph7', 'vh_charge'): out[f] = _charge_at_ph7(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f == 'vh_hydrophobicity':    out[f] = _gravy(vh_seq.upper().replace('-','')) if vh_seq else 0.0
        elif f.startswith('count_'):      out[f] = float(counts.get(f.split('_')[1], 0))
        # Single amino acid counts — e.g. 'R', 'K', 'W', 'D'
        elif len(f) == 1 and f in 'ACDEFGHIKLMNPQRSTVWY':
                                          out[f] = float(counts.get(f, 0))
        else:                             out[f] = 0.0
    return out

def compute_kmer_features(seq: str, k_list: list, normalize: bool=True) -> dict:
    seq = seq.upper().replace('-','')
    out = {}
    for k in k_list:
        all_kmers = [''.join(p) for p in product(AMINO_ACIDS, repeat=k)]
        cnts = Counter(seq[i:i+k] for i in range(max(0, len(seq)-k+1))) if len(seq)>=k else {}
        total = sum(cnts.values()) or 1
        for km in all_kmers:
            v = cnts.get(km, 0)
            out[f'{k}mer_{km}'] = v/total if normalize else float(v)
    return out



# ── ASCII lookup for fast one-hot (same as transformer_onehot.py) ─────────────
_AA_LOOKUP_RF = np.full(128, 20, dtype=np.int32)
for _aa_rf, _idx_rf in {aa: i for i, aa in enumerate(AMINO_ACIDS)}.items():
    _AA_LOOKUP_RF[ord(_aa_rf)] = _idx_rf


def compute_onehot_features(seq: str, max_len: int) -> np.ndarray:
    """
    One-hot encode a sequence into a flattened vector of shape (max_len * 20,).
    Unknown / padding positions → all-zero row.
    """
    seq = seq.upper().replace('-', '')[:max_len]
    out = np.zeros((max_len, len(AMINO_ACIDS)), dtype=np.float32)
    if seq:
        chars  = np.frombuffer(seq.encode('ascii'), dtype=np.uint8)
        aa_idx = _AA_LOOKUP_RF[np.clip(chars, 0, 127)]
        pos    = np.arange(len(chars))
        valid  = aa_idx < 20
        out[pos[valid], aa_idx[valid]] = 1.0
    return out.flatten()   # (max_len * 20,)


def onehot_feature_names(prefix: str, max_len: int) -> list:
    """
    Generate feature names for a one-hot encoded sequence.
    Format: {prefix}_pos{N}_{AA}  e.g. H_pos001_A, H_pos001_C, ...
    """
    names = []
    for pos in range(1, max_len + 1):
        for aa in AMINO_ACIDS:
            names.append(f"{prefix}_pos{pos:03d}_{aa}")
    return names


# ══════════════════════════════════════════════════════════════════════════════
# 3.  FEATURE BUILDER
# ══════════════════════════════════════════════════════════════════════════════

class FeatureBuilder:
    """
    Builds combined feature matrix for RandomForestModel.

    Supported feature types (any combination):
      • embedding   – PLM vector (ablang 480d, antiberta2 1024d, ...)
      • kmer        – k-mer frequencies from VH | CDR3 | VHVL
      • biophysical – pI, charge, hydrophobicity, AA counts, VH-level
      • onehot      – position × AA one-hot from HCDR3 | VH | VHVL

    YAML keys:
      features.embedding / kmer / biophysical / onehot  : bool
      kmer.sequence     : CDR3 | VH | VHVL
      kmer.k            : [1, 2]
      kmer.normalize    : true
      biophysical.features : [length, pi, charge_ph7, ...]
      onehot.sequence   : HCDR3 | VH | VHVL  (default HCDR3)
    """

    def __init__(self, config: dict):
        self.cfg        = config
        self.feat_cfg   = config.get('features', {})
        self.km_cfg     = config.get('kmer',     {})
        self.bp_cfg     = config.get('biophysical', {})
        self.oh_cfg     = config.get('onehot',   {})
        # Ensure onehot key exists in feat_cfg — may be absent from old YAML
        self.feat_cfg.setdefault('onehot', False)

        self.feature_names_ : list = None
        self._emb_dim               = None

        # Max sequence lengths — fitted from training data
        self._oh_max_vh   = 0
        self._oh_max_vl   = 0
        self._oh_max_cdr3 = 0

    # ── Sequence accessors ────────────────────────────────────────────────────

    def _get_kmer_seq(self, row) -> str:
        """Sequence used for k-mer features (kmer.sequence in YAML)."""
        src = self.km_cfg.get('sequence', 'CDR3').upper()
        vh  = str(row.get('HSEQ', '') or '')
        vl  = str(row.get('LSEQ', '') or '')
        cdr = str(row.get('CDR3', '') or '')
        if src == 'VH':   return vh
        if src == 'VHVL': return vh + vl
        return cdr  # default CDR3

    def _get_cdr3(self, row) -> str:
        return str(row.get('CDR3', '') or '')

    def _get_vh(self, row) -> str:
        return str(row.get('HSEQ', '') or '')

    def _get_vl(self, row) -> str:
        return str(row.get('LSEQ', '') or '')

    def _oh_segs(self, row) -> list:
        """
        Return list of (seq, max_len, prefix) for one-hot encoding.
        Controlled by onehot.sequence in YAML: HCDR3 | VH | VHVL
        """
        mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
        vh   = self._get_vh(row)
        vl   = self._get_vl(row)
        cdr  = self._get_cdr3(row)
        if mode == 'VH':
            return [(vh, self._oh_max_vh, 'oh_vh')]
        elif mode == 'VHVL':
            return [(vh, self._oh_max_vh, 'oh_vh'),
                    (vl, self._oh_max_vl, 'oh_vl')]
        else:  # HCDR3 (default)
            return [(cdr, self._oh_max_cdr3, 'oh_hcdr3')]

    # ── Chain tag for output file naming ─────────────────────────────────────

    @property
    def chain_tag(self) -> str:
        """
        Short tag added to output filenames for kmer/onehot modes.
        e.g. 'HCDR3', 'VH', 'VHVL'
        """
        if self.feat_cfg.get('kmer') and not self.feat_cfg.get('onehot'):
            return self.km_cfg.get('sequence', 'CDR3').upper()
        if self.feat_cfg.get('onehot'):
            return self.oh_cfg.get('sequence', 'HCDR3').upper()
        return ''

    # ── Fit ───────────────────────────────────────────────────────────────────

    def fit(self, X_df: pd.DataFrame, embeddings: np.ndarray = None):
        names = []

        # Embedding
        if self.feat_cfg.get('embedding'):
            if embeddings is None:
                raise ValueError(
                    "features.embedding=True but embeddings=None.\n"
                    "  Pass embeddings=X.values to train()/kfold_validation().\n"
                    "  Or set features.embedding: false in YAML.")
            self._emb_dim = embeddings.shape[1]
            names += [f'emb_{i}' for i in range(self._emb_dim)]

        # K-mer
        if self.feat_cfg.get('kmer'):
            kf    = compute_kmer_features(
                        self._get_kmer_seq(X_df.iloc[0]),
                        self.km_cfg.get('k', [1, 2]),
                        self.km_cfg.get('normalize', True))
            names += list(kf.keys())

        # Biophysical
        if self.feat_cfg.get('biophysical'):
            bf = compute_biophysical_features(
                     self._get_cdr3(X_df.iloc[0]),
                     self.bp_cfg.get('features', []),
                     vh_seq=self._get_vh(X_df.iloc[0]))
            def _rename(k):
                k = k.replace('count_', '').replace('_ph7', '')
                if k.startswith('vh_') or k.startswith('cdr3_'):
                    return k
                return f'cdr3_{k}'
            names += [_rename(k) for k in bf.keys()]

        # One-hot — fit max lengths from training data, capped by YAML max_lengths.
        # Priority:
        #   1. YAML onehot.max_lengths.{vh|vl|hcdr3} — hard cap (safety ceiling)
        #   2. Actual max observed in training data    — fitted value (stored in pkl)
        #   3. Fallback defaults (135/135/25)
        #
        # At predict time the saved _oh_max_* values from the pkl are used directly,
        # so predict sequences longer than training max are truncated with a warning.
        if self.feat_cfg.get('onehot'):
            mode     = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _max_cfg = self.oh_cfg.get('max_lengths', {})
            # YAML caps (0 or missing → no cap applied)
            _cap_vh   = int(_max_cfg.get('vh',   0) or 0)
            _cap_vl   = int(_max_cfg.get('vl',   0) or 0)
            _cap_cdr3 = int(_max_cfg.get('hcdr3',0) or 0)

            if mode in ('VH', 'VHVL'):
                _data_max_vh   = max((len(str(r.get('HSEQ','') or ''))
                                      for _, r in X_df.iterrows()), default=135)
                self._oh_max_vh  = _cap_vh  if _cap_vh  > 0 else _data_max_vh
                print(f"  [onehot] VH  max_len={self._oh_max_vh:3d}  "
                      f"(data_max={_data_max_vh}"
                      + (f"  capped at {_cap_vh}" if _cap_vh > 0 else "") + ")")
            if mode == 'VHVL':
                _data_max_vl   = max((len(str(r.get('LSEQ','') or ''))
                                      for _, r in X_df.iterrows()), default=135)
                self._oh_max_vl  = _cap_vl  if _cap_vl  > 0 else _data_max_vl
                print(f"  [onehot] VL  max_len={self._oh_max_vl:3d}  "
                      f"(data_max={_data_max_vl}"
                      + (f"  capped at {_cap_vl}" if _cap_vl > 0 else "") + ")")
            if mode == 'HCDR3':
                _data_max_cdr3 = max((len(str(r.get('CDR3','') or ''))
                                      for _, r in X_df.iterrows()), default=25)
                self._oh_max_cdr3 = _cap_cdr3 if _cap_cdr3 > 0 else _data_max_cdr3
                print(f"  [onehot] CDR3 max_len={self._oh_max_cdr3:3d}  "
                      f"(data_max={_data_max_cdr3}"
                      + (f"  capped at {_cap_cdr3}" if _cap_cdr3 > 0 else "") + ")")

            # Feature names
            row0 = X_df.iloc[0]
            for _seq, max_len, pfx in self._oh_segs(row0):
                for pos in range(1, max_len + 1):
                    for aa in AMINO_ACIDS:
                        names.append(f"{pfx}_{pos}_{aa}")

        if not names:
            raise ValueError(
                "No features selected — enable at least one of: "
                "embedding | kmer | biophysical | onehot")
        self.feature_names_ = names

        # ── Print feature dimension report ────────────────────────────────────
        _AA = len(AMINO_ACIDS)   # 20
        _SEP = '─' * 52
        print(f"[FeatureBuilder] Feature dimensions:")
        print(f"  {_SEP}")
        _total = 0

        if self.feat_cfg.get('embedding') and self._emb_dim:
            _d = self._emb_dim
            print(f"  embedding       : {_d:>6,} d  (PLM vector)")
            _total += _d

        if self.feat_cfg.get('kmer'):
            _k   = self.km_cfg.get('k', [1, 2])
            _seq = self.km_cfg.get('sequence', 'CDR3')
            _d   = len([n for n in names if 'mer_' in n])
            print(f"  kmer            : {_d:>6,} d  "
                  f"(k={_k}  seq={_seq})")
            _total += _d

        if self.feat_cfg.get('biophysical'):
            _d = len([n for n in names
                      if n.startswith('cdr3_') or n.startswith('vh_')])
            print(f"  biophysical     : {_d:>6,} d  "
                  f"(CDR3 + VH properties)")
            _total += _d

        if self.feat_cfg.get('onehot'):
            _mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            if _mode == 'VHVL':
                _vh = self._oh_max_vh
                _vl = self._oh_max_vl
                _d  = (_vh + _vl) * _AA
                print(f"  onehot (VHVL)   : {_d:>6,} d  "
                      f"= 2 chains × pos × 20 AA")
                print(f"    VH  : {_vh:>3} pos × {_AA} AA = {_vh*_AA:,}")
                print(f"    VL  : {_vl:>3} pos × {_AA} AA = {_vl*_AA:,}")
            elif _mode == 'VH':
                _vh = self._oh_max_vh
                _d  = _vh * _AA
                print(f"  onehot (VH)     : {_d:>6,} d  "
                      f"= {_vh} pos × {_AA} AA")
            else:   # HCDR3
                _cdr3 = self._oh_max_cdr3
                _d    = _cdr3 * _AA
                print(f"  onehot (HCDR3)  : {_d:>6,} d  "
                      f"= {_cdr3} pos × {_AA} AA")
            _total += _d

        print(f"  {_SEP}")
        print(f"  TOTAL           : {len(names):>6,} d")
        print(f"  {_SEP}")
        return self

    # ── Transform ─────────────────────────────────────────────────────────────

    def transform(self, X_df: pd.DataFrame,
                  embeddings: np.ndarray = None) -> np.ndarray:
        parts = []
        n     = len(X_df)

        if self.feat_cfg.get('embedding') and embeddings is not None:
            parts.append(embeddings.astype(np.float32))

        if self.feat_cfg.get('kmer'):
            k_list = self.km_cfg.get('k', [1, 2])
            norm   = self.km_cfg.get('normalize', True)
            rows   = [list(compute_kmer_features(
                               self._get_kmer_seq(X_df.iloc[i]),
                               k_list, norm).values())
                      for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))

        if self.feat_cfg.get('biophysical'):
            fl   = self.bp_cfg.get('features', [])
            rows = [list(compute_biophysical_features(
                            self._get_cdr3(X_df.iloc[i]),
                            fl,
                            vh_seq=self._get_vh(X_df.iloc[i])).values())
                    for i in range(n)]
            parts.append(np.array(rows, dtype=np.float32))

        if self.feat_cfg.get('onehot'):
            # Warn if any predict sequence exceeds the training max length
            # (sequences will be silently truncated — could shift residue positions)
            mode = self.oh_cfg.get('sequence', 'HCDR3').upper()
            _warned = False
            for i in range(min(n, 5)):   # check first 5 rows only for speed
                row = X_df.iloc[i]
                if mode in ('VH', 'VHVL') and self._oh_max_vh > 0:
                    _l = len(str(row.get('HSEQ','') or ''))
                    if _l > self._oh_max_vh and not _warned:
                        print(f"  [onehot] WARNING: VH sequence length {_l} > "
                              f"training max {self._oh_max_vh} — truncated. "
                              f"Consider increasing onehot.max_lengths.vh in YAML.")
                        _warned = True
                if mode == 'VHVL' and self._oh_max_vl > 0:
                    _l = len(str(row.get('LSEQ','') or ''))
                    if _l > self._oh_max_vl and not _warned:
                        print(f"  [onehot] WARNING: VL sequence length {_l} > "
                              f"training max {self._oh_max_vl} — truncated. "
                              f"Consider increasing onehot.max_lengths.vl in YAML.")
                        _warned = True
                if mode == 'HCDR3' and self._oh_max_cdr3 > 0:
                    _l = len(str(row.get('CDR3','') or ''))
                    if _l > self._oh_max_cdr3 and not _warned:
                        print(f"  [onehot] WARNING: CDR3 length {_l} > "
                              f"training max {self._oh_max_cdr3} — truncated. "
                              f"Re-train with this data or increase max_lengths.hcdr3.")
                        _warned = True
            oh_rows = []
            for i in range(n):
                row  = X_df.iloc[i]
                vecs = [compute_onehot_features(seq, max_len)
                        for seq, max_len, _ in self._oh_segs(row)]
                oh_rows.append(np.concatenate(vecs))
            parts.append(np.array(oh_rows, dtype=np.float32))

        return np.hstack(parts) if parts else np.zeros((n, 0), dtype=np.float32)

    def fit_transform(self, X_df, embeddings=None):
        return self.fit(X_df, embeddings).transform(X_df, embeddings)

    # ── Properties ────────────────────────────────────────────────────────────

    @property
    def n_features(self):
        return len(self.feature_names_ or [])

    @property
    def non_embedding_feature_names(self):
        """kmer + biophysical + onehot names — interpretable, used for SHAP."""
        return [n for n in (self.feature_names_ or [])
                if not n.startswith('emb_')]

    @property
    def non_embedding_indices(self):
        return [i for i, n in enumerate(self.feature_names_ or [])
                if not n.startswith('emb_')]

    @property
    def onehot_feature_names(self):
        """One-hot feature names (oh_hcdr3_1_A, oh_vh_1_A, ...)."""
        return [n for n in (self.feature_names_ or [])
                if n.startswith('oh_')]

    @property
    def onehot_indices(self):
        return [i for i, n in enumerate(self.feature_names_ or [])
                if n.startswith('oh_')]


# ══════════════════════════════════════════════════════════════════════════════
# UTILITIES
# ══════════════════════════════════════════════════════════════════════════════

def _deep_merge(base: dict, override: dict) -> dict:
    """Recursively merge override into base dict (in-place). Returns base."""
    for k, v in override.items():
        if isinstance(v, dict) and isinstance(base.get(k), dict):
            _deep_merge(base[k], v)
        else:
            base[k] = v
    return base

def _setup_logger(log_path: str) -> logging.Logger:
    logger = logging.getLogger(f'RF_{os.path.basename(log_path)}')
    logger.setLevel(logging.DEBUG)
    if not logger.handlers:
        fh = logging.FileHandler(log_path, mode='a', encoding='utf-8')
        fh.setFormatter(logging.Formatter('%(message)s'))
        logger.addHandler(fh)
    return logger

def _log(logger, msg: str):
    print(msg)
    if logger: logger.info(msg)


# ══════════════════════════════════════════════════════════════════════════════
# 5.  MAIN CLASS
# ══════════════════════════════════════════════════════════════════════════════



# ══════════════════════════════════════════════════════════════════════════════
# DEFAULT CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════
_XGB_DEFAULT_CONFIG = {
    'mode': 'manual',
    'task': 'classification',
    'model': {
        'n_estimators':       3000,
        'max_depth':             6,
        'learning_rate':      0.01,
        'subsample':           0.8,
        'colsample_bytree':    0.8,  # analogous to max_features in RF
        'gamma':               0.0,
        'min_child_weight':      1,
        'reg_alpha':           0.0,
        'reg_lambda':          1.0,
        'tree_method':       'hist',
        'n_jobs':               -1,
        'random_state':         42,
    },
    'training': {
        'use_random_search':  False,  # True = RandomizedSearchCV
        'n_iter':               30,
        'cv':                    5,
        'early_stopping_rounds': 0,   # 0 = disabled (same as RF, use kfold for selection)
                                       # set to 50 when using --split 0.8
        'eval_metric':        'auc',
        'class_weight':   'balanced', # auto scale_pos_weight
    },
    'features': {
        'embedding':   False,
        'kmer':        False,
        'biophysical': False,
        'onehot':      False,
    },
    'kmer': {
        'k':        [1, 2],
        'sequence': 'CDR3',
        'normalize': True,
    },
    'onehot': {
        'sequence':    'HCDR3',
        'max_lengths': {'vh': 0, 'vl': 0, 'hcdr3': 0},
    },
    # LM-specific colsample_bytree overrides (mirrors max_features in RF)
    'lm_profiles': {
        'biophysical':     {'colsample_bytree': 1.0},
        'kmer':            {'colsample_bytree': 0.5},
        'onehot':          {'colsample_bytree': 0.3, 'max_depth': 8},
        'onehot_vh':       {'colsample_bytree': 0.3, 'max_depth': 8},
        'onehot_hcdr3':    {'colsample_bytree': 0.5},
        'onehot_cdr3':     {'colsample_bytree': 0.5},
        'ablang':          {'colsample_bytree': 0.25},
        'antiberty':       {'colsample_bytree': 0.25},
        'antiberta2':      {'colsample_bytree': 0.20},
        'antiberta2-cssp': {'colsample_bytree': 0.20},
        'igbert':          {'colsample_bytree': 0.20},
    },
    'shap': {
        'enabled':               True,
        'top_features':            30,
        'plot_types': ['bar', 'beeswarm', 'heatmap'],
        'max_waterfall_samples':   50,
        'waterfall_top_n':         20,
        'waterfall_fig_size':   [8, 6],
        'pub_dpi':               300,
        'waterfall_fmt':        'tiff',
        'make_ppt':              True,
    },
    'mutagenesis': {
        'max_samples':  50,
        'format':     'tiff',
        'pub_dpi':     300,
        'make_ppt':   True,
    },
}


# ══════════════════════════════════════════════════════════════════════════════
# XGBoostModel
# ══════════════════════════════════════════════════════════════════════════════

class XGBoostModel:
    """
    XGBoost classifier/regressor with full feature + analysis parity vs RF.
    Identical external interface — plug-in replacement in predict_developability.py.
    """

    # ── static helpers ─────────────────────────────────────────────────────


    @staticmethod
    def auto_detect_config(n: int, pos_rate: float,
                           task: str = 'classification') -> dict:
        """
        Derive RF config from dataset size, class balance and task.

        Size tiers:
          xs  n <  5k  →  500 trees, depth 10, min_leaf 3,  max_feat sqrt
          sm  5–20k    → 1000 trees, depth 15, min_leaf 5,  max_feat 0.3
          md  20–80k   → 2000 trees, depth 20, min_leaf 8,  max_feat 0.2
          lg  80–200k  → 2000 trees, depth 25, min_leaf 10, max_feat 0.15
          xl  >200k    → 3000 trees, depth 30, min_leaf 15, max_feat 0.1
        """
        if   n <  5_000: tier='xs'; n_e,d,ml,mf = 500, 10, 3,  'sqrt'
        elif n < 20_000: tier='sm'; n_e,d,ml,mf = 1000,15, 5,  0.3
        elif n < 80_000: tier='md'; n_e,d,ml,mf = 2000,20, 8,  0.2
        elif n <200_000: tier='lg'; n_e,d,ml,mf = 2000,25,10,  0.15
        else:            tier='xl'; n_e,d,ml,mf = 3000,30,15,  0.1

        criterion = 'squared_error' if task=='regression' else 'entropy'
        if task == 'regression':
            cw = None
        else:
            min_rate = min(pos_rate, 1-pos_rate)
            if   min_rate >= 0.40: cw = 'balanced'
            elif min_rate >= 0.10: cw = 'balanced_subsample'
            else:
                ratio = (1-pos_rate)/(pos_rate+1e-8)
                cw    = {0:1, 1:ratio}
                print(f"  [WARN] Severe imbalance → explicit weights {cw}")

        cfg = copy.deepcopy(_XGB_DEFAULT_CONFIG)
        cfg['mode'] = 'auto'; cfg['task'] = task
        cfg['model'].update({'n_estimators':n_e,'max_depth':d,
                             'learning_rate': 0.01,
                             'colsample_bytree': mf if isinstance(mf, float) else 0.8})
        cfg['training']['class_weight'] = cw
        cfg['_auto'] = {'n':n,'pos_rate':round(pos_rate,4),'size_tier':tier,'task':task}
        return cfg

    @staticmethod
    def print_config_report(cfg: dict) -> None:
        a    = cfg.get('_auto', {});  m  = cfg.get('model', {})
        t    = cfg.get('training', {}); ft = cfg.get('features', {})
        km   = cfg.get('kmer', {});   bp  = cfg.get('biophysical', {})
        sh   = cfg.get('shap', {});   task = cfg.get('task', 'classification')
        W = 62; sep = '═' * W; sep2 = '─' * W
        print(f"\n{sep}")
        if a:
            print(f"  XGBoost  ·  AUTO  ({task})  tier={a.get('size_tier')}  n={a.get('n'):,}")
        else:
            print(f"  XGBoost  ·  MANUAL  ({task})")
        print(sep2)
        print(f"  MODEL          n_est={m.get('n_estimators',3000)}"
              f"  depth={m.get('max_depth',6)}"
              f"  lr={m.get('learning_rate',0.01):.3f}"
              f"  subsample={m.get('subsample',0.8)}"
              f"  colsample={m.get('colsample_bytree',0.8)}")
        print(f"                 gamma={m.get('gamma',0.0)}"
              f"  min_child_w={m.get('min_child_weight',1)}"
              f"  reg_α={m.get('reg_alpha',0.0)}"
              f"  reg_λ={m.get('reg_lambda',1.0)}")
        print(sep2)
        print(f"  FEATURES       embedding={ft.get('embedding')}  "
              f"kmer={ft.get('kmer')}  biophysical={ft.get('biophysical')}  "
              f"onehot={ft.get('onehot', False)}")
        if ft.get('kmer'):
            print(f"                 k-mer k={km.get('k')}  seq={km.get('sequence')}")
        if ft.get('biophysical'):
            print(f"                 bio features={len(bp.get('features', []))}  "
                  f"seq={bp.get('sequence')}")
        print(sep2)
        _search = t.get('use_random_search', False)
        _esr    = t.get('early_stopping_rounds', 0)
        if task == 'regression':
            print(f"  TRAINING       search={_search}  early_stop={_esr}")
        else:
            print(f"  TRAINING       class_weight={t.get('class_weight')}  "
                  f"search={_search}  early_stop={_esr}")
        print(f"  SHAP           enabled={sh.get('enabled')}  "
              f"top={sh.get('top_features')}  plots={sh.get('plot_types')}")
        print(f"{sep}\n")

    # ── init ──────────────────────────────────────────────────────────────────


    def __init__(self, config: dict = None,
                 config_path: str = "config/xgboost.yaml",
                 verbose: bool = True):
        import yaml
        self.config = copy.deepcopy(_XGB_DEFAULT_CONFIG)

        # Load YAML if exists
        if config_path and os.path.exists(config_path):
            with open(config_path) as f:
                user_cfg = yaml.safe_load(f) or {}
            _deep_merge(self.config, user_cfg)
            if verbose:
                print(f"[XGBoostModel] config ← {config_path}")

        # Override with programmatic config
        if config is not None:
            _deep_merge(self.config, config)

        self.model = None
        self.fb_   = None
        self.task  = self.config.get('task', 'classification')

        if verbose:
            self.print_config_report(self.config)

    # ── config helpers ─────────────────────────────────────────────────────

    def apply_lm_profile(self, lm: str, logger=None) -> None:
        """
        Apply LM-specific hyperparameter overrides from lm_profiles in YAML.
        Called before _build_sklearn_model() with the --lm name.

        Priority: lm_profiles.{lm} > model defaults
        Only keys present in the profile are overridden.
        """
        profiles = self.config.get('lm_profiles', {})
        profile  = profiles.get(lm, {})
        if not profile:
            return
        m = self.config['model']
        for k, v in profile.items():
            m[k] = v
        _log(logger, f"[XGB] lm_profile '{lm}' applied: "
                     + "  ".join(f"{k}={v}" for k, v in profile.items()))


    def _auto_fix_features(self, embeddings, X_df=None):
        """Auto-enable embedding features when embeddings provided."""
        ft = self.config.setdefault('features', {})
        if embeddings is not None and not ft.get('embedding', False):
            ft['embedding'] = True
        if embeddings is None and ft.get('embedding', True):
            if not any(ft.get(k) for k in ('kmer', 'biophysical', 'onehot')):
                ft['biophysical'] = True

    # ── build XGB model ────────────────────────────────────────────────────

    def _build_xgb(self, y=None):
        """Build XGBClassifier or XGBRegressor from config."""
        m = self.config['model']
        t = self.config['training']

        # Scale pos weight for class imbalance (replaces class_weight='balanced')
        # scale_pos_weight: classification only
        # For regression: spw=1.0 (not applicable)
        _y_for_spw = np.asarray(y, dtype=float) if y is not None else None
        _is_bin_spw = (_y_for_spw is not None and
                       len(set(_y_for_spw.tolist())) <= 2 and
                       set(_y_for_spw.tolist()).issubset({0, 1, 0.0, 1.0}))
        if (y is not None and self.task == 'classification'
                and _is_bin_spw
                and t.get('class_weight') == 'balanced'):
            counts = np.bincount(np.asarray(y, dtype=int))
            spw    = (counts[0] / counts[1]
                      if len(counts) > 1 and counts[1] > 0 else 1.0)
        else:
            spw = 1.0

        # NOTE: early_stopping_rounds is NOT set in constructor —
        # it is passed to fit() only when eval_set is available.
        # This avoids crashes when RandomizedSearchCV calls fit() without eval_set.
        common = dict(
            tree_method      = m.get('tree_method', 'hist'),
            n_estimators     = m.get('n_estimators', 3000),
            max_depth        = m.get('max_depth', 6),
            learning_rate    = m.get('learning_rate', 0.01),
            subsample        = m.get('subsample', 0.8),
            colsample_bytree = m.get('colsample_bytree', 0.8),
            gamma            = m.get('gamma', 0.0),
            min_child_weight = m.get('min_child_weight', 1),
            reg_alpha        = m.get('reg_alpha', 0.0),
            reg_lambda       = m.get('reg_lambda', 1.0),
            n_jobs           = m.get('n_jobs', -1),
            random_state     = m.get('random_state', 42),
            verbosity        = 0,
        )
        # Auto-detect task from y if config says 'auto' or if y provided
        _task = self.task
        if y is not None:
            _yf = np.asarray(y, dtype=float)
            _is_bin = (len(set(_yf.tolist())) <= 2 and
                       set(_yf.tolist()).issubset({0, 1, 0.0, 1.0}))
            if _task == 'auto':
                _task = 'classification' if _is_bin else 'regression'
            elif _task == 'classification' and not _is_bin:
                _task = 'regression'   # override: y is continuous
        if _task == 'classification':
            return xgb.XGBClassifier(
                objective        = 'binary:logistic',
                eval_metric      = t.get('eval_metric', 'auc'),
                scale_pos_weight = spw,
                **common)
        else:
            return xgb.XGBRegressor(
                objective = 'reg:squarederror',
                **common)

    # ── train ──────────────────────────────────────────────────────────────

    def train(self, X_df: pd.DataFrame, y,
              embeddings:   np.ndarray = None,
              target:       str        = "model",
              target_col:   str        = "",
              embedding_lm: str        = "",
              val_X:        pd.DataFrame = None,
              val_y                       = None,
              val_embeddings: np.ndarray  = None,
              logger                      = None):
        """
        Train XGBoost with FeatureBuilder pipeline.
        Same interface as RandomForestModel.train().
        """
        # Auto-detect regression if config says classification but y is continuous
        _y_tmp    = np.asarray(y, dtype=float)
        _unique_y = len(set(_y_tmp.tolist()))
        _is_bin   = (_unique_y <= 2 and set(_y_tmp.tolist()).issubset({0, 1, 0.0, 1.0}))
        if self.task == 'classification' and not _is_bin:
            self.task = 'regression'
            _log(logger, f"[train] Auto-detected task=regression "
                         f"(y has {_unique_y} unique values, not binary)")
        y_arr = np.asarray(y, dtype=float if self.task == 'regression' else int)

        self._auto_fix_features(embeddings, X_df)
        self.fb_ = FeatureBuilder(self.config)
        self.fb_.fit(X_df, embeddings)
        self.apply_lm_profile(embedding_lm, logger=logger)

        X_feat = self.fb_.transform(X_df, embeddings)
        _log(logger, f"[train] n={len(y_arr):,}  features={X_feat.shape[1]}"
                     f"  task={self.task}  lm={embedding_lm}")

        t   = self.config['training']
        esr = t.get('early_stopping_rounds', 0)

        if t.get('use_random_search', False):
            _log(logger, "[train] RandomizedSearchCV ...")
            param_dist = {
                # Original IPI distributions — validated on ipi_psr_trainset
                'n_estimators':     randint(3000, 6000),
                'max_depth':        randint(4,  12),
                'learning_rate':    uniform(0.005, 0.090),  # 0.005–0.095
                'subsample':        uniform(0.6,  0.4),     # 0.6–1.0
                'colsample_bytree': uniform(0.6,  0.4),     # 0.6–1.0
                'gamma':            uniform(0,    0.5),
                'min_child_weight': randint(1,  10),
                'reg_alpha':        uniform(0,    1.0),
                'reg_lambda':       uniform(0.1,  1.9),     # 0.1–2.0
            }
            # Build base model WITHOUT early_stopping_rounds:
            # RandomizedSearchCV uses its own CV and cannot pass eval_set per fold
            _base = self._build_xgb(y_arr)
            try:
                _base.set_params(early_stopping_rounds=None)
            except Exception:
                pass   # older XGBoost may not support None — ignore
            _cv_scoring = ('r2' if self.task == 'regression'
                          else 'roc_auc')
            search = RandomizedSearchCV(
                _base, param_dist,
                n_iter      = t.get('n_iter', 30),
                cv          = t.get('cv', 5),
                scoring     = _cv_scoring,
                n_jobs      = -1,
                random_state= 42, verbose=0)
            search.fit(X_feat, y_arr)
            self.model = search.best_estimator_
            _metric_name = 'R²' if self.task == 'regression' else 'AUC'
            _log(logger, f"[train] best CV {_metric_name} = {search.best_score_:.4f}")
            _bp = search.best_params_
            _log(logger, "[best_params] " + "  ".join(f"{k}={v}" for k, v in sorted(_bp.items())))
            # Store best params into config so save() persists them
            self.config['model'].update(_bp)
            _log(logger, "[best_params] Stored in config — copy to xgboost.yaml to reuse")

        elif esr > 0 and val_X is not None and val_y is not None:
            # Early stopping with explicit val set — esr passed to fit()
            X_va = self.fb_.transform(val_X, val_embeddings)
            y_va = np.asarray(val_y, dtype=float if self.task == 'regression' else int)
            self.model = self._build_xgb(y_arr)
            self.model.fit(X_feat, y_arr,
                           eval_set              = [(X_va, y_va)],
                           early_stopping_rounds = esr,
                           verbose               = False)
            _log(logger, f"[train] best iteration: {self.model.best_iteration}")

        elif esr > 0:
            # Early stopping with internal split — esr passed to fit()
            from sklearn.model_selection import train_test_split
            X_tr, X_va, y_tr, y_va = train_test_split(
                X_feat, y_arr, test_size=0.1,
                stratify=(y_arr if self.task == 'classification' else None),
                random_state=42)
            self.model = self._build_xgb(y_arr)
            self.model.fit(X_tr, y_tr,
                           eval_set              = [(X_va, y_va)],
                           early_stopping_rounds = esr,
                           verbose               = False)
            _log(logger, f"[train] best iteration: {self.model.best_iteration}")

        else:
            self.model = self._build_xgb(y_arr)
            self.model.fit(X_feat, y_arr)

        # Log final hyperparameters used
        try:
            _p = self.model.get_params()
            def _fmt(v, fmt='.4f'):
                try:    return format(float(v), fmt)
                except: return str(v)
            _log(logger, f"[params] n_est={_p.get('n_estimators')}"
                         f"  depth={_p.get('max_depth')}"
                         f"  lr={_fmt(_p.get('learning_rate'))}"
                         f"  subsample={_fmt(_p.get('subsample'))}"
                         f"  colsample={_fmt(_p.get('colsample_bytree'))}"
                         f"  gamma={_fmt(_p.get('gamma'))}"
                         f"  min_child_w={_p.get('min_child_weight')}"
                         f"  reg_α={_fmt(_p.get('reg_alpha'))}"
                         f"  reg_λ={_fmt(_p.get('reg_lambda'))}"
                         f"  scale_pos_w={_fmt(_p.get('scale_pos_weight', 1.0), '.2f')}")
            _bi = getattr(self.model, 'best_iteration', None)
            if _bi:
                _log(logger, f"[params] early_stop → best_iteration={_bi}")
        except Exception as _pe:
            _log(logger, f"[params] WARNING: could not log params — {_pe}")

        # ── Validation metrics ─────────────────────────────────────────────
        if val_X is not None and val_y is not None:
            X_va  = self.fb_.transform(val_X, val_embeddings)
            y_va  = np.asarray(val_y, dtype=float if self.task == 'regression' else int)
            if self.task == 'classification':
                from sklearn.metrics import roc_auc_score, accuracy_score, f1_score, recall_score
                val_probs = self.model.predict_proba(X_va)[:, 1]
                val_preds = (val_probs >= 0.5).astype(int)
                val_auc   = roc_auc_score(y_va, val_probs)
                val_acc   = accuracy_score(y_va, val_preds)
                val_f1    = f1_score(y_va, val_preds, zero_division=0)
                val_rcf   = recall_score(y_va, val_preds, pos_label=0, zero_division=0)
                _log(logger, f"  val_auc={val_auc:.4f}  val_acc={val_acc:.4f}"
                             f"  val_f1={val_f1:.4f}  val_rec_fail={val_rcf:.4f}")
            else:
                from sklearn.metrics import r2_score, mean_absolute_error
                from scipy.stats import pearsonr, spearmanr
                val_preds_r = self.model.predict(X_va)
                _log(logger, f"  val_r2={r2_score(y_va, val_preds_r):.4f}"
                             f"  val_pearson={pearsonr(y_va, val_preds_r)[0]:.4f}"
                             f"  val_spearman={spearmanr(y_va, val_preds_r)[0]:.4f}")

        _log(logger, "[train] completed.")

        # ── Auto SHAP ─────────────────────────────────────────────────────
        if self.config.get('shap', {}).get('enabled', False) and _SHAP_AVAILABLE:
            sh_top  = self.config.get('shap', {}).get('top_features', 30)
            _parts  = str(target).split('_xgboost_')
            _shap_lm = embedding_lm
            _shap_db = _parts[1] if len(_parts) > 1 else target
            _log(logger, f"\n[SHAP] Running on TRAIN set (top {sh_top}) ...")
            try:
                self.shap_analysis(X_df, embeddings,
                                   output_prefix   = target,
                                   split_tag       = "train",
                                   top_n           = sh_top,
                                   barcodes        = list(X_df.index.astype(str)),
                                   actual_labels   = list(y_arr),
                                   actual_col_name = target_col or target,
                                   lm_name         = _shap_lm,
                                   db_name         = _shap_db,
                                   logger          = logger)
            except Exception as _se:
                import traceback
                _log(logger, f"[SHAP] train SHAP failed — {_se}\n{traceback.format_exc()}")

            if val_X is not None and val_y is not None:
                _log(logger, f"\n[SHAP] Running on VAL set (top {sh_top}) ...")
                try:
                    self.shap_analysis(val_X, val_embeddings,
                                       output_prefix   = target,
                                       split_tag       = "val",
                                       top_n           = sh_top,
                                       barcodes        = list(val_X.index.astype(str)),
                                       actual_labels   = list(val_y),
                                       actual_col_name = target_col or target,
                                       lm_name         = _shap_lm,
                                       db_name         = _shap_db,
                                       logger          = logger)
                except Exception as _se:
                    import traceback
                    _log(logger, f"[SHAP] val SHAP failed — {_se}\n{traceback.format_exc()}")
        elif not _SHAP_AVAILABLE:
            _log(logger, "[SHAP] skipped — pip install shap to enable")

        return self

    # ── predict ────────────────────────────────────────────────────────────

    def predict_proba(self, X_df, embeddings=None) -> np.ndarray:
        if self.task == 'regression':
            raise RuntimeError("predict_proba not available for regression. Use predict().")
        return self.model.predict_proba(
            self.fb_.transform(X_df, embeddings))[:, 1]

    def predict(self, X_df, embeddings=None, threshold: float = None) -> np.ndarray:
        X_feat = self.fb_.transform(X_df, embeddings)
        if self.task == 'regression':
            return self.model.predict(X_feat)
        t = threshold if threshold is not None else getattr(self, 'recommended_threshold', 0.5)
        return (self.model.predict_proba(X_feat)[:, 1] >= t).astype(int)

    def predict_raw(self, X_df, embeddings=None) -> np.ndarray:
        """Raw XGB score (regression) or proba (classification)."""
        X_feat = self.fb_.transform(X_df, embeddings)
        if self.task == 'regression':
            return self.model.predict(X_feat)
        return self.model.predict_proba(X_feat)[:, 1]

    # ── save / load ────────────────────────────────────────────────────────

    def save(self, path: str):
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        rec_thresh = getattr(self, 'recommended_threshold', None)
        joblib.dump({
            'model':                 self.model,
            'fb_':                   self.fb_,
            'config':                self.config,
            'task':                  self.task,
            'recommended_threshold': rec_thresh,
        }, path)
        thresh_note = (f"  threshold={rec_thresh:.4f}"
                       if rec_thresh is not None else "  threshold=None (0.5 at predict)")
        print(f"[XGBoostModel] saved → {path}{thresh_note}")

    @classmethod
    def load(cls, path, config_path: str = "config/xgboost.yaml"):
        instance = cls(config_path=config_path, verbose=False)
        payload  = joblib.load(path)
        if isinstance(payload, dict) and 'model' in payload:
            instance.model  = payload['model']
            instance.fb_    = payload.get('fb_', None)
            _deep_merge(instance.config, payload.get('config', {}))
            instance.task   = payload.get('task', 'classification')
            rt              = payload.get('recommended_threshold', None)
        else:
            instance.model = payload
            rt = None
        instance.recommended_threshold = float(rt) if rt is not None else 0.5
        flag = "  (embedded by kfold)" if rt and rt != 0.5 else ""
        print(f"[XGBoostModel] loaded ← {path}")
        print(f"  recommended_threshold={instance.recommended_threshold:.4f}{flag}")
        return instance


    # ── CDR3 mutagenesis ─────────────────────────────────────────────────

    def cdr3_mutagenesis_heatmap(
        self,
        X_df: pd.DataFrame,
        embeddings: np.ndarray = None,
        output_dir: str = None,
        pub_dpi: int = 300,
        fig_size: tuple = (10, 4),
        logger = None,
    ) -> dict:
        """
        In-silico CDR3 single-point mutagenesis for every antibody in X_df.
        For each position × amino acid substitution, predicts score/probability
        using the trained RF model and generates a heatmap per antibody.

        Saves:
          {output_dir}/{barcode}_cdr3_mutagenesis.tiff   (Nature 300 DPI)
          {output_dir}/cdr3_mutagenesis_all.pptx         (one slide per antibody)

        Returns dict: {barcode: heatmap_matrix (20 × cdr3_len)}
        """
        if self.model is None or self.fb_ is None:
            _log(logger, "[MUTA] model not trained/loaded"); return {}

        out_dir = output_dir or MODEL_DIR
        os.makedirs(out_dir, exist_ok=True)
        results = {}
        saved_paths = []

        # Get prediction scores for WT sequences first
        for i in range(len(X_df)):
            row = X_df.iloc[i]
            bc  = str(X_df.index[i]) if hasattr(X_df, 'index') else str(i)
            cdr3_seq = str(row.get('CDR3', '') or '').upper().replace('-', '')
            if not cdr3_seq:
                _log(logger, f"  [MUTA] {bc}: empty CDR3 — skipped")
                continue

            n_pos  = len(cdr3_seq)
            n_aa   = len(AMINO_ACIDS)
            matrix = np.zeros((n_aa, n_pos), dtype=np.float32)

            # Build mutation DataFrame for all substitutions at once (vectorised)
            mut_rows = []
            mut_info = []   # (aa_idx, pos_idx)
            for pos in range(n_pos):
                for aa_idx, aa in enumerate(AMINO_ACIDS):
                    mutant_cdr3 = cdr3_seq[:pos] + aa + cdr3_seq[pos+1:]
                    mut_row = dict(row)
                    mut_row['CDR3'] = mutant_cdr3
                    # Also mutate HSEQ if available (replace CDR3 region)
                    if 'HSEQ' in mut_row and mut_row['HSEQ']:
                        hseq = str(mut_row['HSEQ'])
                        cdr3_start = hseq.find(cdr3_seq)
                        if cdr3_start >= 0:
                            mut_row['HSEQ'] = (hseq[:cdr3_start] + mutant_cdr3 +
                                               hseq[cdr3_start + n_pos:])
                    mut_rows.append(mut_row)
                    mut_info.append((aa_idx, pos))

            # Predict all mutants in one batch
            mut_df = pd.DataFrame(mut_rows)
            mut_df.index = [f"{bc}_mut_{j}" for j in range(len(mut_df))]

            # Build embeddings for mutants if needed
            mut_emb = None
            if embeddings is not None and self.feat_cfg.get('embedding', False):
                _log(logger,
                     f"  [MUTA] {bc}: embedding mode — mutant embeddings not available.\n"
                     "         Mutagenesis requires kmer/biophysical/onehot features.\n"
                     "         Set features.embedding=false in YAML.")
                continue

            try:
                X_mut = self.fb_.transform(mut_df, mut_emb)
                if self.task == 'classification':
                    scores = self.model.predict_proba(X_mut)[:, 1]
                else:
                    scores = self.model.predict(X_mut)
            except Exception as _e:
                _log(logger, f"  [MUTA] {bc}: prediction failed — {_e}"); continue

            for j, (aa_idx, pos_idx) in enumerate(mut_info):
                matrix[aa_idx, pos_idx] = scores[j]

            results[bc] = matrix

            # ── WT score (diagonal reference line) ────────────────────────────
            wt_row = pd.DataFrame([dict(row)])
            wt_row.index = [bc]
            try:
                X_wt = self.fb_.transform(wt_row, None)
                wt_score = (self.model.predict_proba(X_wt)[:, 1][0]
                            if self.task == 'classification'
                            else self.model.predict(X_wt)[0])
            except Exception:
                wt_score = None

            # ── Plot heatmap ─────────────────────────────────────────────────
            fig, ax = plt.subplots(figsize=(max(fig_size[0], n_pos * 0.45), fig_size[1]))

            score_label = "P(PASS)" if self.task == 'classification' else "Predicted score"
            cmap = 'RdBu'
            import matplotlib.colors as _mc
            if self.task == 'classification':
                _norm = _mc.TwoSlopeNorm(vmin=0.0, vcenter=0.5, vmax=1.0)
            else:
                _vmin = float(matrix.min()); _vmax = float(matrix.max())
                _norm = _mc.TwoSlopeNorm(vmin=_vmin, vcenter=(_vmin+_vmax)/2, vmax=_vmax)

            im = ax.imshow(matrix, aspect='auto', cmap=cmap, norm=_norm)

            # ── Cell annotations ─────────────────────────────────────────────
            _cell_fsz = max(4.0, min(7.0, 120.0 / max(n_pos, 1)))
            for _ai in range(n_aa):
                for _pi in range(n_pos):
                    _v = matrix[_ai, _pi]
                    if np.isnan(_v): continue
                    _tc = 'white' if (_v < 0.35 or _v > 0.65) else '#333333'
                    ax.text(_pi, _ai, f"{_v:.2f}", ha='center', va='center',
                            fontsize=_cell_fsz, color=_tc,
                            fontweight='bold' if abs(_v - 0.5) > 0.3 else 'normal')

            plt.colorbar(im, ax=ax, label=score_label, fraction=0.03, pad=0.02)

            ax.set_xticks(range(n_pos))
            ax.set_xticklabels([f"{cdr3_seq[p]}\n{p+1}" for p in range(n_pos)],
                               fontsize=8)
            ax.set_yticks(range(n_aa))
            ax.set_yticklabels(list(AMINO_ACIDS), fontsize=7)
            ax.set_xlabel('CDR3 position  (WT residue shown above position number)',
                          fontsize=8)
            ax.set_ylabel('Substituted AA', fontsize=8)

            wt_str = f"  WT {score_label}={wt_score:.4f}" if wt_score is not None else ""
            ax.set_title(
                f"IPI MLAbDev · CDR3 Mutagenesis Heatmap\n"
                f"ID: {bc}{wt_str}\n"
                f"RF | {self.task.upper()} | CDR3={cdr3_seq}",
                fontsize=8, loc='center', pad=6
            )

            # Mark WT residues on diagonal
            for pos in range(n_pos):
                wt_aa_idx = AMINO_ACIDS.find(cdr3_seq[pos])
                if wt_aa_idx >= 0:
                    ax.add_patch(plt.Rectangle(
                        (pos - 0.5, wt_aa_idx - 0.5), 1, 1,
                        fill=False, edgecolor='black', lw=1.5, zorder=5
                    ))

            plt.tight_layout()
            bc_safe  = bc.replace('/', '_').replace(' ', '_')
            img_path = os.path.join(out_dir, f"{bc_safe}_cdr3_mutagenesis.tiff")
            plt.savefig(img_path, dpi=pub_dpi, format='tiff', bbox_inches='tight')
            plt.close()
            saved_paths.append((bc, img_path, wt_score))
            _log(logger, f"  [MUTA] {bc} → {os.path.basename(img_path)}")

        # ── Assemble PPT ──────────────────────────────────────────────────────
        if saved_paths:
            try:
                from pptx import Presentation as _Prs
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor

                prs = _Prs()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                for bc, img_path, wt_score in saved_paths:
                    slide = prs.slides.add_slide(blank)
                    img_w = Inches(fig_size[0])
                    img_h = Inches(fig_size[1])
                    left  = (prs.slide_width - img_w) / 2
                    slide.shapes.add_picture(img_path, left, Inches(0.6),
                                             width=img_w, height=img_h)
                    txb = slide.shapes.add_textbox(
                        Inches(0.15), Inches(7.1), Inches(13.0), Inches(0.35))
                    tf  = txb.text_frame
                    wt_str = f"  WT score={wt_score:.4f}" if wt_score is not None else ""
                    tf.text = f"CDR3 Mutagenesis  |  {bc}{wt_str}"
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    p.runs[0].font.size = Pt(7)
                    p.runs[0].font.color.rgb = RGBColor(0x88, 0x87, 0x80)

                ppt_path = os.path.join(out_dir, "cdr3_mutagenesis_all.pptx")
                prs.save(ppt_path)
                _log(logger, f"[MUTA] PPT ({len(saved_paths)} slides) → {ppt_path}")
            except ImportError:
                _log(logger, "[MUTA] python-pptx not installed — PPT skipped")
            except Exception as _pe:
                _log(logger, f"[MUTA] PPT failed — {_pe}")

        _log(logger, f"[MUTA] complete  ({len(results)} antibodies)")
        return results


    # ── SHAP analysis ──────────────────────────────────────────────────

    def shap_analysis(self, X_df, embeddings=None,
                      output_prefix="rf_shap",
                      split_tag="train",
                      top_n=50,
                      barcodes=None,
                      actual_labels=None,
                      actual_col_name="label",
                      lm_name="",
                      db_name="",
                      logger=None):
        """
        SHAP TreeExplainer — works for ALL feature modes.
        Generates beeswarm/bar/heatmap for all samples, plus a multi-page
        PDF waterfall (one page per antibody) with full title info.

        Parameters
        ----------
        X_df          : DataFrame with HSEQ/CDR3 columns (index = BARCODE)
        embeddings    : (n, emb_dim) PLM array or None
        output_prefix : filename stem
        split_tag     : 'train' | 'val' | 'predict'
        top_n         : top features shown per plot
        barcodes      : list of antibody IDs — if None, inferred from X_df.index
        actual_labels : list/array of ground-truth labels — shown in title if provided
        actual_col_name: target column name (e.g. 'psr_filter') — shown in title
        lm_name       : LM name for title (e.g. 'ablang', 'biophysical')
        db_name       : database stem for title (e.g. 'ipi_psr_trainset_train')

        YAML shap config keys used here:
          top_features          : features per waterfall panel  (default 30)
          waterfall_top_features: override top_n for waterfall only (default = top_features)
          max_waterfall_samples : max antibodies in PDF          (default 50)
          waterfall_per_page    : 1 | 2 | 4 panels per page     (default 1)
          plot_dpi              : DPI for PNG outputs            (default 150)
          pub_dpi               : DPI for publication PDF        (default 300)
        """
        if not _SHAP_AVAILABLE:
            _log(logger, "[SHAP] shap not installed — pip install shap"); return
        if self.model is None or self.fb_ is None:
            _log(logger, "[SHAP] model not trained"); return

        # Infer barcodes from X_df index or provided list
        if barcodes is None:
            barcodes = list(X_df.index.astype(str)) if hasattr(X_df, 'index') else [str(i) for i in range(len(X_df))]

        sh_cfg     = self.config.get('shap', {})
        max_s      = sh_cfg.get('max_samples', 500)
        plots      = sh_cfg.get('plot_types', ['bar', 'beeswarm'])
        top_n      = sh_cfg.get('top_features', top_n)        # YAML overrides arg
        wf_top_n   = sh_cfg.get('waterfall_top_features', top_n)  # waterfall-specific
        max_wf     = sh_cfg.get('max_waterfall_samples', 50)  # max antibodies in PDF
        wf_per_pg  = sh_cfg.get('waterfall_per_page', 1)      # 1 | 2 | 4 per page
        plot_dpi   = sh_cfg.get('plot_dpi', 150)               # PNG dpi
        pub_dpi    = sh_cfg.get('pub_dpi',  300)               # PDF dpi (publication)
        # beeswarm is always generated regardless of plots list
        if 'beeswarm' not in plots:
            plots = list(plots) + ['beeswarm']

        # Build feature matrix and names
        X_feat   = self.fb_.transform(X_df, embeddings)
        n_feat   = X_feat.shape[1]
        ne_names = self.fb_.non_embedding_feature_names   # kmer + biophysical
        ne_idx   = self.fb_.non_embedding_indices

        # ── Embedding-only mode: SHAP is not meaningful ───────────────────────
        # PLM dimensions (emb_0 … emb_N) are entangled latent features —
        # no single dimension maps to a biological property.
        # Use Integrated Gradients from transformer_lm for sequence-level
        # attribution on PLM embeddings.
        if not ne_idx:
            _log(logger,
                 "[SHAP] Skipped — embedding-only mode.\n"
                 "  PLM dimensions (emb_0..emb_N) are latent and not interpretable.\n"
                 "  For sequence-level attribution use transformer_lm IG instead:\n"
                 "    model.global_ig_analysis(dataset)  in transformer_lm.py")
            return

        # Slice to interpretable features only (kmer + biophysical).
        # Embedding dims are excluded — they have no biological name.
        # SHAP runs only on ne_idx columns so no empty labels appear.
        ne_idx_arr = np.array(ne_idx)
        X_ne   = X_feat[:, ne_idx_arr]          # (n, n_interpretable)
        all_names = ne_names                     # names align 1:1 with X_ne columns

        # Subsample for speed
        # Keep full (unsubsampled) matrix for predict_proba — correct scores
        if len(X_ne) > max_s:
            idx_s   = np.random.choice(len(X_ne), max_s, replace=False)
            X_shap  = X_ne[idx_s]
            # For waterfall predictions: use same subsample but from full X_feat
            _X_full_wf = X_feat[idx_s]
        else:
            X_shap     = X_ne
            _X_full_wf = X_feat

        _log(logger, f"\n[SHAP] {split_tag.upper()} — TreeExplainer  "
                     f"n={len(X_shap)}  interpretable features={len(ne_idx)}  "
                     f"(embedding dims excluded)  top={top_n}")

        explainer = _shap_lib.TreeExplainer(self.model)
        try:
            shap_values = explainer.shap_values(X_shap, check_additivity=True)
        except Exception:
            shap_values = explainer.shap_values(X_shap, check_additivity=False)

        # Robustly extract (n_samples, n_features) array for class-1 or regression.
        # SHAP returns different shapes across versions:
        #   regression           : ndarray (n, f)
        #   binary classification:
        #     shap < 0.41        : list[ (n,f), (n,f) ]           → take [1]
        #     shap >= 0.42 binary: ndarray (n, f)                  → direct
        #     shap >= 0.42 multi : ndarray (n, f, n_classes)       → take [:,:,1]
        #     some versions      : ndarray (n_classes, n, f)       → take [1]
        if isinstance(shap_values, list):
            # list of arrays — take class 1
            sv = np.array(shap_values[1], dtype=np.float64)
        elif isinstance(shap_values, np.ndarray):
            if shap_values.ndim == 3:
                # could be (n, f, classes) or (classes, n, f)
                if shap_values.shape[0] == len(X_shap):
                    sv = shap_values[:, :, 1]        # (n, f, classes) → class 1
                else:
                    sv = shap_values[1]              # (classes, n, f) → class 1
            else:
                sv = shap_values                     # (n, f) — regression or binary
        else:
            sv = np.array(shap_values, dtype=np.float64)

        # Ensure sv is always a plain float64 numpy array (n_samples, n_features)
        sv = np.asarray(sv, dtype=np.float64)
        if sv.ndim != 2:
            raise ValueError(f"[SHAP] Unexpected sv shape after extraction: {sv.shape}")
        _log(logger, f"  shap_values shape={sv.shape}  dtype={sv.dtype}")

        mean_abs = np.mean(np.abs(sv), axis=0)
        # Convert to plain Python ints to avoid numpy scalar indexing errors
        top_idx  = [int(i) for i in np.argsort(mean_abs)[::-1][:top_n]]
        top_vals = [float(mean_abs[i]) for i in top_idx]
        top_lbls = [all_names[i] for i in top_idx]

        os.makedirs(MODEL_DIR, exist_ok=True)
        prefix = f"{output_prefix}_{split_tag}"

        # ── Print top-N table ─────────────────────────────────────────────────
        _log(logger, f"  {'Rank':>4}  {'Feature':35s}  {'Mean |SHAP|':>11}")
        _log(logger, f"  {'─'*4}  {'─'*35}  {'─'*11}")
        for rank, (lbl, val) in enumerate(zip(top_lbls, top_vals), 1):
            _log(logger, f"  {rank:4d}  {lbl:35s}  {val:11.6f}")

        # ── CSV export ────────────────────────────────────────────────────────
        import csv
        csv_path = os.path.join(MODEL_DIR, f"{prefix}_shap_top{top_n}.csv")
        with open(csv_path, 'w', newline='') as f:
            w = csv.writer(f)
            w.writerow(['rank', 'feature', 'mean_abs_shap'])
            for rank, (lbl, val) in enumerate(zip(top_lbls, top_vals), 1):
                w.writerow([rank, lbl, f"{val:.8f}"])
        _log(logger, f"[SHAP] csv  → {csv_path}")

        # ── Bar chart ─────────────────────────────────────────────────────────
        if 'bar' in plots:
            colours = []
            for l in top_lbls:
                if l.startswith('emb_'):              colours.append('#378ADD')  # PLM embedding
                elif l.startswith('1mer_'):            colours.append('#534AB7')  # 1-mer
                elif l.startswith('2mer_'):            colours.append('#7F77DD')  # 2-mer
                elif l.startswith('3mer_'):            colours.append('#AFA9EC')  # 3-mer
                elif l.startswith('oh_vh'):            colours.append('#E07B39')  # one-hot VH
                elif l.startswith('oh_vl'):            colours.append('#F5B942')  # one-hot VL
                elif l.startswith('oh_hcdr3') or l.startswith('oh_cdr3'):
                                                       colours.append('#C94A8C')  # one-hot HCDR3
                elif l.startswith('oh_'):              colours.append('#E07B39')  # one-hot other
                elif l.startswith('cdr3_') or l.startswith('vh_'):
                                                       colours.append('#1D9E75')  # biophysical
                else:                                  colours.append('#1D9E75')  # biophysical

            fig, ax = plt.subplots(figsize=(8.27, min(11.69, max(5, top_n * 0.20))))
            ax.barh(range(len(top_idx)), top_vals, color=colours)
            ax.set_yticks(range(len(top_idx)))
            ax.set_yticklabels(top_lbls, fontsize=7)
            ax.invert_yaxis()
            ax.set_xlabel('Mean |SHAP value|')
            ax.set_title(f'SHAP top {top_n} — {output_prefix} [{split_tag}]')
            patches = [mpatches.Patch(color='#378ADD', label='embedding'),
                       mpatches.Patch(color='#534AB7', label='1-mer'),
                       mpatches.Patch(color='#7F77DD', label='2-mer'),
                       mpatches.Patch(color='#AFA9EC', label='3-mer'),
                       mpatches.Patch(color='#1D9E75', label='biophysical'),
                       mpatches.Patch(color='#E07B39', label='one-hot VH'),
                       mpatches.Patch(color='#F5B942', label='one-hot VL'),
                       mpatches.Patch(color='#C94A8C', label='one-hot HCDR3')]
            # Only show legend entries that are actually present
            _present = set(colours)
            _color_map = {'#378ADD': 'embedding', '#534AB7': '1-mer',
                          '#7F77DD': '2-mer', '#AFA9EC': '3-mer',
                          '#1D9E75': 'biophysical', '#E07B39': 'one-hot VH',
                          '#F5B942': 'one-hot VL', '#C94A8C': 'one-hot HCDR3'}
            patches = [mpatches.Patch(color=c, label=l)
                       for c, l in _color_map.items() if c in _present]
            ax.legend(handles=patches, loc='lower right', fontsize=8)
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_bar.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] bar  → {out}")

        # ── Beeswarm ──────────────────────────────────────────────────────────
        if 'beeswarm' in plots:
            _shap_lib.summary_plot(sv[:, top_idx],
                                   X_shap[:, top_idx] if hasattr(X_shap, '__getitem__') else X_shap,
                                   feature_names=top_lbls,
                                   show=False, max_display=top_n)
            plt.title(f'SHAP beeswarm — {output_prefix} [{split_tag}]')
            plt.gcf().set_size_inches(8.27, min(11.69, max(5, top_n * 0.22)))
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_beeswarm.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] beeswarm → {out}")

        # ── Heatmap ───────────────────────────────────────────────────────────
        if 'heatmap' in plots:
            n_heat = min(top_n, 20)
            sub    = sv[:min(100, len(sv)), :][:, top_idx[:n_heat]]   # top_idx is plain int list
            lbls   = top_lbls[:n_heat]
            fig, ax = plt.subplots(figsize=(11, 7))
            vmax    = np.abs(sub).max()
            im      = ax.imshow(sub.T, aspect='auto', cmap='RdBu_r',
                                vmin=-vmax, vmax=vmax)
            plt.colorbar(im, ax=ax, label='SHAP value')
            ax.set_yticks(range(len(lbls)))
            ax.set_yticklabels(lbls, fontsize=7)
            ax.set_xlabel('Sample index')
            ax.set_title(f'SHAP heatmap top {n_heat} — {output_prefix} [{split_tag}]')
            plt.tight_layout()
            out = os.path.join(MODEL_DIR, f"{prefix}_shap_heatmap.png")
            plt.savefig(out, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"[SHAP] heatmap → {out}")

        # ── Waterfall — individual TIFF/JPEG + PPT slide per antibody ────────
        # YAML controls:
        #   max_waterfall_samples  : max antibodies          (default 50)
        #   waterfall_top_features : features per plot       (default 20)
        #   waterfall_fig_size     : [w, h] in inches        (default [8, 6])
        #   waterfall_format       : 'tiff' | 'jpeg' | 'png' (default 'tiff')
        #   pub_dpi                : image DPI               (default 300)
        #   waterfall_ppt          : also build PPT          (default true)
        _log(logger, f"\n[SHAP] Generating waterfall plots ...")
        _log(logger, f"  shap enabled={sh_cfg.get('enabled', False)}  "
                     f"max_wf={max_wf}  wf_top_n={wf_top_n}  "
                     f"fmt={sh_cfg.get('waterfall_format','tiff')}  "
                     f"ppt={sh_cfg.get('waterfall_ppt', True)}")
        try:
            import matplotlib.patches as _mp

            # ── Read YAML settings ────────────────────────────────────────────
            n_wf      = min(len(sv), max_wf)
            wf_size   = sh_cfg.get('waterfall_fig_size', [8, 6])
            wf_fmt    = sh_cfg.get('waterfall_format', 'tiff').lower().strip('.')
            make_ppt  = sh_cfg.get('waterfall_ppt', True)
            if wf_fmt not in ('tiff', 'jpeg', 'jpg', 'png'):
                wf_fmt = 'tiff'

            if n_wf < len(sv):
                _log(logger, f"  Waterfall limited to first {n_wf}/{len(sv)} "
                             f"(max_waterfall_samples={max_wf})")

            # ── Output folder: {MODEL_DIR}/{prefix}_waterfalls/ ───────────────
            wf_dir = os.path.join(MODEL_DIR, f"{prefix}_shap_waterfalls")
            os.makedirs(wf_dir, exist_ok=True)

            # ── Predictions — use full feature matrix for correct probabilities ────
            # X_shap is sliced to ne_idx only. For predict_proba, use _X_full_wf
            # (full feature matrix) so P(PASS) matches model's actual output.
            if self.task == 'classification':
                _thresh         = getattr(self, 'recommended_threshold', 0.5)
                all_probs       = self.model.predict_proba(_X_full_wf[:n_wf])[:, 1]
                all_labels_pred = (all_probs >= _thresh).astype(int)
            else:
                all_probs       = self.model.predict(_X_full_wf[:n_wf])
                all_labels_pred = None

            expected_val = float(
                explainer.expected_value[1]
                if isinstance(explainer.expected_value, (list, np.ndarray))
                else explainer.expected_value
            )
            task_str = "CL"  if self.task == 'classification' else "REG"
            lm_str   = lm_name if lm_name else "unknown_lm"
            db_str   = db_name if db_name else output_prefix

            saved_paths = []

            for s_idx in range(n_wf):
                bc  = barcodes[s_idx]                       if barcodes and s_idx < len(barcodes) else str(s_idx)
                act = actual_labels[s_idx]                       if actual_labels is not None                       and s_idx < len(actual_labels) else None

                # Features
                n_f       = min(wf_top_n, len(top_idx))
                shap_vals = [float(sv[s_idx][i])       for i in top_idx[:n_f]]
                feat_vals = [float(X_shap[s_idx, i])   for i in top_idx[:n_f]]
                lbls_s    = top_lbls[:n_f]

                triples = sorted(zip(shap_vals, lbls_s, feat_vals),
                                 key=lambda x: abs(x[0]), reverse=True)
                wf_s    = [x[0] for x in triples]
                wf_l    = [f"{x[1]} = {x[2]:.3g}" for x in triples]
                cols    = ['#D62728' if v < 0 else '#1F77B4' for v in wf_s]

                cumsum  = expected_val
                bottoms = []
                for v in wf_s:
                    bottoms.append(cumsum)
                    cumsum += v

                # ── Fixed figure size — no scaling ───────────────────────────
                fig, ax = plt.subplots(figsize=(wf_size[0], wf_size[1]))

                # Bar height and font size fixed to figure size / n_features
                bar_h = 0.70   # fixed — same solid bar height regardless of n features
                fsz   = 8.5    # fixed font size — consistent across all plots

                # Compute x range once for inside/outside text decision
                _all_ends = [b + v for b, v in zip(bottoms, wf_s)]
                _xmin  = min(bottoms + [expected_val] + _all_ends)
                _xmax  = max(bottoms + [expected_val] + _all_ends)
                _xrng  = max(abs(_xmax - _xmin), 1e-6)

                for i, (v, lbl, bot, col) in enumerate(
                        zip(wf_s, wf_l, bottoms, cols)):
                    ax.barh(i, v, left=bot, color=col,
                            height=bar_h, edgecolor='white', linewidth=0.4)
                    bw = abs(v)
                    # Put text inside bar only if bar is wide enough to read
                    inside = bw / _xrng > 0.08
                    if inside:
                        xt = bot + v / 2
                        ha = 'center'
                        tc = 'white'
                    else:
                        pad = _xrng * 0.012
                        xt  = bot + v + pad * (1 if v >= 0 else -1)
                        ha  = 'left' if v >= 0 else 'right'
                        tc  = col
                    # Only show text if value is meaningful (> 0.5% of range)
                    if bw / _xrng > 0.005:
                        ax.text(xt, i, f"{v:+.4f}", va='center', ha=ha,
                                fontsize=max(6.5, fsz - 0.5),
                                color=tc,
                                fontweight='bold' if bw / _xrng > 0.05 else 'normal')

                ax.set_yticks(range(n_f))
                ax.set_yticklabels(wf_l, fontsize=max(7.0, fsz))
                ax.invert_yaxis()
                # Pad left so outside-bar text never overlaps y-axis labels
                _xpad = _xrng * 0.18
                ax.set_xlim(_xmin - _xpad, _xmax + _xrng * 0.06)
                ax.axvline(expected_val, color='#888',    lw=0.8, ls='--')
                ax.axvline(cumsum,       color='#222299', lw=2.0, ls='-')
                ax.set_xlabel('SHAP value  (cumulative impact on model output)', fontsize=max(8.0, fsz))
                ax.tick_params(axis='x', labelsize=max(7.0, fsz - 0.5))

                # Title
                _col_name = actual_col_name if actual_col_name else "label"
                act_part  = (f"  |  Actual ({actual_col_name}) = "
                             f"{'PASS' if int(act)==1 else 'FAIL'}"
                             ) if act is not None else ""
                if self.task == 'classification':
                    prob    = float(all_probs[s_idx])
                    ml_lbl  = "PASS" if all_labels_pred[s_idx] == 1 else "FAIL"
                    sc_part = f"P(PASS)={prob:.4f}  →  {ml_lbl}"
                else:
                    sc_part = f"Predicted={float(all_probs[s_idx]):.4f}"

                # Line 1: "SHAP waterfall — {ID}"
                # Line 2: P(PASS)=X.XXXX → PASS/FAIL  |  Actual (col)=PASS/FAIL  |  RF|lm|CL|db
                _act_str = ""
                if act is not None:
                    _act_val  = "PASS" if int(act) == 1 else "FAIL"
                    _col_name = actual_col_name if actual_col_name else "label"
                    _act_str  = f"  |  Actual ({_col_name}) = {_act_val}"

                _id_part   = f"IPI MLAbDev Platform: Prediction interpretability — {bc}"
                _info_part = f"{sc_part}{_act_str}  |  RandomForest | {lm_str} | {task_str} | {db_str}"
                ax.set_title(
                    f"{_id_part}\n{_info_part}",
                    fontsize=8.5, loc='center', pad=10
                )
                ax.legend(handles=[
                    _mp.Patch(color='#1F77B4', label='toward PASS  (class 1 +)'),
                    _mp.Patch(color='#D62728', label='toward FAIL  (class 0 −)'),
                    plt.Line2D([0],[0], color='#999',    ls='--', lw=1.0,
                               label=f'baseline={expected_val:.3f}'),
                    plt.Line2D([0],[0], color='#1A237E', ls='-',  lw=2.0,
                               label=f'final={cumsum if abs(cumsum) > 1e-4 else 0.0:.4f}'),
                ], fontsize=max(5.0, fsz - 1.0),
                   loc='lower right', framealpha=0.85,
                   edgecolor='#ccc')

                plt.tight_layout()

                # Safe filename from barcode
                bc_safe = str(bc).replace('/', '_').replace('\\', '_').replace(' ', '_')
                img_path = os.path.join(
                    wf_dir,
                    f"{s_idx+1:04d}_{bc_safe}_waterfall.{wf_fmt}"
                )
                save_kw = dict(dpi=pub_dpi, bbox_inches='tight')
                if wf_fmt in ('jpeg', 'jpg'):
                    save_kw['format']  = 'jpeg'
                    save_kw['pil_kwargs'] = {'quality': 95}
                elif wf_fmt == 'tiff':
                    save_kw['format'] = 'tiff'
                plt.savefig(img_path, **save_kw)
                plt.close()
                saved_paths.append(img_path)

            _log(logger, f"[SHAP] {n_wf} waterfall images ({wf_fmt}, "
                         f"{pub_dpi} DPI) → {wf_dir}/")

            # ── PPT — one slide per antibody ──────────────────────────────────
            if make_ppt:
                try:
                    from pptx import Presentation as _Prs
                    from pptx.util import Inches, Pt
                    from pptx.enum.text import PP_ALIGN
                    import io as _io

                    prs   = _Prs()
                    # Widescreen 16:9 slide
                    prs.slide_width  = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    blank_layout = prs.slide_layouts[6]  # blank

                    for img_path, s_idx in zip(saved_paths,
                                               range(len(saved_paths))):
                        bc  = barcodes[s_idx]                               if barcodes and s_idx < len(barcodes) else str(s_idx)
                        slide = prs.slides.add_slide(blank_layout)

                        # Image — centered, leaving margin for title
                        img_w = Inches(wf_size[0])
                        img_h = Inches(wf_size[1])
                        left  = (prs.slide_width  - img_w) / 2
                        top   = Inches(0.6)
                        slide.shapes.add_picture(
                            img_path, left, top,
                            width=img_w, height=img_h
                        )

                        # Slide number + barcode footer
                        txb = slide.shapes.add_textbox(
                            Inches(0.15), Inches(7.1),
                            Inches(13.0), Inches(0.35)
                        )
                        tf  = txb.text_frame
                        tf.text = (f"{s_idx+1}/{n_wf}  |  {bc}  |  "
                                   f"RF | {lm_str} | {task_str} | {db_str}  |  "
                                   f"{split_tag}")
                        from pptx.dml.color import RGBColor
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0]
                        run.font.size = Pt(7)
                        run.font.color.rgb = RGBColor(0x88, 0x87, 0x80)

                    ppt_path = os.path.join(MODEL_DIR,
                                            f"{prefix}_shap_waterfall.pptx")
                    prs.save(ppt_path)
                    _log(logger,
                         f"[SHAP] waterfall PPT ({n_wf} slides) → {ppt_path}")

                except ImportError:
                    _log(logger,
                         "[SHAP] python-pptx not installed — PPT skipped.\n"
                         "  pip install python-pptx")
                except Exception as _pe:
                    import traceback as _tb
                    _log(logger, f"[SHAP] PPT failed — {_pe}")
                    _log(logger, _tb.format_exc())

        except Exception as _we:
            import traceback as _tb
            _log(logger, f"[SHAP] waterfall failed — {_we}")
            _log(logger, _tb.format_exc())
        _log(logger, f"[SHAP] {split_tag} complete.")

    # ── k-fold validation ──────────────────────────────────────────────────

    @classmethod
    def kfold_validation(cls,
                         data,
                         X_df:          pd.DataFrame,
                         y,
                         embeddings:    np.ndarray  = None,
                         embedding_lm:  str         = '',
                         title:         str         = "XGBoost",
                         kfold:         int         = 10,
                         target:        str         = "psr_filter",
                         cluster_col:   str         = "HCDR3_CLUSTER_0.8",
                         db_stem:       str         = "",
                         override_features: dict    = None,
                         cost_fn:       float       = 3.0,
                         cost_fp:       float       = 1.0):
        """
        CDR3-cluster stratified k-fold cross-validation.
        Identical interface to RandomForestModel.kfold_validation().
        """
        import yaml as _yaml

        os.makedirs(MODEL_DIR, exist_ok=True)

        # ── Logging ───────────────────────────────────────────────────────
        ts       = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        _db_tag  = f"_{db_stem}" if db_stem else ""
        log_name = (f"kfold_{target}_{embedding_lm}_xgboost"
                    f"{_db_tag}_k{kfold}_{ts}.log")
        logger   = _setup_logger(os.path.join(MODEL_DIR, log_name))
        _log(logger, f"[log] {'='*58}")
        _log(logger, f"[log] Started : {datetime.datetime.now()}")
        _log(logger, f"[log] kfold_validation  target={target}  "
                     f"lm={embedding_lm}  k={kfold}")
        _log(logger, f"[log] {'='*58}\n")

        # ── Base config ───────────────────────────────────────────────────
        _holder        = cls.__new__(cls)
        _holder.config = copy.deepcopy(_XGB_DEFAULT_CONFIG)
        _cfg_path      = "config/xgboost.yaml"
        if os.path.exists(_cfg_path):
            with open(_cfg_path) as _f:
                _deep_merge(_holder.config, _yaml.safe_load(_f) or {})
        resolved_cfg = _holder.config

        # ── Auto-detect task from y values ────────────────────────────────
        _kfold_task = resolved_cfg.get('task', 'classification')
        _y_tmp_kf   = np.asarray(y, dtype=float)
        _is_bin_kf  = (len(set(_y_tmp_kf.tolist())) <= 2 and
                       set(_y_tmp_kf.tolist()).issubset({0, 1, 0.0, 1.0}))
        # Always auto-detect: binary {0,1} → classification, else → regression
        # This overrides yaml task: auto / classification / anything
        if _is_bin_kf:
            _kfold_task = 'classification'
        else:
            _kfold_task = 'regression'
            if resolved_cfg.get('task') != 'regression':
                _log(logger, f"[kfold] Auto-detected task=regression "
                             f"({len(set(_y_tmp_kf.tolist()))} unique values — not binary)")
            resolved_cfg['task'] = 'regression'   # propagate to inst.config
        y_arr = np.asarray(y, dtype=float if _kfold_task == 'regression' else int)
        n     = len(y_arr)
        mean_fpr = np.linspace(0, 1, 100)

        # ── Splitter ──────────────────────────────────────────────────────
        # StratifiedGroupKFold/StratifiedKFold require integer class labels
        # → use GroupKFold/KFold for regression targets
        from sklearn.model_selection import GroupKFold, KFold
        kfold_actual = kfold
        if cluster_col in data.columns:
            groups          = data[cluster_col].values
            n_unique_groups = len(np.unique(groups))
            if n_unique_groups < kfold_actual:
                kfold_actual = n_unique_groups
                _log(logger, f"[kfold] {n_unique_groups} clusters → "
                             f"reducing folds to {kfold_actual}")
            if _kfold_task == 'regression':
                splitter   = GroupKFold(n_splits=kfold_actual)
                split_iter = splitter.split(np.arange(n), y_arr, groups)
                _log(logger, f"[kfold] GroupKFold on '{cluster_col}' (regression) "
                             f"({n_unique_groups} clusters, {kfold_actual} folds)")
            elif n_unique_groups == n:
                _log(logger, "[kfold] all-singleton clusters → StratifiedKFold")
                splitter   = StratifiedKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr)
            else:
                splitter   = StratifiedGroupKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr, groups)
                _log(logger, f"[kfold] StratifiedGroupKFold on '{cluster_col}' "
                             f"({n_unique_groups} clusters, {kfold_actual} folds)")
        else:
            _log(logger, f"[kfold] WARNING: '{cluster_col}' not found → "
                         f"{'KFold' if _kfold_task=='regression' else 'StratifiedKFold'}")
            if _kfold_task == 'regression':
                splitter   = KFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n))
            else:
                splitter   = StratifiedKFold(n_splits=kfold_actual, shuffle=True, random_state=42)
                split_iter = splitter.split(np.arange(n), y_arr)

        # ── Fold loop ─────────────────────────────────────────────────────
        tprs, aucs_list, fold_metrics, all_records = [], [], [], []
        best_fold_auc = -1.0
        best_fold_num = -1
        best_fold_state = best_fold_cfg = best_fold_fb = None

        _log(logger, f"\n{'='*62}")
        _log(logger, f"  {kfold_actual}-FOLD CV — XGBoost | {target.upper()} | "
                     f"{embedding_lm} | task={_kfold_task}")
        _log(logger, f"{'─'*62}")
        if _kfold_task == 'regression':
            _log(logger, f"  {'Fold':>5}  {'R2':>7}  {'MAE':>7}  {'Pearson':>8}  {'Spearman':>8}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*8}  {'─'*8}")
        else:
            _log(logger, f"  {'Fold':>5}  {'AUC':>7}  {'Acc':>7}  {'F1':>7}  "
                         f"{'Prec':>7}  {'Rec':>7}  {'Rec(F)':>7}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")

        plt.figure(figsize=(8, 7))

        for fold, (tr_idx, va_idx) in enumerate(split_iter, 1):
            # Leakage check
            if cluster_col in data.columns:
                leaked = set(groups[tr_idx]) & set(groups[va_idx])
                status = f"[WARN] {len(leaked)} cluster(s) leaked" if leaked else "[OK]  No CDR3 leakage"
                _log(logger, f"\n── Fold {fold}/{kfold_actual} ── {status}")
            else:
                _log(logger, f"\n── Fold {fold}/{kfold_actual} ──")

            y_tr = y_arr[tr_idx]; y_va = y_arr[va_idx]
            if _kfold_task == 'regression':
                _log(logger, f"  Train={len(tr_idx):,} mu={y_tr.mean():.4f}  "
                             f"Val={len(va_idx):,} mu={y_va.mean():.4f}")
            else:
                _log(logger, f"  Train={len(tr_idx):,} pos={y_tr.mean():.1%}  "
                             f"Val={len(va_idx):,} pos={y_va.mean():.1%}")

            X_df_tr = X_df.iloc[tr_idx]; X_df_va = X_df.iloc[va_idx]
            emb_tr  = embeddings[tr_idx] if embeddings is not None else None
            emb_va  = embeddings[va_idx] if embeddings is not None else None

            inst      = cls(config=resolved_cfg, verbose=False)
            inst.task = _kfold_task   # already auto-detected above
            if override_features:
                _oh_seq = override_features.pop('_onehot_sequence', None)
                inst.config['features'].update(override_features)
                if _oh_seq:
                    inst.config.setdefault('onehot', {})['sequence'] = _oh_seq
            inst.apply_lm_profile(embedding_lm, logger=logger)
            inst.fb_ = FeatureBuilder(inst.config)
            X_tr = inst.fb_.fit_transform(X_df_tr, emb_tr)
            X_va = inst.fb_.transform(X_df_va, emb_va)

            inst.model = inst._build_xgb(y_tr)

            # Log hyperparameters used for this fold (safe for any numeric type)
            try:
                _p = inst.model.get_params()
                def _fmt(v, f='.4f'):
                    try:    return format(float(v), f)
                    except: return str(v)
                _log(logger, f"  [params] n_est={_p.get('n_estimators')}"
                             f"  depth={_p.get('max_depth')}"
                             f"  lr={_fmt(_p.get('learning_rate'))}"
                             f"  subsample={_fmt(_p.get('subsample'))}"
                             f"  colsample={_fmt(_p.get('colsample_bytree'))}"
                             f"  gamma={_fmt(_p.get('gamma'))}"
                             f"  min_child_w={_p.get('min_child_weight')}"
                             f"  reg_α={_fmt(_p.get('reg_alpha'))}"
                             f"  reg_λ={_fmt(_p.get('reg_lambda'))}"
                             f"  scale_pos_w={_fmt(_p.get('scale_pos_weight', 1.0), '.2f')}")
            except Exception as _pe:
                _log(logger, f"  [params] WARNING: could not log params — {_pe}")

            esr = inst.config['training'].get('early_stopping_rounds', 0)
            if esr > 0:
                inst.model.fit(X_tr, y_tr,
                               eval_set              = [(X_va, y_va)],
                               early_stopping_rounds = esr,
                               verbose               = False)
                _log(logger, f"  [early_stop] best_iteration={inst.model.best_iteration}")
            else:
                inst.model.fit(X_tr, y_tr)

            # ── Classification metrics ────────────────────────────────
            if inst.task == 'classification':
                probs = inst.model.predict_proba(X_va)[:, 1]
                preds = (probs >= 0.5).astype(int)

                if len(set(y_va)) < 2:
                    _log(logger, f"  Skipping fold {fold} — only one class in val.")
                    continue

                fold_auc  = roc_auc_score(y_va,  probs)
                fold_acc  = accuracy_score(y_va,  preds)
                fold_f1   = f1_score(y_va,        preds, zero_division=0)
                fold_prec = precision_score(y_va, preds, zero_division=0)
                fold_rec  = recall_score(y_va,    preds, zero_division=0)
                fold_recf = recall_score(y_va,    preds, pos_label=0, zero_division=0)

                fold_metrics.append({'fold': fold, 'auc': fold_auc, 'acc': fold_acc,
                                     'f1': fold_f1, 'precision': fold_prec,
                                     'recall': fold_rec, 'rec_fail': fold_recf})
                aucs_list.append(fold_auc)

                fpr, tpr, _ = roc_curve(y_va, probs)
                tprs.append(np.interp(mean_fpr, fpr, tpr)); tprs[-1][0] = 0.0
                plt.plot(fpr, tpr, lw=1, alpha=0.3,
                         label=f'Fold {fold} ({fold_auc:.3f})')

                marker = " ←" if fold == best_fold_num else ""
                _log(logger, f"  {fold:5d}  {fold_auc:7.4f}  {fold_acc:7.4f}  "
                             f"{fold_f1:7.4f}  {fold_prec:7.4f}  {fold_rec:7.4f}  "
                             f"{fold_recf:7.4f}{marker}")

            # ── Regression metrics ────────────────────────────────────────
            else:
                from sklearn.metrics import r2_score, mean_absolute_error
                from scipy.stats import pearsonr, spearmanr
                preds_r  = inst.model.predict(X_va)
                fold_r2  = r2_score(y_va, preds_r)
                fold_mae = mean_absolute_error(y_va, preds_r)
                fold_rp  = pearsonr(y_va,  preds_r)[0]
                fold_rs  = spearmanr(y_va, preds_r)[0]

                fold_metrics.append({'fold': fold, 'r2': fold_r2,
                                     'mae': fold_mae, 'pearson': fold_rp,
                                     'spearman': fold_rs})
                aucs_list.append(fold_rs)  # use Spearman as primary metric

                # Scatter plot per fold
                plt.scatter(y_va, preds_r, alpha=0.3, s=10,
                            label=f'Fold {fold} (ρ={fold_rs:.3f})')

                marker = " ←" if fold == best_fold_num else ""
                _log(logger, f"  {fold:5d}  {fold_r2:7.4f}  {fold_mae:7.4f}  "
                             f"{fold_rp:7.4f}  {fold_rs:7.4f}{marker}")

            # Save fold checkpoint
            _task_tag  = "_regression" if _kfold_task == 'regression' else ""
            fold_path = os.path.join(
                MODEL_DIR,
                f"xgboost_{target}_{embedding_lm}{_db_tag}_fold{fold}_k{kfold_actual}{_task_tag}.pkl")
            inst.save(fold_path)

            # Primary metric: AUC (classification) or Spearman (regression)
            _fold_primary = fold_auc if inst.task == 'classification' else fold_rs
            if _fold_primary > best_fold_auc:
                best_fold_auc   = _fold_primary
                best_fold_num   = fold
                best_fold_state = copy.deepcopy(inst.model)
                best_fold_cfg   = copy.deepcopy(inst.config)
                best_fold_fb    = copy.deepcopy(inst.fb_)

            bcs = X_df_va.index.astype(str).tolist()
            if inst.task == 'classification':
                for bc, true, pred, prob in zip(bcs, y_va, preds, probs):
                    all_records.append({'BARCODE': bc, 'fold': fold,
                                        'true': true, 'pred': pred, 'prob': prob})
            else:
                for bc, true, pred in zip(bcs, y_va, preds_r):
                    all_records.append({'BARCODE': bc, 'fold': fold,
                                        'true': true, 'pred': pred, 'prob': pred})

        if not aucs_list:
            _log(logger, "[kfold] No valid folds — check class distribution.")
            return

        # ── Summary + plot (task-aware) ──────────────────────────────────
        mean_auc = float(np.mean(aucs_list))
        std_auc  = float(np.std(aucs_list))

        _is_reg  = (fold_metrics and 'r2' in fold_metrics[0])

        if not _is_reg:
            # ── Classification summary ────────────────────────────────────
            mean_acc  = float(np.mean([m['acc']       for m in fold_metrics]))
            mean_f1   = float(np.mean([m['f1']        for m in fold_metrics]))
            mean_prec = float(np.mean([m['precision'] for m in fold_metrics]))
            mean_rec  = float(np.mean([m['recall']    for m in fold_metrics]))
            mean_recf = float(np.mean([m['rec_fail']  for m in fold_metrics]))

            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            for m in fold_metrics:
                mark = " ←" if m['fold'] == best_fold_num else ""
                _log(logger, f"  {m['fold']:5d}  {m['auc']:7.4f}  {m['acc']:7.4f}  "
                             f"{m['f1']:7.4f}  {m['precision']:7.4f}  {m['recall']:7.4f}  "
                             f"{m['rec_fail']:7.4f}{mark}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            _log(logger, f"  {'Mean':>5}  {mean_auc:7.4f}  {mean_acc:7.4f}  "
                         f"{mean_f1:7.4f}  {mean_prec:7.4f}  {mean_rec:7.4f}  {mean_recf:7.4f}")
            _log(logger, f"  {'±Std':>5}  {std_auc:7.4f}")
            _log(logger, f"  Best fold : {best_fold_num}  (AUC={best_fold_auc:.4f})")
            _log(logger, f"  Rec(Fail) : {mean_recf:.4f}  ← minority-class recall")

            # ROC plot
            mean_tpr = np.mean(tprs, axis=0); mean_tpr[-1] = 1.0
            std_tpr  = np.std(tprs,  axis=0)
            plt.plot(mean_fpr, mean_tpr, 'b', lw=3,
                     label=f'Mean ROC (AUC={mean_auc:.3f}±{std_auc:.3f})')
            plt.fill_between(mean_fpr,
                             np.maximum(mean_tpr - std_tpr, 0),
                             np.minimum(mean_tpr + std_tpr, 1),
                             color='lightblue', alpha=0.3, label='±1 std')
            plt.plot([0, 1], [0, 1], '--', color='gray', lw=0.8)
            plt.xlim([0, 1]); plt.ylim([0, 1.05])
            plt.xlabel('False Positive Rate'); plt.ylabel('True Positive Rate')
            plt.title(
                f'XGBoost — {target.upper()}\n'
                f'{kfold_actual}-Fold SGKF ROC  |  {embedding_lm}\n'
                f'Acc={mean_acc:.3f}  F1={mean_f1:.3f}  Prec={mean_prec:.3f}  '
                f'Rec={mean_rec:.3f}  Rec(Fail)={mean_recf:.3f}', fontsize=9)
            plt.legend(loc='lower right', fontsize=7)
            plt.grid(alpha=0.3); plt.tight_layout()
            plot_path = os.path.join(
                MODEL_DIR,
                f"CV_ROC_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}{_task_tag}.png")
            plt.savefig(plot_path, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"  ROC plot  : {plot_path}")

        else:
            # ── Regression summary ────────────────────────────────────────
            mean_r2  = float(np.mean([m['r2']       for m in fold_metrics]))
            mean_mae = float(np.mean([m['mae']      for m in fold_metrics]))
            mean_rp  = float(np.mean([m['pearson']  for m in fold_metrics]))
            mean_rs  = float(np.mean([m['spearman'] for m in fold_metrics]))
            std_rs   = float(np.std([m['spearman']  for m in fold_metrics]))

            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            _log(logger, f"  {'Fold':>5}  {'R2':>7}  {'MAE':>7}  "
                         f"{'Pearson':>7}  {'Spearman':>8}")
            for m in fold_metrics:
                mark = " ←" if m['fold'] == best_fold_num else ""
                _log(logger, f"  {m['fold']:5d}  {m['r2']:7.4f}  {m['mae']:7.4f}  "
                             f"{m['pearson']:7.4f}  {m['spearman']:8.4f}{mark}")
            _log(logger, f"  {'─'*5}  {'─'*7}  {'─'*7}  {'─'*7}  {'─'*7}")
            _log(logger, f"  {'Mean':>5}  {mean_r2:7.4f}  {mean_mae:7.4f}  "
                         f"{mean_rp:7.4f}  {mean_rs:8.4f}")
            _log(logger, f"  {'±Std':>5}  {std_auc:7.4f}  (Spearman ρ)")
            _log(logger, f"  Best fold : {best_fold_num}  (ρ={best_fold_auc:.4f})")

            # Scatter plot with identity line
            _all_y   = [r['true'] for r in all_records]
            _all_p   = [r['pred'] for r in all_records]
            plt.scatter(_all_y, _all_p, alpha=0.3, s=10, color='steelblue')
            _mn = min(_all_y + _all_p); _mx = max(_all_y + _all_p)
            plt.plot([_mn, _mx], [_mn, _mx], '--', color='gray', lw=1)
            plt.xlabel('True'); plt.ylabel('Predicted')
            plt.title(
                f'XGBoost Regression — {target.upper()}\n'
                f'{kfold_actual}-Fold  |  {embedding_lm}\n'
                f'R²={mean_r2:.3f}  MAE={mean_mae:.3f}  '
                f'Pearson={mean_rp:.3f}  Spearman={mean_rs:.3f}', fontsize=9)
            plt.grid(alpha=0.3); plt.tight_layout()
            plot_path = os.path.join(
                MODEL_DIR,
                f"CV_scatter_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}{_task_tag}.png")
            plt.savefig(plot_path, dpi=150, bbox_inches='tight'); plt.close()
            _log(logger, f"  Scatter plot : {plot_path}")
            mean_acc = mean_f1 = mean_prec = mean_rec = mean_recf = float('nan')

        # ── Best-fold checkpoint ──────────────────────────────────────────
        _task_tag = "_regression" if _kfold_task == 'regression' else ""
        best_path = None
        if best_fold_state is not None:
            best_path = os.path.join(
                MODEL_DIR,
                f"BEST_{target}_{embedding_lm}_xgboost{_db_tag}"
                f"_k{kfold_actual}_fold{best_fold_num}{_task_tag}.pkl")
            _best               = cls(config=best_fold_cfg, verbose=False)
            _best.model         = best_fold_state
            _best.fb_           = best_fold_fb
            _best.task          = 'classification'
            _best.save(best_path)
            _log(logger, f"\n[kfold] Best fold → {best_path}"
                         f"  (fold={best_fold_num}, AUC={best_fold_auc:.4f})")

        # ── Fold predictions CSV ──────────────────────────────────────────
        pred_path = None
        if all_records:
            pred_path = os.path.join(
                MODEL_DIR,
                f"fold_preds_{target}_{embedding_lm}_xgboost{_db_tag}_k{kfold_actual}{_task_tag}.csv")
            df_preds = pd.DataFrame(all_records)
            df_preds['best_fold'] = (df_preds['fold'] == best_fold_num).astype(int)
            df_preds.to_csv(pred_path, index=False)
            _log(logger, f"[kfold] Fold predictions → {pred_path}")

        # ── Threshold optimisation ────────────────────────────────────────
        if _THRESHOLD_OPT_AVAILABLE and pred_path and best_path:
            _log(logger, f"\n[threshold] Optimising ...")
            try:
                stability = run_full_threshold_pipeline(
                    fold_preds_csv = pred_path,
                    target         = target,
                    lm             = embedding_lm,
                    model          = 'xgboost',
                    db_stem        = db_stem,
                    best_ckpt_path = best_path,
                    output_dir     = MODEL_DIR,
                    cost_fp        = cost_fp,
                    cost_fn        = cost_fn)
                rec_t = stability.get('pooled_threshold', 0.5)
                cls.recommended_threshold = float(rec_t)
                _log(logger, f"  Pooled OOF threshold : {rec_t:.4f}")
            except Exception as e:
                _log(logger, f"[threshold] WARNING: {e} — defaulting to 0.5")

        _log(logger, f"\n[log] Finished : {datetime.datetime.now()}")
        return mean_auc, std_auc, mean_acc, mean_f1, mean_prec, mean_rec